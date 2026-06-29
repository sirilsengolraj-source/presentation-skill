#!/usr/bin/env python3
"""Targeted QA for layout-polish issues that generic geometry checks miss."""

from __future__ import annotations

import argparse
import json
import posixpath
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "r": "http://schemas.openxmlformats.org/package/2006/relationships",
}
DEFAULT_BANNED = [
    "follow previous tool",
    "external pptx",
    "updated skill",
    "demo deck",
    "sample deck",
    "placeholder",
]

DEFAULT_READABILITY = {
    "min_title_pt": 24.0,
    "min_body_pt": 11.0,
    "min_caption_pt": 7.5,
    "chart_label_min_pt": 8.0,
    "footer_reserved_inches": 0.25,
}


def _box(shape):
    return (
        shape.left.inches,
        shape.top.inches,
        shape.width.inches,
        shape.height.inches,
    )


def _contains(outer, inner, pad=0.02):
    ox, oy, ow, oh = outer
    ix, iy, iw, ih = inner
    return (
        ix >= ox + pad
        and iy >= oy + pad
        and ix + iw <= ox + ow - pad
        and iy + ih <= oy + oh - pad
    )


def _overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    overlap_x = min(ax + aw, bx + bw) - max(ax, bx)
    overlap_y = min(ay + ah, by + bh) - max(ay, by)
    return max(0.0, overlap_x), max(0.0, overlap_y)


def _center(box):
    x, y, w, h = box
    return (x + w / 2.0, y + h / 2.0)


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").strip()


def _load_json(path: Path):
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def _readability_contract(design_brief_path: Path | None):
    contract = dict(DEFAULT_READABILITY)
    if design_brief_path is None:
        return contract, False
    payload = _load_json(design_brief_path)
    if not isinstance(payload, dict):
        return contract, False
    brief_contract = payload.get("readability_contract")
    if not isinstance(brief_contract, dict):
        return contract, False
    for key in (
        "min_title_pt",
        "min_body_pt",
        "min_caption_pt",
        "chart_label_min_pt",
        "footer_reserved_inches",
    ):
        value = brief_contract.get(key)
        if isinstance(value, (int, float)):
            contract[key] = float(value)
    return contract, True


def _font_sizes(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    sizes = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size is not None:
            sizes.append(float(paragraph.font.size.pt))
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return sizes


def _text_role(shape, text, slide_h):
    box = _box(shape)
    top = box[1]
    height = box[3]
    lower = text.lower()
    if top >= slide_h - 0.75:
        return "caption"
    if height <= 0.36:
        return "caption"
    if lower.startswith(("source", "sources", "ref", "refs")):
        return "caption"
    if top <= 1.15 and len(text) <= 160:
        return "title"
    return "body"


def _shape_kind(shape):
    try:
        if not hasattr(shape, "auto_shape_type") or shape.auto_shape_type is None:
            return ""
        return str(shape.auto_shape_type).upper()
    except Exception:
        return ""


def _has_visible_fill(shape):
    try:
        return getattr(shape.fill, "type", None) is not None
    except Exception:
        return False


def _iter_text_shapes(slide):
    for idx, shape in enumerate(slide.shapes, start=1):
        text = _shape_text(shape)
        if text:
            yield idx, shape, text


def _iter_auto_shapes(slide):
    for idx, shape in enumerate(slide.shapes, start=1):
        kind = _shape_kind(shape)
        if kind:
            yield idx, shape, kind


def check_branding(slide_idx, text_shapes, banned):
    issues = []
    for shape_id, _, text in text_shapes:
        lowered = text.lower()
        for phrase in banned:
            if phrase in lowered:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_id": f"shape-{shape_id}",
                        "type": "residual_branding",
                        "severity": "error",
                        "text": text[:160],
                        "phrase": phrase,
                    }
                )
    return issues


def check_footer_overlap(slide_idx, text_shapes, slide_h, contract):
    issues = []
    footer_top = max(0.0, slide_h - 0.70)
    bottom_band = [
        (shape_id, shape, text)
        for shape_id, shape, text in text_shapes
        if shape.top.inches >= footer_top
    ]
    for i in range(len(bottom_band)):
        for j in range(i + 1, len(bottom_band)):
            left = bottom_band[i]
            right = bottom_band[j]
            overlap_x, overlap_y = _overlap(_box(left[1]), _box(right[1]))
            if overlap_x > 0.02 and overlap_y > 0.02:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_ids": [f"shape-{left[0]}", f"shape-{right[0]}"],
                        "type": "footer_text_overlap",
                        "severity": "error",
                        "delta_inches": round(min(overlap_x, overlap_y), 3),
                    }
                )
    try:
        footer_reserved = float(contract.get("footer_reserved_inches", DEFAULT_READABILITY["footer_reserved_inches"]))
    except (TypeError, ValueError):
        footer_reserved = float(DEFAULT_READABILITY["footer_reserved_inches"])
    reserve_top = max(0.0, slide_h - max(0.0, footer_reserved))
    for shape_id, shape, text in text_shapes:
        if getattr(shape, "has_table", False):
            continue
        box = _box(shape)
        bottom = box[1] + box[3]
        if bottom <= reserve_top + 0.02:
            continue
        role = _text_role(shape, text, slide_h)
        if role == "caption":
            continue
        issues.append(
            {
                "slide_index": slide_idx,
                "shape_id": f"shape-{shape_id}",
                "type": "footer_reserved_space_intrusion",
                "severity": "warning",
                "reserved_inches": round(footer_reserved, 2),
                "intrusion_inches": round(bottom - reserve_top, 3),
                "text": text[:120],
            }
        )
    return issues


def check_text_readability(slide_idx, text_shapes, slide_h, contract):
    issues = []
    thresholds = {
        "title": float(contract["min_title_pt"]),
        "body": float(contract["min_body_pt"]),
        "caption": float(contract["min_caption_pt"]),
    }
    for shape_id, shape, text in text_shapes:
        if getattr(shape, "has_table", False):
            continue
        sizes = _font_sizes(shape)
        if not sizes:
            continue
        min_font = min(sizes)
        role = _text_role(shape, text, slide_h)
        threshold = thresholds[role]
        if min_font + 0.05 < threshold:
            issues.append(
                {
                    "slide_index": slide_idx,
                    "shape_id": f"shape-{shape_id}",
                    "type": f"{role}_font_too_small",
                    "severity": "warning",
                    "font_pt": round(min_font, 1),
                    "min_allowed_pt": round(threshold, 1),
                    "text": text[:120],
                }
            )
    return issues


def check_table_readability(slide_idx, slide, contract):
    issues = []
    min_table_font = max(7.8, float(contract["min_caption_pt"]))
    for shape_id, shape in enumerate(slide.shapes, start=1):
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        row_count = len(table.rows)
        col_count = len(table.columns)
        if row_count > 10 or col_count > 6:
            issues.append(
                {
                    "slide_index": slide_idx,
                    "shape_id": f"shape-{shape_id}",
                    "type": "table_density_risk",
                    "severity": "warning",
                    "rows": row_count,
                    "columns": col_count,
                }
            )
        min_font = 99.0
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    if paragraph.font.size is not None:
                        min_font = min(min_font, float(paragraph.font.size.pt))
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            min_font = min(min_font, float(run.font.size.pt))
        if min_font != 99.0 and min_font < min_table_font:
            issues.append(
                {
                    "slide_index": slide_idx,
                    "shape_id": f"shape-{shape_id}",
                    "type": "table_font_too_small",
                    "severity": "warning",
                    "font_pt": round(min_font, 1),
                    "min_allowed_pt": round(min_table_font, 1),
                }
            )
    return issues


def check_stacked_text_gaps(slide_idx, auto_shapes, text_shapes):
    issues = []
    text_boxes = [(shape_id, shape, text, _box(shape)) for shape_id, shape, text in text_shapes]
    for shape_id, shape, kind in auto_shapes:
        if "RECTANGLE" not in kind:
            continue
        box = _box(shape)
        if box[1] >= 6.7:
            continue
        if not (0.6 <= box[3] <= 4.2 and 1.0 <= box[2] <= 7.0):
            continue
        inside = [item for item in text_boxes if _contains(box, item[3])]
        if len(inside) < 2:
            continue
        inside.sort(key=lambda item: item[3][1])
        for current, nxt in zip(inside, inside[1:]):
            current_box = current[3]
            next_box = nxt[3]
            overlap_x, _ = _overlap(current_box, next_box)
            min_w = min(current_box[2], next_box[2])
            if overlap_x < min_w * 0.35:
                continue
            gap = next_box[1] - (current_box[1] + current_box[3])
            if gap < 0.08:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_ids": [f"shape-{current[0]}", f"shape-{nxt[0]}"],
                        "container_shape_id": f"shape-{shape_id}",
                        "type": "stack_gap_too_small",
                        "severity": "warning",
                        "delta_inches": round(0.08 - gap, 3),
                    }
                )
                break
    return issues


def check_marker_centering(slide_idx, auto_shapes, text_shapes):
    issues = []
    text_candidates = [
        (shape_id, shape, text, _box(shape))
        for shape_id, shape, text in text_shapes
        if len(text.strip()) <= 3
    ]
    for shape_id, shape, kind in auto_shapes:
        if "ELLIPSE" not in kind:
            continue
        box = _box(shape)
        if not (0.2 <= box[2] <= 0.8 and abs(box[2] - box[3]) <= 0.08):
            continue
        circle_center = _center(box)
        matches = []
        for text_id, _, _, text_box in text_candidates:
            text_center = _center(text_box)
            if abs(text_center[0] - circle_center[0]) <= 0.35 and abs(text_center[1] - circle_center[1]) <= 0.35:
                matches.append((text_id, text_box, text_center))
        if not matches:
            continue
        text_id, text_box, text_center = min(
            matches,
            key=lambda item: abs(item[2][0] - circle_center[0]) + abs(item[2][1] - circle_center[1]),
        )
        dx = abs(text_center[0] - circle_center[0])
        dy = abs(text_center[1] - circle_center[1])
        if dx > 0.03 or dy > 0.03:
            issues.append(
                {
                    "slide_index": slide_idx,
                    "shape_ids": [f"shape-{shape_id}", f"shape-{text_id}"],
                    "type": "marker_label_off_center",
                    "severity": "error",
                    "delta_inches": round(max(dx, dy), 3),
                }
            )
    return issues


def _chart_slide_index_map(pptx_path: Path):
    chart_to_slide: dict[str, int] = {}
    with zipfile.ZipFile(pptx_path, "r") as archive:
        rel_names = [
            name
            for name in archive.namelist()
            if name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels")
        ]
        for rel_name in rel_names:
            slide_name = Path(rel_name).name
            raw_number = slide_name.removeprefix("slide").removesuffix(".xml.rels")
            try:
                slide_index = int(raw_number) - 1
            except ValueError:
                continue
            root = ET.fromstring(archive.read(rel_name))
            for rel in root.findall(".//r:Relationship", NS):
                target = rel.attrib.get("Target", "")
                rel_type = rel.attrib.get("Type", "")
                if "chart" not in rel_type and "charts/" not in target:
                    continue
                chart_part = posixpath.normpath(posixpath.join("ppt/slides", target)).lstrip("/")
                if chart_part.startswith("../"):
                    chart_part = posixpath.normpath(posixpath.join("ppt", chart_part[3:]))
                if chart_part.startswith("charts/"):
                    chart_part = f"ppt/{chart_part}"
                chart_to_slide[chart_part] = slide_index
    return chart_to_slide


def check_chart_headroom(pptx_path: Path, chart_slide_indexes: dict[str, int] | None = None):
    issues = []
    chart_slide_indexes = chart_slide_indexes or {}
    with zipfile.ZipFile(pptx_path, "r") as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/charts/chart") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            show_val = any(
                node.attrib.get("val") in {"1", "true", "True"}
                for node in root.findall(".//c:dLbls/c:showVal", NS)
            )
            if not show_val:
                continue
            axis_max_values = [
                float(node.attrib.get("val"))
                for node in root.findall(".//c:valAx/c:scaling/c:max", NS)
                if node.attrib.get("val")
            ]
            if not axis_max_values:
                continue
            axis_max = max(axis_max_values)
            point_values = [
                float(node.text)
                for node in root.findall(".//c:ser//c:val//c:v", NS)
                if node.text
            ]
            if point_values and max(point_values) >= axis_max:
                issues.append(
                    {
                        "chart_part": name,
                        "type": "chart_value_label_headroom_risk",
                        "severity": "warning",
                        **({"slide_index": chart_slide_indexes[name]} if name in chart_slide_indexes else {}),
                        "axis_max": axis_max,
                        "max_value": max(point_values),
                    }
                )
    return issues


def _ooxml_font_pt(node: ET.Element) -> float | None:
    raw = node.attrib.get("sz")
    if raw is None:
        return None
    try:
        return float(raw) / 100.0
    except ValueError:
        return None


def _chart_text_sizes(root: ET.Element, xpath: str) -> list[float]:
    sizes: list[float] = []
    for node in root.findall(xpath, NS):
        value = _ooxml_font_pt(node)
        if value is not None:
            sizes.append(value)
    return sizes


def check_chart_readability(
    pptx_path: Path,
    contract: dict[str, float],
    chart_slide_indexes: dict[str, int] | None = None,
):
    issues = []
    chart_slide_indexes = chart_slide_indexes or {}
    min_chart_font = float(contract.get("chart_label_min_pt", DEFAULT_READABILITY["chart_label_min_pt"]))
    role_paths = {
        "axis_label": [
            ".//c:catAx/c:txPr//a:defRPr",
            ".//c:valAx/c:txPr//a:defRPr",
            ".//c:serAx/c:txPr//a:defRPr",
            ".//c:dateAx/c:txPr//a:defRPr",
        ],
        "axis_title": [
            ".//c:catAx/c:title/c:txPr//a:defRPr",
            ".//c:valAx/c:title/c:txPr//a:defRPr",
            ".//c:serAx/c:title/c:txPr//a:defRPr",
            ".//c:dateAx/c:title/c:txPr//a:defRPr",
        ],
        "legend_label": [".//c:legend/c:txPr//a:defRPr"],
        "data_label": [".//c:dLbls/c:txPr//a:defRPr"],
    }
    with zipfile.ZipFile(pptx_path, "r") as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/charts/chart") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            for role, paths in role_paths.items():
                sizes: list[float] = []
                for xpath in paths:
                    sizes.extend(_chart_text_sizes(root, xpath))
                if not sizes:
                    continue
                min_font = min(sizes)
                if min_font + 0.05 < min_chart_font:
                    issues.append(
                        {
                            "chart_part": name,
                            "type": "chart_label_font_too_small",
                            "severity": "warning",
                            **({"slide_index": chart_slide_indexes[name]} if name in chart_slide_indexes else {}),
                            "role": role,
                            "font_pt": round(min_font, 1),
                            "min_allowed_pt": round(min_chart_font, 1),
                        }
                    )
    return issues


def main() -> int:
    parser = argparse.ArgumentParser(description="Targeted design QA")
    parser.add_argument("--input", required=True, help="Input PPTX")
    parser.add_argument("--report", help="Optional JSON report path")
    parser.add_argument(
        "--banned-phrase",
        action="append",
        default=[],
        help="Additional banned phrase to flag",
    )
    parser.add_argument(
        "--design-brief",
        help=(
            "Optional design_brief.json. When present, design QA uses its "
            "readability_contract thresholds for title/body/caption/table/chart text."
        ),
    )
    args = parser.parse_args()

    pptx_path = Path(args.input).expanduser().resolve()
    prs = Presentation(str(pptx_path))
    design_brief_path = Path(args.design_brief).expanduser().resolve() if args.design_brief else None
    readability_contract, enforce_text_readability = _readability_contract(design_brief_path)

    banned = [item.lower() for item in (DEFAULT_BANNED + args.banned_phrase)]
    issues = []
    slide_summaries = []

    for slide_idx, slide in enumerate(prs.slides):
        text_shapes = list(_iter_text_shapes(slide))
        auto_shapes = list(_iter_auto_shapes(slide))
        slide_h = prs.slide_height.inches
        slide_issues = []
        slide_issues.extend(check_branding(slide_idx, text_shapes, banned))
        slide_issues.extend(check_footer_overlap(slide_idx, text_shapes, slide_h, readability_contract))
        if enforce_text_readability:
            slide_issues.extend(check_text_readability(slide_idx, text_shapes, slide_h, readability_contract))
        slide_issues.extend(check_table_readability(slide_idx, slide, readability_contract))
        slide_issues.extend(check_stacked_text_gaps(slide_idx, auto_shapes, text_shapes))
        slide_issues.extend(check_marker_centering(slide_idx, auto_shapes, text_shapes))
        issues.extend(slide_issues)
        slide_summaries.append({"slide_index": slide_idx, "issue_count": len(slide_issues)})

    chart_slide_indexes = _chart_slide_index_map(pptx_path)
    chart_issues = check_chart_headroom(pptx_path, chart_slide_indexes)
    issues.extend(chart_issues)
    chart_readability_issues = check_chart_readability(pptx_path, readability_contract, chart_slide_indexes)
    issues.extend(chart_readability_issues)

    payload = {
        "input": str(pptx_path),
        "issue_count": len(issues),
        "error_count": sum(1 for item in issues if item.get("severity") == "error"),
        "warning_count": sum(1 for item in issues if item.get("severity") == "warning"),
        "readability_contract": readability_contract,
        "readability_contract_enforced": enforce_text_readability,
        "slides": slide_summaries,
        "issues": issues,
        "passed": not issues,
    }

    if args.report:
        report_path = Path(args.report).expanduser().resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    print(
        f"Design rules QA: {pptx_path}\n"
        f"  {payload['issue_count']} issue(s) | "
        f"errors={payload['error_count']} warnings={payload['warning_count']}"
    )
    for issue in issues:
        location = (
            f"slide {issue.get('slide_index', 0) + 1}"
            if "slide_index" in issue
            else issue.get("chart_part", "chart")
        )
        print(f"  - {location}: {issue.get('type')}")

    return 0 if payload["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
