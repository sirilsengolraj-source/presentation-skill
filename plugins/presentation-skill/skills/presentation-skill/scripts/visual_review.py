#!/usr/bin/env python3
"""Create a rendered-slide review packet and catch polish risks.

This is the deterministic half of the "look at the slides" loop. It does
not replace human/model visual judgment, but it gives agents a concrete packet:
contact sheet, per-slide thumbnails, JSON findings, and a markdown punch list.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import shutil
import subprocess
import sys
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:  # pragma: no cover - optional dependency error path
    Image = None  # type: ignore[assignment]
    ImageDraw = None  # type: ignore[assignment]
    ImageFont = None  # type: ignore[assignment]


EMU_PER_INCH = 914400.0
SLIDE_BOTTOM_SAFE = 0.35
SLIDE_SIDE_SAFE = 0.35


def _inches(value: int | float) -> float:
    return float(value) / EMU_PER_INCH if value else 0.0


def _run(cmd: list[str]) -> tuple[int, str]:
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    return result.returncode, result.stdout


def _load_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {}


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    chunks: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)
    return "\n".join(chunks).strip()


def _font_size_pt(shape: Any) -> float:
    if not getattr(shape, "has_text_frame", False):
        return 12.0
    for paragraph in shape.text_frame.paragraphs:
        paragraph_font = getattr(paragraph, "font", None)
        if paragraph_font is not None and paragraph_font.size is not None:
            return float(paragraph_font.size.pt)
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return 12.0


def _shape_record(slide_num: int, shape_idx: int, shape: Any) -> dict[str, Any] | None:
    text = _shape_text(shape)
    if not text:
        return None
    x = _inches(shape.left)
    y = _inches(shape.top)
    w = _inches(shape.width)
    h = _inches(shape.height)
    if w <= 0 or h <= 0:
        return None
    font_pt = _font_size_pt(shape)
    return {
        "slide": slide_num,
        "shape_id": f"shape-{shape_idx}",
        "name": getattr(shape, "name", f"shape-{shape_idx}"),
        "text": text,
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "font_pt": font_pt,
        "right": x + w,
        "bottom": y + h,
    }


def _line_capacity(width_in: float, font_pt: float) -> int:
    # Deliberately conservative. PowerPoint layout varies by font; this errs
    # toward finding risky wraps, not guaranteeing exact line breaks.
    char_w = max(0.055, (font_pt / 72.0) * 0.52)
    usable = max(0.25, width_in - 0.12)
    return max(5, int(usable / char_w))


def _wrap_words(text: str, width_in: float, font_pt: float) -> list[str]:
    capacity = _line_capacity(width_in, font_pt)
    lines: list[str] = []
    for raw_paragraph in text.splitlines() or [text]:
        words = raw_paragraph.strip().split()
        if not words:
            continue
        line = words[0]
        for word in words[1:]:
            if len(line) + 1 + len(word) <= capacity:
                line += " " + word
            else:
                lines.append(line)
                line = word
        lines.append(line)
    return lines


def _estimated_text_height(lines: list[str], font_pt: float) -> float:
    if not lines:
        return 0.0
    return len(lines) * (font_pt / 72.0) * 1.28


def _issue(
    slide: int,
    issue_type: str,
    severity: str,
    message: str,
    *,
    shape_id: str | None = None,
    suggestion: str | None = None,
    extra: dict[str, Any] | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "slide": slide,
        "type": issue_type,
        "severity": severity,
        "message": message,
    }
    if shape_id:
        payload["shape_id"] = shape_id
    if suggestion:
        payload["suggestion"] = suggestion
    if extra:
        payload.update(extra)
    return payload


def _analyze_text_shapes(prs: Presentation) -> list[dict[str, Any]]:
    slide_w = _inches(prs.slide_width)
    slide_h = _inches(prs.slide_height)
    issues: list[dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        records: list[dict[str, Any]] = []
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            record = _shape_record(slide_num, shape_idx, shape)
            if record:
                records.append(record)

        for record in records:
            text = str(record["text"])
            font_pt = float(record["font_pt"])
            x = float(record["x"])
            y = float(record["y"])
            w = float(record["w"])
            h = float(record["h"])
            right = float(record["right"])
            bottom = float(record["bottom"])
            lines = _wrap_words(text, w, font_pt)
            estimated_h = _estimated_text_height(lines, font_pt)
            is_title_like = y <= 1.35 and font_pt >= 22
            is_footer_like = y >= slide_h - 0.75 and font_pt <= 13
            is_caption_like = font_pt <= 12.5

            if not is_footer_like:
                if is_caption_like and font_pt < 8.0 and (
                    is_footer_like or (y > 1.40 and len(text) >= 28)
                ):
                    issues.append(
                        _issue(
                            slide_num,
                            "caption_font_too_small",
                            "warning",
                            f"Caption-like text is {font_pt:.1f}pt.",
                            shape_id=str(record["shape_id"]),
                            suggestion="Use at least 8pt for captions and source lines.",
                            extra={"font_pt": round(font_pt, 1)},
                        )
                    )
                elif not is_title_like and not is_footer_like and y > 1.35 and font_pt < 10.0:
                    issues.append(
                        _issue(
                            slide_num,
                            "body_font_small",
                            "info",
                            f"Body-like text is {font_pt:.1f}pt.",
                            shape_id=str(record["shape_id"]),
                            suggestion="Check readability; use larger text or split dense content if this is not a table/caption.",
                            extra={"font_pt": round(font_pt, 1)},
                        )
                    )

            if is_title_like and len(lines) >= 3:
                issues.append(
                    _issue(
                        slide_num,
                        "title_wrap_risk",
                        "warning",
                        f"Title-like text is estimated at {len(lines)} lines.",
                        shape_id=str(record["shape_id"]),
                        suggestion="Shorten the title or move part of it into the subtitle.",
                        extra={"estimated_lines": len(lines), "font_pt": round(font_pt, 1)},
                    )
                )

            if not is_caption_like and len(lines) >= 2 and len(text) >= 48:
                last_words = lines[-1].split()
                previous_len = len(lines[-2]) if len(lines) >= 2 else 0
                last_len = len(lines[-1])
                if len(last_words) <= 2 and last_len <= 5 and previous_len >= 18:
                    issues.append(
                        _issue(
                            slide_num,
                            "orphan_word_risk",
                            "warning",
                            f"Last estimated line is short: \"{lines[-1]}\".",
                            shape_id=str(record["shape_id"]),
                            suggestion="Shorten the sentence, widen the text box, or rebalance line breaks.",
                            extra={"estimated_lines": lines[-4:], "font_pt": round(font_pt, 1)},
                        )
                    )

            fill_ratio = estimated_h / h if h > 0 else 0.0
            if fill_ratio >= 1.05:
                issues.append(
                    _issue(
                        slide_num,
                        "text_box_clip_risk",
                        "warning",
                        f"Estimated text height is {fill_ratio:.0%} of its box height.",
                        shape_id=str(record["shape_id"]),
                        suggestion="Increase the box height, reduce body copy, or reduce font size.",
                        extra={
                            "estimated_height_in": round(estimated_h, 2),
                            "box_height_in": round(h, 2),
                        },
                    )
                )
            elif fill_ratio >= 0.88 and not is_footer_like:
                issues.append(
                    _issue(
                        slide_num,
                        "text_box_tight",
                        "info",
                        f"Estimated text height uses {fill_ratio:.0%} of its box.",
                        shape_id=str(record["shape_id"]),
                        suggestion="Check the rendered slide for clipping or awkward wrap.",
                        extra={
                            "estimated_height_in": round(estimated_h, 2),
                            "box_height_in": round(h, 2),
                        },
                    )
                )

            if not is_footer_like and (
                x < SLIDE_SIDE_SAFE
                or right > slide_w - SLIDE_SIDE_SAFE
                or bottom > slide_h - SLIDE_BOTTOM_SAFE
            ):
                issues.append(
                    _issue(
                        slide_num,
                        "safe_area_risk",
                        "warning",
                        "Text is close to the slide edge or bottom safe area.",
                        shape_id=str(record["shape_id"]),
                        suggestion="Move the text into the grid safe area or reserve more footer space.",
                        extra={
                            "x": round(x, 2),
                            "right": round(right, 2),
                            "bottom": round(bottom, 2),
                        },
                    )
                )

        title_records = [
            r for r in records if float(r["y"]) <= 1.35 and float(r["font_pt"]) >= 22
        ]
        for title in title_records:
            title_lines = _wrap_words(str(title["text"]), float(title["w"]), float(title["font_pt"]))
            title_bottom = float(title["y"]) + _estimated_text_height(title_lines, float(title["font_pt"]))
            below = [
                r
                for r in records
                if r["shape_id"] != title["shape_id"]
                and float(r["y"]) > float(title["y"]) + 0.02
                and float(r["y"]) < slide_h - 0.75
            ]
            if not below:
                continue
            next_record = min(below, key=lambda r: float(r["y"]))
            gap = float(next_record["y"]) - title_bottom
            if gap < 0.10:
                issues.append(
                    _issue(
                        slide_num,
                        "title_clearance_risk",
                        "warning",
                        f"Estimated title clearance to the next text box is {gap:.2f}\".",
                        shape_id=str(title["shape_id"]),
                        suggestion="Shorten the title, reduce title font, or increase measured header spacing.",
                        extra={
                            "gap_in": round(gap, 2),
                            "next_shape_id": next_record["shape_id"],
                        },
                    )
                )

        footer_records = [r for r in records if float(r["y"]) >= slide_h - 0.78]
        body_records = [r for r in records if float(r["bottom"]) < slide_h - 0.55]
        if footer_records and body_records:
            footer_top = min(float(r["y"]) for r in footer_records)
            nearest_body = max(float(r["bottom"]) for r in body_records)
            gap = footer_top - nearest_body
            if gap < 0.16:
                issues.append(
                    _issue(
                        slide_num,
                        "footer_clearance_risk",
                        "warning",
                        f"Footer/content clearance is only {gap:.2f}\".",
                        suggestion="Move body content up, shorten the footer, or remove one content row.",
                        extra={"gap_in": round(gap, 2)},
                    )
                )

    return issues


def _outline_rhythm_issues(outline_path: Path | None) -> list[dict[str, Any]]:
    if not outline_path or not outline_path.exists():
        return [
            _issue(
                0,
                "layout_rhythm_unchecked",
                "info",
                "No outline.json was available, so variant/family rhythm checks were skipped.",
                suggestion="Pass --outline or keep outline.json next to the deck/workspace build.",
            )
        ]
    outline = _load_json(outline_path)
    slides = outline.get("slides")
    if not isinstance(slides, list):
        return []

    issues: list[dict[str, Any]] = []
    variants: list[tuple[int, str]] = []
    visual_anchors = 0
    image_anchors = 0
    icon_slides = 0
    family_runs: list[tuple[int, str]] = []
    card_family = {"cards-2", "cards-3", "timeline", "matrix", "stats"}
    deck_style = outline.get("deck_style") if isinstance(outline.get("deck_style"), dict) else {}
    compliance = outline.get("compliance") if isinstance(outline.get("compliance"), dict) else {}
    research_visual_mode = bool(
        (isinstance(deck_style, dict) and deck_style.get("research_visual_mode"))
        or (isinstance(compliance, dict) and compliance.get("require_attribution"))
    )

    def family_for(variant: str) -> str:
        if variant in {"cards-2", "cards-3", "matrix", "stats", "timeline"}:
            return "container-grid"
        if variant in {"table", "lab-run-results"}:
            return "table-evidence"
        if variant in {"image-sidebar", "scientific-figure", "generated-image"}:
            return "image-evidence"
        if variant in {"split", "comparison-2col"}:
            return "two-column"
        if variant in {"flow", "chart"}:
            return "diagram-data"
        if variant == "kpi-hero":
            return "hero-metric"
        return "open-body"

    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        if str(slide.get("type") or "content").strip().lower() != "content":
            continue
        variant = str(slide.get("variant") or "standard").strip().lower() or "standard"
        variants.append((idx, variant))
        family_runs.append((idx, family_for(variant)))
        assets = slide.get("assets") if isinstance(slide.get("assets"), dict) else {}
        figures = slide.get("figures")
        has_image_like_asset = isinstance(assets, dict) and any(
            assets.get(k)
            for k in ("hero_image", "generated_image", "image")
        )
        has_figures = isinstance(figures, list) and bool(figures)
        if variant in {"image-sidebar", "scientific-figure", "generated-image"} or has_image_like_asset or has_figures:
            image_anchors += 1
        if variant in {
            "chart",
            "table",
            "lab-run-results",
            "flow",
            "scientific-figure",
            "generated-image",
            "image-sidebar",
            "kpi-hero",
            "comparison-2col",
        } or (
            isinstance(assets, dict)
            and any(assets.get(k) for k in ("hero_image", "generated_image", "diagram", "mermaid_source", "image"))
        ) or has_figures:
            visual_anchors += 1
        if isinstance(assets, dict) and isinstance(assets.get("icons"), list) and assets.get("icons"):
            icon_slides += 1

    if not variants:
        return []

    counts = Counter(variant for _, variant in variants)
    for variant, count in sorted(counts.items()):
        if count >= 4 and len(variants) <= 10:
            issues.append(
                _issue(
                    0,
                    "variant_overuse",
                    "warning",
                    f"`{variant}` appears {count} times across {len(variants)} content slides.",
                    suggestion="Replace at least one repeated layout with a topic-specific rhythm breaker.",
                    extra={"variant": variant, "count": count},
                )
            )

    card_family_count = sum(1 for _, variant in variants if variant in card_family)
    if len(variants) >= 5 and card_family_count / len(variants) >= 0.70:
        issues.append(
            _issue(
                0,
                "visual_family_overuse",
                "warning",
                f"Card-family layouts appear on {card_family_count}/{len(variants)} content slides.",
                suggestion="Add a figure, table, chart, flow, comparison, or KPI rhythm breaker.",
                extra={"card_family_count": card_family_count, "content_slide_count": len(variants)},
            )
        )

    family_counts = Counter(family for _, family in family_runs)
    top_family, top_family_count = family_counts.most_common(1)[0]
    if (
        len(variants) >= 6
        and top_family_count / len(variants) >= 0.62
        and top_family not in {"table-evidence"}
    ):
        issues.append(
            _issue(
                0,
                "composition_family_repetition",
                "warning",
                f"`{top_family}` compositions account for {top_family_count}/{len(variants)} content slides.",
                suggestion=(
                    "Change at least one slide's visual role before rendering: "
                    "source-backed image, figure/table evidence, comparison, flow, "
                    "or a simple open-body slide."
                ),
                extra={"family": top_family, "count": top_family_count},
            )
        )

    if len(variants) >= 4 and visual_anchors == 0:
        issues.append(
            _issue(
                0,
                "visual_anchor_absent",
                "warning",
                "No content slide uses a figure, table, chart, diagram, image, or KPI anchor.",
                suggestion="Add at least one topic-specific visual anchor; for lab/data decks prefer image-sidebar, table, chart, or flow.",
            )
        )

    if research_visual_mode and image_anchors == 0 and len(variants) >= 4:
        issues.append(
            _issue(
                0,
                "research_visual_mode_without_images",
                "warning",
                "research_visual_mode is enabled, but no slide uses a source-backed image or figure.",
                suggestion=(
                    "Run build_workspace.py with --plan-research-assets "
                    "--allow-network-assets, or add assets.image / figures entries "
                    "that reference staged aliases from asset_plan.json."
                ),
            )
        )
    elif len(variants) >= 7 and image_anchors == 0:
        issues.append(
            _issue(
                0,
                "source_image_absent",
                "info",
                "No content slide uses a photographic/source-backed image or figure.",
                suggestion=(
                    "For public-topic decks, consider one image-sidebar or "
                    "scientific-figure slide with attribution rather than relying "
                    "only on cards, tables, or synthetic charts."
                ),
            )
        )

    if len(variants) >= 5 and icon_slides == 0 and card_family_count >= 3:
        issues.append(
            _issue(
                0,
                "icon_system_absent",
                "info",
                "Several card-family slides are present but none use assets.icons.",
                suggestion="Use react-icons slugs on conceptual cards/timelines when icons would clarify categories.",
            )
        )

    run_variant = ""
    run_slides: list[int] = []
    for slide_num, variant in variants + [(-1, "__sentinel__")]:
        if variant == run_variant:
            run_slides.append(slide_num)
            continue
        if len(run_slides) >= 3:
            issues.append(
                _issue(
                    run_slides[0],
                    "variant_run",
                    "warning",
                    f"`{run_variant}` repeats for {len(run_slides)} consecutive content slides.",
                    suggestion="Break the run with a KPI, timeline, comparison, image, or section slide.",
                    extra={"variant": run_variant, "slides": run_slides},
                )
            )
        run_variant = variant
        run_slides = [slide_num]

    run_family = ""
    run_family_slides: list[int] = []
    for slide_num, family in family_runs + [(-1, "__sentinel__")]:
        if family == run_family:
            run_family_slides.append(slide_num)
            continue
        if len(run_family_slides) >= 4 and run_family != "table-evidence":
            issues.append(
                _issue(
                    run_family_slides[0],
                    "composition_family_run",
                    "warning",
                    f"`{run_family}` composition repeats for {len(run_family_slides)} consecutive content slides.",
                    suggestion="Break the run with a different evidence format or a simpler open-body slide.",
                    extra={"family": run_family, "slides": run_family_slides},
                )
            )
        run_family = family
        run_family_slides = [slide_num]

    return issues


def _discover_outline(pptx_path: Path) -> Path | None:
    candidates = [
        pptx_path.with_name("outline.json"),
        pptx_path.parent / "outline.json",
        pptx_path.parent.parent / "outline.json",
    ]
    if pptx_path.parent.name == "build":
        candidates.append(pptx_path.parent.parent / "outline.json")
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    return None


def _render_if_needed(pptx_path: Path, renders_dir: Path) -> tuple[bool, str]:
    existing = _rendered_paths(renders_dir)
    if existing:
        return True, "existing"

    base = Path(__file__).resolve().parent
    cmd = [
        sys.executable,
        str(base / "render_slides.py"),
        "--input",
        str(pptx_path),
        "--outdir",
        str(renders_dir),
        "--dpi",
        "150",
        "--format",
        "jpeg",
    ]
    rc, output = _run(cmd)
    return rc == 0 and bool(_rendered_paths(renders_dir)), output


def _rendered_paths(renders_dir: Path) -> list[Path]:
    if not renders_dir.exists():
        return []
    paths = (
        list(renders_dir.glob("slide-*.jpg"))
        + list(renders_dir.glob("slide-*.jpeg"))
        + list(renders_dir.glob("slide-*.png"))
    )
    return sorted(paths, key=_slide_sort_key)


def _slide_sort_key(path: Path) -> tuple[int, str]:
    match = re.search(r"slide-(\d+)", path.stem)
    return (int(match.group(1)) if match else 10**9, path.name)


def _issue_badges_by_slide(issues: list[dict[str, Any]]) -> dict[int, tuple[int, int]]:
    badges: dict[int, tuple[int, int]] = {}
    for issue in issues:
        slide_raw = issue.get("slide")
        if not isinstance(slide_raw, int) or slide_raw <= 0:
            continue
        warning, info = badges.get(slide_raw, (0, 0))
        if issue.get("severity") == "warning":
            warning += 1
        elif issue.get("severity") == "info":
            info += 1
        badges[slide_raw] = (warning, info)
    return badges


def _make_contact_sheet(rendered_paths: list[Path], output_path: Path, issues: list[dict[str, Any]] | None = None) -> str | None:
    if Image is None or ImageDraw is None or ImageFont is None:
        return None
    if not rendered_paths:
        return None

    thumbs: list[Image.Image] = []
    target_w = 420
    label_h = 32
    pad = 18
    badges = _issue_badges_by_slide(issues or [])
    for path in rendered_paths:
        with Image.open(path) as img:
            img = img.convert("RGB")
            ratio = img.height / img.width
            thumb = img.resize((target_w, max(1, int(target_w * ratio))))
            tile = Image.new("RGB", (target_w, thumb.height + label_h), (246, 248, 251))
            tile.paste(thumb, (0, label_h))
            draw = ImageDraw.Draw(tile)
            draw.text((8, 8), path.stem.replace("-", " "), fill=(15, 23, 42))
            slide_num = _slide_sort_key(path)[0]
            warning_count, _info_count = badges.get(slide_num, (0, 0))
            if warning_count:
                badge = f"W{warning_count}"
                fill = (185, 28, 28)
                x0 = target_w - 58
                y0 = 5
                draw.rounded_rectangle([x0, y0, target_w - 8, label_h - 5], radius=7, fill=fill)
                draw.text((x0 + 8, y0 + 4), badge, fill=(255, 255, 255))
            thumbs.append(tile)

    cols = min(3, max(1, math.ceil(math.sqrt(len(thumbs)))))
    rows = math.ceil(len(thumbs) / cols)
    tile_w = target_w
    tile_h = max(tile.height for tile in thumbs)
    sheet = Image.new(
        "RGB",
        (cols * tile_w + (cols + 1) * pad, rows * tile_h + (rows + 1) * pad),
        (226, 232, 240),
    )
    for idx, tile in enumerate(thumbs):
        row = idx // cols
        col = idx % cols
        x = pad + col * (tile_w + pad)
        y = pad + row * (tile_h + pad)
        sheet.paste(tile, (x, y))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path, quality=92)
    return str(output_path)


def _markdown_report(
    *,
    pptx_path: Path,
    rendered_paths: list[Path],
    contact_sheet: str | None,
    issues: list[dict[str, Any]],
) -> str:
    warnings = [item for item in issues if item.get("severity") == "warning"]
    infos = [item for item in issues if item.get("severity") == "info"]
    lines = [
        "# Visual Review Packet",
        "",
        f"- PPTX: `{pptx_path}`",
        f"- Rendered slides: {len(rendered_paths)}",
        f"- Warnings: {len(warnings)}",
        f"- Info: {len(infos)}",
    ]
    if contact_sheet:
        lines.append(f"- Contact sheet: `{contact_sheet}`")
    lines.extend(["", "## Findings", ""])
    if not issues:
        lines.append("No deterministic visual-review issues found. Still inspect the rendered slides.")
    else:
        for item in issues:
            slide = item.get("slide")
            location = "deck" if slide == 0 else f"slide {slide}"
            shape = f" / {item['shape_id']}" if item.get("shape_id") else ""
            lines.append(
                f"- [{str(item.get('severity', 'info')).upper()}] {location}{shape}: "
                f"{item.get('type')} - {item.get('message')}"
            )
            if item.get("suggestion"):
                lines.append(f"  Suggestion: {item['suggestion']}")
    lines.extend(["", "## Rendered Slides", ""])
    for idx, path in enumerate(rendered_paths, start=1):
        lines.append(f"{idx}. `{path}`")
    lines.extend(
        [
            "",
            "## Human/Agent Review Checklist",
            "",
            "- Inspect the contact sheet first for rhythm, density, and repeated layouts.",
            "- Open individual slides for title wrapping, orphan words, weak contrast, and footer clearance.",
            "- Fix `outline.json`, `design_brief.json`, assets, or renderer code, then rebuild.",
            "- Do not patch the generated `.pptx` directly unless the task is an explicit one-off edit.",
        ]
    )
    return "\n".join(lines) + "\n"


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create a rendered PPTX visual-review packet.")
    parser.add_argument("--input", required=True, help="Input .pptx file")
    parser.add_argument("--outdir", required=True, help="Directory for visual review artifacts")
    parser.add_argument("--outline", help="Optional outline.json for layout-rhythm checks")
    parser.add_argument(
        "--renders-dir",
        help="Optional existing renders directory. If omitted, slides are rendered into <outdir>/renders.",
    )
    parser.add_argument("--report", help="JSON report path (default: <outdir>/visual_review.json)")
    parser.add_argument("--markdown", help="Markdown report path (default: <outdir>/visual_review.md)")
    parser.add_argument(
        "--skip-render",
        action="store_true",
        help="Do not render if --renders-dir has no slide images; run text/rhythm checks only.",
    )
    parser.add_argument(
        "--fail-on-warnings",
        action="store_true",
        help="Exit non-zero if warning-level findings are present.",
    )
    return parser.parse_args()


def main() -> int:
    args = _args()
    pptx_path = Path(args.input).expanduser().resolve()
    outdir = Path(args.outdir).expanduser().resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"Input file not found: {pptx_path}")
    outdir.mkdir(parents=True, exist_ok=True)

    renders_dir = (
        Path(args.renders_dir).expanduser().resolve()
        if args.renders_dir
        else outdir / "renders"
    )
    renders_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    issues = _analyze_text_shapes(prs)
    outline_path = Path(args.outline).expanduser().resolve() if args.outline else _discover_outline(pptx_path)
    issues.extend(_outline_rhythm_issues(outline_path))

    render_status = "skipped"
    if not args.skip_render:
        rendered, render_status = _render_if_needed(pptx_path, renders_dir)
        if not rendered:
            issues.append(
                _issue(
                    0,
                    "render_unavailable",
                    "warning",
                    "Slides could not be rendered for visual contact-sheet review.",
                    suggestion="Install LibreOffice/Poppler or rerun with a valid --renders-dir.",
                    extra={"render_output": render_status[-1200:]},
                )
            )

    rendered_paths = _rendered_paths(renders_dir)
    contact_sheet = _make_contact_sheet(rendered_paths, outdir / "contact_sheet.jpg", issues)

    warning_count = sum(1 for item in issues if item.get("severity") == "warning")
    info_count = sum(1 for item in issues if item.get("severity") == "info")
    payload = {
        "input": str(pptx_path),
        "outline": str(outline_path) if outline_path else "",
        "outdir": str(outdir),
        "renders_dir": str(renders_dir),
        "render_status": render_status,
        "rendered_slide_count": len(rendered_paths),
        "contact_sheet": contact_sheet or "",
        "warning_count": warning_count,
        "info_count": info_count,
        "issue_count": len(issues),
        "issues": issues,
        "passed": warning_count == 0,
    }

    report_path = Path(args.report).expanduser().resolve() if args.report else outdir / "visual_review.json"
    markdown_path = Path(args.markdown).expanduser().resolve() if args.markdown else outdir / "visual_review.md"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    markdown_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
    markdown_path.write_text(
        _markdown_report(
            pptx_path=pptx_path,
            rendered_paths=rendered_paths,
            contact_sheet=contact_sheet,
            issues=issues,
        ),
        encoding="utf-8",
    )

    # If renders were supplied from a QA directory, keep them in place. If this
    # script rendered into its own outdir, the paths are already there.
    if args.renders_dir and rendered_paths:
        review_renders = outdir / "renders"
        review_renders.mkdir(parents=True, exist_ok=True)
        if review_renders.resolve() != renders_dir.resolve():
            for path in rendered_paths:
                shutil.copy2(path, review_renders / path.name)

    print(f"Visual review artifacts: {outdir}")
    print(f"Rendered slides: {len(rendered_paths)}")
    print(f"Warnings: {warning_count}")
    print(f"Info: {info_count}")
    if contact_sheet:
        print(f"Contact sheet: {contact_sheet}")
    print(f"Visual review report: {report_path}")
    print(f"Visual review markdown: {markdown_path}")

    if args.fail_on_warnings and warning_count:
        return 1
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
