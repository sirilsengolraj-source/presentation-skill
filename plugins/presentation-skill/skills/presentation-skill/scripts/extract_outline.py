#!/usr/bin/env python3
"""Extract slide titles, body bullets, and notes from a PPTX file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation


def _title_text(slide: Any) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            return (title_shape.text or "").strip()
    except Exception:
        pass
    return ""


def _body_lines(slide: Any, title_text: str) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        shape_text = (shape.text or "").strip()
        if shape_text and title_text and shape_text == title_text:
            continue
        frame = shape.text_frame
        for paragraph in frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            rows.append({"text": text, "level": int(getattr(paragraph, "level", 0))})
    return rows


def _notes_text(slide: Any) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    lines: list[str] = []
    for paragraph in notes_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return "\n".join(lines)


def _to_markdown(data: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append("# Deck Outline")
    lines.append("")
    lines.append(f"- File: {data['file']}")
    lines.append(f"- Slides: {data['slide_count']}")
    lines.append("")

    for slide in data["slides"]:
        title = slide["title"] or "(No title)"
        lines.append(f"## Slide {slide['index']}: {title}")
        lines.append("")

        if slide["body"]:
            lines.append("Body:")
            for item in slide["body"]:
                prefix = "  " * int(item["level"])
                lines.append(f"- {prefix}{item['text']}")
        else:
            lines.append("Body: (empty)")

        if slide["notes"]:
            lines.append("")
            lines.append("Notes:")
            for note_line in slide["notes"].splitlines():
                lines.append(f"- {note_line}")

        lines.append("")

    return "\n".join(lines).strip() + "\n"


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract an outline from a .pptx file.")
    parser.add_argument("--input", required=True, help="Input .pptx path")
    parser.add_argument(
        "--format",
        choices=["markdown", "json"],
        default="markdown",
        help="Output format",
    )
    parser.add_argument("--output", help="Optional output file path")
    return parser.parse_args()


def main() -> int:
    args = _args()
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    presentation = Presentation(str(input_path))
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = _title_text(slide)
        slides.append(
            {
                "index": idx,
                "title": title,
                "body": _body_lines(slide, title),
                "notes": _notes_text(slide),
            }
        )

    payload = {
        "file": str(input_path),
        "slide_count": len(slides),
        "slides": slides,
    }

    if args.format == "json":
        content = json.dumps(payload, indent=2, ensure_ascii=False) + "\n"
    else:
        content = _to_markdown(payload)

    if args.output:
        output_path = Path(args.output).expanduser().resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(content, encoding="utf-8")
    else:
        print(content, end="")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
