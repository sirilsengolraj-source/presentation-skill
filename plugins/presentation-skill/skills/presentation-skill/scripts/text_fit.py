#!/usr/bin/env python3
"""Deterministic text-fit remediation for PPTX decks."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from design_tokens import get_style_preset
from inventory import InventoryData, extract_text_inventory


def _shape_text_len(shape_data: Any) -> int:
    total = 0
    for paragraph in shape_data.paragraphs:
        total += len((paragraph.text or "").strip())
    return total


def _is_heading_shape(shape_data: Any) -> bool:
    font_sizes = [
        paragraph.font_size
        for paragraph in shape_data.paragraphs
        if paragraph.font_size is not None
    ]
    max_font = max(font_sizes) if font_sizes else 0.0
    return max_font >= 18 or shape_data.height <= 0.90


def _line_budget(shape_data: Any, heading: bool) -> int:
    chars_per_line = max(12, int(shape_data.width * (9 if heading else 12)))
    max_lines = 2 if heading else 4
    return chars_per_line * max_lines


def _find_overflow(inventory: InventoryData, tolerance: float) -> list[tuple[str, str, Any]]:
    targets: list[tuple[str, str, Any]] = []
    for slide_key, shapes in inventory.items():
        for shape_key, shape_data in shapes.items():
            overflow = shape_data.frame_overflow_bottom
            if overflow is not None and overflow > tolerance:
                targets.append((slide_key, shape_key, shape_data))
    return targets


def _find_line_budget_violations(
    inventory: InventoryData,
) -> list[tuple[str, str, Any, int]]:
    targets: list[tuple[str, str, Any, int]] = []
    for slide_key, shapes in inventory.items():
        for shape_key, shape_data in shapes.items():
            text_len = _shape_text_len(shape_data)
            if text_len == 0:
                continue
            heading = _is_heading_shape(shape_data)
            budget = _line_budget(shape_data, heading)
            if text_len > budget:
                targets.append((slide_key, shape_key, shape_data, text_len - budget))
    return targets


def _find_slide_overflow(
    inventory: InventoryData,
    tolerance: float,
) -> list[tuple[str, str, Any, float, float]]:
    targets: list[tuple[str, str, Any, float, float]] = []
    for slide_key, shapes in inventory.items():
        for shape_key, shape_data in shapes.items():
            right = float(shape_data.slide_overflow_right or 0.0)
            bottom = float(shape_data.slide_overflow_bottom or 0.0)
            if right > tolerance or bottom > tolerance:
                targets.append((slide_key, shape_key, shape_data, right, bottom))
    return targets


def _reduce_font_once(
    shape_data: Any,
    *,
    heading_min: int,
    body_min: int,
    heading_default: int,
    body_default: int,
) -> tuple[bool, list[dict[str, Any]]]:
    shape = shape_data.shape
    if shape is None or not hasattr(shape, "text_frame"):
        return False, []
    heading = _is_heading_shape(shape_data)
    min_size = heading_min if heading else body_min
    default_size = heading_default if heading else body_default
    changes: list[dict[str, Any]] = []
    changed = False

    for p_index, paragraph in enumerate(shape.text_frame.paragraphs):
        if paragraph.runs:
            for r_index, run in enumerate(paragraph.runs):
                current = run.font.size.pt if run.font.size else float(default_size)
                if current <= min_size:
                    continue
                target = max(float(min_size), current - 1.0)
                if target < current:
                    run.font.size = Pt(target)
                    changed = True
                    changes.append(
                        {
                            "paragraph_index": p_index,
                            "run_index": r_index,
                            "from_pt": round(current, 2),
                            "to_pt": round(target, 2),
                        }
                    )
        else:
            current = (
                paragraph.font.size.pt
                if paragraph.font is not None and paragraph.font.size
                else float(default_size)
            )
            if current > min_size:
                target = max(float(min_size), current - 1.0)
                if target < current:
                    paragraph.font.size = Pt(target)
                    changed = True
                    changes.append(
                        {
                            "paragraph_index": p_index,
                            "run_index": None,
                            "from_pt": round(current, 2),
                            "to_pt": round(target, 2),
                        }
                    )
    return changed, changes


def _grow_shape(shape_data: Any, amount: float) -> bool:
    shape = shape_data.shape
    if shape is None:
        return False
    if amount <= 0:
        return False
    top_in = shape.top / 914400.0
    height_in = shape.height / 914400.0
    slide_h_emu = getattr(shape_data, "slide_height_emu", None)
    if slide_h_emu:
        slide_h_in = slide_h_emu / 914400.0
        # Avoid expanding footer-like shapes into slide overflow.
        if top_in > slide_h_in - 0.55 and height_in <= 0.5:
            return False
        max_height = max(0.05, slide_h_in - top_in - 0.02)
        if height_in >= max_height:
            return False
        amount = min(amount, max_height - height_in)
        if amount <= 0:
            return False
    shape.height += Inches(amount)
    return True


def _to_inches(emu: int) -> float:
    return emu / 914400.0


def _to_emu(inches: float) -> int:
    return int(round(inches * 914400.0))


def _shape_bounds(shape: Any) -> tuple[float, float, float, float]:
    return (
        _to_inches(shape.left),
        _to_inches(shape.top),
        _to_inches(shape.width),
        _to_inches(shape.height),
    )


def _is_card(shape: Any) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    _, _, width, height = _shape_bounds(shape)
    area = width * height
    if width <= 0.05 or height <= 0.05:
        return False
    return width >= 1.4 and height >= 0.7 and area >= 1.0


def _is_rail(shape: Any) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    _, _, width, height = _shape_bounds(shape)
    return width >= 1.2 and height <= 0.16


def _overlap_width(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> float:
    a_left, _, a_width, _ = a
    b_left, _, b_width, _ = b
    left = max(a_left, b_left)
    right = min(a_left + a_width, b_left + b_width)
    return max(0.0, right - left)


def _group_cards_by_top(cards: list[dict[str, Any]], tolerance: float) -> list[list[dict[str, Any]]]:
    groups: list[list[dict[str, Any]]] = []
    for item in sorted(cards, key=lambda entry: entry["top"]):
        placed = False
        for group in groups:
            seed = sum(entry["top"] for entry in group) / len(group)
            if abs(item["top"] - seed) <= tolerance:
                group.append(item)
                placed = True
                break
        if not placed:
            groups.append([item])
    return groups


def _apply_density_autofix(
    prs: Presentation,
    density_target: float,
    min_card_height: float = 1.35,
    max_shrink_fraction: float = 0.08,
) -> list[dict[str, Any]]:
    """Shrink card-like shapes on slides that exceed the density target.

    Density-too-high warnings come from cards/quadrants filling the content
    zone by design. When QA still flags a slide after variant-aware caps,
    this pass reduces card heights by up to `max_shrink_fraction`, floored
    at `min_card_height`. Only card shapes are shrunk (not rails, footers,
    text frames, images).

    Returns a list of change records describing the shrinks applied.
    """
    changes: list[dict[str, Any]] = []
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height
    slide_w = _to_inches(slide_w_emu)
    slide_h = _to_inches(slide_h_emu)
    slide_area = max(0.01, slide_w * slide_h)

    for slide_index, slide in enumerate(prs.slides, start=1):
        cards: list[dict[str, Any]] = []
        covered_area = 0.0
        for shape_index, shape in enumerate(slide.shapes):
            left, top, width, height = _shape_bounds(shape)
            area = width * height
            # Match layout_lint's content-shape filter.
            if area < 0.15:
                continue
            if top <= 0.10:
                continue
            covered_area += area
            if _is_card(shape) and height > min_card_height:
                cards.append(
                    {
                        "shape": shape,
                        "shape_id": f"shape-{shape_index}",
                        "top": top,
                        "height": height,
                    }
                )

        density = min(1.0, covered_area / slide_area)
        if density <= density_target:
            continue
        if not cards:
            continue

        # Compute shrink factor needed to bring density to target, capped
        # at max_shrink_fraction. We only shrink cards, so the recovered
        # area is shrink_fraction * sum(card heights * card widths).
        # Approximation: shrink each card by the same fraction.
        over = density - density_target
        needed_area_drop = over * slide_area
        cards_area = sum(
            _shape_bounds(c["shape"])[2] * c["height"] for c in cards
        )
        if cards_area <= 0.01:
            continue
        shrink_fraction = min(max_shrink_fraction, needed_area_drop / cards_area)
        if shrink_fraction <= 0.005:
            continue

        for card in cards:
            new_height = card["height"] * (1.0 - shrink_fraction)
            if new_height < min_card_height:
                new_height = min_card_height
            if abs(new_height - card["height"]) < 0.01:
                continue
            card["shape"].height = _to_emu(new_height)
            changes.append(
                {
                    "slide": f"slide-{slide_index}",
                    "shape": card["shape_id"],
                    "action": "density_shrink_card",
                    "details": {
                        "from_height": round(card["height"], 4),
                        "to_height": round(new_height, 4),
                        "shrink_fraction": round(shrink_fraction, 4),
                        "density_before": round(density, 4),
                    },
                }
            )

    return changes


def _apply_geometry_autofix(prs: Presentation, edge_tolerance: float) -> list[dict[str, Any]]:
    changes: list[dict[str, Any]] = []
    row_tolerance = 0.16
    for slide_index, slide in enumerate(prs.slides, start=1):
        cards: list[dict[str, Any]] = []
        rails: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes):
            left, top, width, height = _shape_bounds(shape)
            item = {
                "shape": shape,
                "shape_id": f"shape-{shape_index}",
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            }
            if _is_card(shape):
                cards.append(item)
            elif _is_rail(shape):
                rails.append(item)

        # Fix accent rails to exactly match target cards.
        for rail in rails:
            candidate = None
            best_score = 0.0
            rail_bounds = (rail["left"], rail["top"], rail["width"], rail["height"])
            for card in cards:
                if card["top"] < rail["top"] or card["top"] - rail["top"] > 0.55:
                    continue
                card_bounds = (card["left"], card["top"], card["width"], card["height"])
                overlap = _overlap_width(rail_bounds, card_bounds)
                if overlap > best_score:
                    best_score = overlap
                    candidate = card
            if candidate is None:
                continue

            left_delta = abs(rail["left"] - candidate["left"])
            width_delta = abs(rail["width"] - candidate["width"])
            if left_delta > edge_tolerance or width_delta > edge_tolerance:
                rail["shape"].left = _to_emu(candidate["left"])
                rail["shape"].width = _to_emu(candidate["width"])
                changes.append(
                    {
                        "slide": f"slide-{slide_index}",
                        "shape": rail["shape_id"],
                        "action": "rail_align_to_card",
                        "details": {
                            "target_shape": candidate["shape_id"],
                            "left_delta": round(left_delta, 4),
                            "width_delta": round(width_delta, 4),
                        },
                    }
                )

        # Fix row top/height drift across cards.
        for group in _group_cards_by_top(cards, tolerance=row_tolerance):
            if len(group) < 2:
                continue
            top_values = [item["top"] for item in group]
            height_values = [item["height"] for item in group]
            top_spread = max(top_values) - min(top_values)
            height_spread = max(height_values) - min(height_values)

            if top_spread > edge_tolerance:
                target_top = min(top_values)
                for item in group:
                    if abs(item["top"] - target_top) <= edge_tolerance:
                        continue
                    item["shape"].top = _to_emu(target_top)
                    changes.append(
                        {
                            "slide": f"slide-{slide_index}",
                            "shape": item["shape_id"],
                            "action": "row_top_align",
                            "details": {
                                "from_top": round(item["top"], 4),
                                "to_top": round(target_top, 4),
                            },
                        }
                    )

            if height_spread > edge_tolerance:
                target_height = max(height_values)
                for item in group:
                    if abs(item["height"] - target_height) <= edge_tolerance:
                        continue
                    item["shape"].height = _to_emu(target_height)
                    changes.append(
                        {
                            "slide": f"slide-{slide_index}",
                            "shape": item["shape_id"],
                            "action": "row_height_normalize",
                            "details": {
                                "from_height": round(item["height"], 4),
                                "to_height": round(target_height, 4),
                            },
                        }
                    )
    return changes


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Apply deterministic text-fit fixes.")
    parser.add_argument("--input", required=True, help="Input .pptx path")
    parser.add_argument("--output", help="Output .pptx path (defaults to input in-place)")
    parser.add_argument(
        "--style-preset",
        default="executive-clinical",
        help="Style preset for font and rhythm bounds",
    )
    parser.add_argument(
        "--max-font-iterations",
        type=int,
        default=6,
        help="Maximum font-reduction passes",
    )
    parser.add_argument(
        "--max-grow-steps",
        type=int,
        default=3,
        help="Maximum box-growth passes for unresolved overflow",
    )
    parser.add_argument(
        "--overflow-tolerance",
        type=float,
        default=0.01,
        help="Overflow threshold (inches) to trigger remediation",
    )
    parser.add_argument(
        "--disable-geometry-fix",
        action="store_true",
        help="Disable deterministic geometry autofix pass",
    )
    parser.add_argument(
        "--disable-density-fix",
        action="store_true",
        help="Disable density autofix pass (card-shrink for density_too_high)",
    )
    parser.add_argument(
        "--density-target",
        type=float,
        default=0.92,
        help="Trigger density autofix when slide density exceeds this (default: 0.92)",
    )
    parser.add_argument(
        "--density-min-card-height",
        type=float,
        default=1.35,
        help="Minimum card height after density autofix (default: 1.35 inches)",
    )
    parser.add_argument(
        "--density-max-shrink",
        type=float,
        default=0.08,
        help="Maximum per-card shrink fraction from density autofix (default: 0.08)",
    )
    parser.add_argument(
        "--density-only",
        action="store_true",
        help=(
            "Skip font-reduce, shape-grow, overflow-reposition, and "
            "geometry-autofix passes. Only the density autofix runs. "
            "Useful when QA is clean except for density warnings."
        ),
    )
    parser.add_argument("--report", help="Optional output JSON report path")
    return parser.parse_args()


def main() -> int:
    args = _args()
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Input not found: {input_path}")
    output_path = (
        Path(args.output).expanduser().resolve() if args.output else input_path
    )
    preset = get_style_preset(args.style_preset)
    typography = preset.typography
    layout = preset.layout

    prs = Presentation(str(input_path))
    report_changes: list[dict[str, Any]] = []
    density_only = bool(args.density_only)

    # Phase 1: bounded font reduction on overflows and line-budget violations.
    loops_executed = 0
    for iteration in range(1, args.max_font_iterations + 1):
        if density_only:
            break
        loops_executed = iteration
        inventory = extract_text_inventory(input_path, prs)
        overflow_targets = _find_overflow(inventory, args.overflow_tolerance)
        budget_targets = _find_line_budget_violations(inventory)
        if not overflow_targets and not budget_targets:
            break

        changed_any = False
        target_map: dict[tuple[str, str], Any] = {
            (slide_key, shape_key): shape_data
            for slide_key, shape_key, shape_data in overflow_targets
        }
        for slide_key, shape_key, shape_data, _ in budget_targets:
            target_map.setdefault((slide_key, shape_key), shape_data)

        for (slide_key, shape_key), shape_data in target_map.items():
            changed, changes = _reduce_font_once(
                shape_data,
                heading_min=typography.section_min,
                body_min=typography.body_min,
                heading_default=typography.section_max,
                body_default=typography.body_max,
            )
            if changed:
                changed_any = True
                report_changes.append(
                    {
                        "slide": slide_key,
                        "shape": shape_key,
                        "action": "font_reduce",
                        "details": changes,
                    }
                )

        if not changed_any:
            break

    # Phase 2: grow unresolved overflow shapes by rhythm increments.
    if density_only:
        args.max_grow_steps = 0
    for _ in range(args.max_grow_steps):
        inventory = extract_text_inventory(input_path, prs)
        overflow_targets = _find_overflow(inventory, args.overflow_tolerance)
        if not overflow_targets:
            break
        changed_any = False
        for slide_key, shape_key, shape_data in overflow_targets:
            overflow = float(shape_data.frame_overflow_bottom or 0.0)
            growth = min(layout.rhythm, max(layout.rhythm / 2.0, overflow + 0.03))
            grew = _grow_shape(shape_data, growth)
            if grew:
                changed_any = True
                report_changes.append(
                    {
                        "slide": slide_key,
                        "shape": shape_key,
                        "action": "height_grow",
                        "details": {"delta_inches": round(growth, 4)},
                    }
                )
        if not changed_any:
            break

    # Phase 3: reposition shapes that exceed slide boundaries.
    slide_overflow_targets = []
    if not density_only:
        inventory = extract_text_inventory(input_path, prs)
        slide_overflow_targets = _find_slide_overflow(inventory, args.overflow_tolerance)
    for slide_key, shape_key, shape_data, overflow_right, overflow_bottom in slide_overflow_targets:
        shape = shape_data.shape
        if shape is None:
            continue
        moved = False
        details: dict[str, float] = {}
        if overflow_right > 0:
            shift_x = overflow_right + 0.02
            shape.left -= Inches(shift_x)
            details["shift_left_inches"] = round(shift_x, 4)
            moved = True
        if overflow_bottom > 0:
            shift_y = overflow_bottom + 0.02
            shape.top -= Inches(shift_y)
            details["shift_up_inches"] = round(shift_y, 4)
            moved = True
        if moved:
            report_changes.append(
                {
                    "slide": slide_key,
                    "shape": shape_key,
                    "action": "slide_overflow_reposition",
                    "details": details,
                }
            )

    geometry_changes: list[dict[str, Any]] = []
    if not args.disable_geometry_fix and not density_only:
        geometry_changes = _apply_geometry_autofix(prs, edge_tolerance=layout.edge_tolerance)
        report_changes.extend(geometry_changes)

    # Phase 4: density autofix. Targets slides that still exceed the
    # density cap after variant-aware bumping in layout_lint. Shrinks
    # card heights by up to 8% with a min-height floor so it can't
    # create empty_ratio_too_high violations.
    density_changes: list[dict[str, Any]] = []
    if not args.disable_density_fix:
        density_changes = _apply_density_autofix(
            prs,
            density_target=args.density_target,
            min_card_height=args.density_min_card_height,
            max_shrink_fraction=args.density_max_shrink,
        )
        report_changes.extend(density_changes)

    if density_only:
        remaining = []
    else:
        final_inventory = extract_text_inventory(input_path, prs)
        remaining = [
            {
                "slide": slide_key,
                "shape": shape_key,
                "overflow_inches": round(float(shape_data.frame_overflow_bottom or 0.0), 4),
            }
            for slide_key, shape_key, shape_data in _find_overflow(
                final_inventory, args.overflow_tolerance
            )
        ]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    payload = {
        "input": str(input_path),
        "output": str(output_path),
        "style_preset": preset.name,
        "iterations": loops_executed,
        "change_count": len(report_changes),
        "changes": report_changes,
        "geometry_fix_count": len(geometry_changes),
        "density_fix_count": len(density_changes),
        "remaining_overflow": remaining,
    }
    if args.report:
        report_path = Path(args.report).expanduser().resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        print(f"Text-fit report: {report_path}")
    print(
        f"Applied {len(report_changes)} text-fit change(s). "
        f"Remaining overflow shapes: {len(remaining)}"
    )

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
