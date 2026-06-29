#!/usr/bin/env python3
"""Detect visually weak but mechanically valid slide layouts.

This script focuses on sparse, underfilled text regions that are easy to miss
when overflow checks are the only QA signal.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation


def _inches(emu: int | float) -> float:
    return round(float(emu) / 914400.0, 2) if emu else 0.0


def _font_size_pt(shape: Any) -> float:
    if not getattr(shape, "has_text_frame", False):
        return 12.0
    for paragraph in shape.text_frame.paragraphs:
        paragraph_font = getattr(paragraph, "font", None)
        if paragraph_font is not None and paragraph_font.size is not None:
            return float(paragraph_font.size.pt)
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return 12.0


def _estimate_line_count(text: str, font_pt: float, box_w_in: float) -> int:
    if not text or box_w_in <= 0:
        return 0
    char_w = max(0.055, (font_pt / 72.0) * 0.55)
    chars_per_line = max(1, int(max(0.2, box_w_in - 0.12) / char_w))
    lines = 0
    for paragraph in text.split("\n"):
        lines += max(1, -(-len(paragraph.strip()) // chars_per_line))
    return lines


def _estimate_text_height(text: str, font_pt: float, box_w_in: float) -> float:
    lines = _estimate_line_count(text, font_pt, box_w_in)
    return lines * (font_pt / 72.0) * 1.24


def _is_section_title(shape: Any, font_pt: float) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    text = (shape.text or "").strip()
    return font_pt >= 28 and len(text) < 60 and text.count("\n") <= 1


def analyze_slide(slide: Any, slide_num: int, slide_w_emu: int, slide_h_emu: int) -> list[dict[str, Any]]:
    slide_w = _inches(slide_w_emu)
    slide_h = _inches(slide_h_emu)
    slide_area = slide_w * slide_h if slide_w and slide_h else 0.0
    shape_count = 0
    total_area = 0.0
    underfills: list[dict[str, Any]] = []

    for shape in slide.shapes:
        x = _inches(shape.left)
        y = _inches(shape.top)
        w = _inches(shape.width)
        h = _inches(shape.height)
        if w < 0.1 or h < 0.1:
            continue

        shape_count += 1
        total_area += w * h

        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text or "").strip()
        if not text or h <= 0.8 or w <= 0.5:
            continue
        font_pt = _font_size_pt(shape)
        if _is_section_title(shape, font_pt):
            continue

        estimated_h = _estimate_text_height(text, font_pt, w)
        fill_pct = (estimated_h / h) * 100 if h > 0 else 100
        if fill_pct < 25:
            underfills.append(
                {
                    "shape": getattr(shape, "name", "TextBox"),
                    "x": x,
                    "y": y,
                    "w": w,
                    "h": h,
                    "est_h": round(estimated_h, 2),
                    "fill_pct": round(fill_pct),
                }
            )

    issues: list[dict[str, Any]] = []
    grouped: defaultdict[tuple[float, float], list[dict[str, Any]]] = defaultdict(list)
    for entry in underfills:
        grouped[(round(entry["y"], 1), round(entry["h"], 1))].append(entry)

    for (_, _), members in grouped.items():
        if len(members) >= 2:
            avg_fill = round(sum(item["fill_pct"] for item in members) / len(members))
            max_est_h = max(item["est_h"] for item in members)
            suggested_h = round(max(max_est_h * 1.55, 0.6), 2)
            issues.append(
                {
                    "slide": slide_num,
                    "severity": "warning",
                    "type": "underfilled_card_row",
                    "count": len(members),
                    "y": members[0]["y"],
                    "box_h": members[0]["h"],
                    "avg_fill_pct": avg_fill,
                    "suggestion": (
                        f"{len(members)} cards in a row are visually underfilled "
                        f"(~{avg_fill}% text fill). Reduce height toward {suggested_h}\" "
                        f"or add supporting content."
                    ),
                }
            )
        else:
            member = members[0]
            suggested_h = round(max(member["est_h"] * 1.55, 0.6), 2)
            issues.append(
                {
                    "slide": slide_num,
                    "severity": "warning",
                    "type": "underfilled_textbox",
                    "shape": member["shape"],
                    "position": f"({member['x']}, {member['y']})",
                    "box_h": member["h"],
                    "est_text_h": member["est_h"],
                    "fill_pct": member["fill_pct"],
                    "suggestion": (
                        f"Text fills ~{member['fill_pct']}% of the box. "
                        f"Reduce height toward {suggested_h}\" or densify the content."
                    ),
                }
            )

    coverage_pct = (total_area / slide_area * 100) if slide_area > 0 else 100
    if coverage_pct < 25 and shape_count >= 3:
        issues.append(
            {
                "slide": slide_num,
                "severity": "info",
                "type": "sparse_slide",
                "shape_count": shape_count,
                "coverage_pct": round(coverage_pct),
                "suggestion": (
                    f"Only ~{round(coverage_pct)}% of the slide area is used. "
                    "Consider a denser pattern, larger evidence block, or a secondary visual anchor."
                ),
            }
        )

    return issues


def main() -> int:
    parser = argparse.ArgumentParser(description="Visual QA for PPTX slides")
    parser.add_argument("--input", required=True, help="Path to the .pptx file")
    parser.add_argument("--json", action="store_true", help="Print machine-readable JSON")
    args = parser.parse_args()

    pptx_path = Path(args.input).expanduser().resolve()
    prs = Presentation(str(pptx_path))
    issues: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides, start=1):
        issues.extend(analyze_slide(slide, index, prs.slide_width, prs.slide_height))

    if args.json:
        print(json.dumps(issues, indent=2))
    else:
        if not issues:
            print(f"Visual QA: {pptx_path} — all slides look good")
        else:
            warning_count = sum(1 for item in issues if item["severity"] == "warning")
            info_count = sum(1 for item in issues if item["severity"] == "info")
            print(f"Visual QA: {pptx_path}")
            print(f"  {warning_count} warning(s), {info_count} info message(s)\n")
            for issue in issues:
                if issue["type"] == "underfilled_card_row":
                    print(f"  [WARNING] Slide {issue['slide']}: {issue['suggestion']}\n")
                elif issue["type"] == "underfilled_textbox":
                    print(
                        f"  [WARNING] Slide {issue['slide']}: \"{issue['shape']}\" at {issue['position']}\n"
                        f"    Box: {issue['box_h']}\" tall, text fills ~{issue['fill_pct']}%\n"
                        f"    -> {issue['suggestion']}\n"
                    )
                else:
                    print(
                        f"  [INFO] Slide {issue['slide']}: {issue['shape_count']} shapes, "
                        f"~{issue['coverage_pct']}% coverage\n"
                        f"    -> {issue['suggestion']}\n"
                    )

    return 1 if any(item["severity"] == "warning" for item in issues) else 0


if __name__ == "__main__":
    raise SystemExit(main())
