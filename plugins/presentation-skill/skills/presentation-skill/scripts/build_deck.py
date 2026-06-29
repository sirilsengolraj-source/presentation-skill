#!/usr/bin/env python3
"""Build a styled PPTX deck from a structured JSON outline."""

from __future__ import annotations

import argparse
import csv
import json
import math
import re
import subprocess
import sys
import zlib
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from design_tokens import DEFAULT_FONT_PAIR_KEY, get_font_pair, get_style_preset

EMU_PER_INCH = 914400.0
EMOJI_PATTERN = re.compile(r"[\U0001F300-\U0001FAFF\u2600-\u27BF]")
RENDER_MODES = {"reliable", "express", "auto"}
EXPRESS_ALLOWED_SLIDE_TYPES = {"title", "section"}
EXPRESS_ALLOWED_VISUALS = {"hero", "timeline", "comparison"}
EXPRESS_ALLOWED_INTENTS = {"section", "process"}
DENSITY_VALUES = {"low", "medium", "high"}
EMOJI_MODES = {"none", "selective"}

PALETTE_LIBRARY: dict[str, dict[str, str]] = {
    "climate_coastal_v1": {
        "bg_primary": "ECFEFF",
        "bg_dark": "082F49",
        "surface": "FFFFFF",
        "text_primary": "0C4A6E",
        "text_muted": "155E75",
        "accent_primary": "0EA5E9",
        "accent_secondary": "14B8A6",
        "line": "BAE6FD",
    },
    "energy_sunset_v1": {
        "bg_primary": "FFF7ED",
        "bg_dark": "431407",
        "surface": "FFFFFF",
        "text_primary": "7C2D12",
        "text_muted": "9A3412",
        "accent_primary": "EA580C",
        "accent_secondary": "F59E0B",
        "line": "FED7AA",
    },
    "enterprise_graphite_v1": {
        "bg_primary": "F8FAFC",
        "bg_dark": "111827",
        "surface": "FFFFFF",
        "text_primary": "111827",
        "text_muted": "4B5563",
        "accent_primary": "2563EB",
        "accent_secondary": "0891B2",
        "line": "D1D5DB",
    },
}

_STAGED_ASSET_LOOKUP: dict[str, dict[str, str]] = {}


TITLE_MOTIF_CHOICES = ("orbit", "stripe", "wave", "arc", "geometric")


@dataclass(frozen=True)
class DeckStyleConfig:
    font_pair: str = DEFAULT_FONT_PAIR_KEY
    palette_key: str = ""
    visual_density: str = "medium"
    emoji_mode: str = "none"
    title_motif: str = ""
    # When True, builder renders "<slide_index> / <total>" on the right
    # side of each slide's footer. Default off to avoid changing every
    # existing deck's footer treatment. Recommended for editorial-minimal
    # decks and any deck the user expects to read like a printed report.
    show_page_numbers: bool = False


@dataclass(frozen=True)
class ComplianceConfig:
    attribution_file: Path
    require_attribution: bool


@dataclass(frozen=True)
class BuildConfig:
    outline_dir: Path
    deck_style: DeckStyleConfig
    compliance: ComplianceConfig


@dataclass(frozen=True)
class RuntimePreset:
    name: str
    palette: dict[str, str]
    typography: Any
    layout: Any
    font_pair: dict[str, str]


def _default_attribution_file(outline_dir: Path) -> Path:
    return (outline_dir / "assets" / "attribution.csv").resolve()


def _as_str(value: Any, default: str = "") -> str:
    text = str(value or "").strip()
    return text if text else default


def _parse_deck_style(data: dict[str, Any]) -> DeckStyleConfig:
    raw = data.get("deck_style")
    if not isinstance(raw, dict):
        return DeckStyleConfig()
    font_pair = _as_str(raw.get("font_pair"), DEFAULT_FONT_PAIR_KEY).lower()
    palette_key = _as_str(raw.get("palette_key"), "").lower()
    visual_density = _as_str(raw.get("visual_density"), "medium").lower()
    if visual_density not in DENSITY_VALUES:
        visual_density = "medium"
    emoji_mode = _as_str(raw.get("emoji_mode"), "none").lower()
    if emoji_mode not in EMOJI_MODES:
        emoji_mode = "none"
    title_motif = _as_str(raw.get("title_motif"), "").lower()
    # Accept the sentinel "none" as an explicit opt-out (skips motif draw
    # entirely — useful when hero_image is the visual anchor).
    if title_motif and title_motif != "none" and title_motif not in TITLE_MOTIF_CHOICES:
        title_motif = ""
    show_page_numbers = bool(raw.get("show_page_numbers"))
    return DeckStyleConfig(
        font_pair=font_pair,
        palette_key=palette_key,
        visual_density=visual_density,
        emoji_mode=emoji_mode,
        title_motif=title_motif,
        show_page_numbers=show_page_numbers,
    )


def _resolve_path(raw_path: str, outline_dir: Path) -> Path:
    candidate = Path(raw_path).expanduser()
    if not candidate.is_absolute():
        candidate = outline_dir / candidate
    return candidate.resolve()


def _default_staged_manifest(outline_dir: Path) -> Path:
    return (outline_dir / "assets" / "staged" / "staged_manifest.json").resolve()


def _staged_asset_lookup(outline_dir: Path) -> dict[str, str]:
    manifest_path = _default_staged_manifest(outline_dir)
    cache_key = str(manifest_path)
    if cache_key in _STAGED_ASSET_LOOKUP:
        return _STAGED_ASSET_LOOKUP[cache_key]

    lookup: dict[str, str] = {}
    if manifest_path.exists():
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
        sections = (
            ("images", ("asset", "image")),
            ("backgrounds", ("asset", "background")),
            ("charts", ("asset", "chart")),
            ("generated_images", ("asset", "image", "generated")),
        )
        for section, prefixes in sections:
            entries = payload.get(section)
            if not isinstance(entries, list):
                continue
            for entry in entries:
                if not isinstance(entry, dict):
                    continue
                name = _as_str(entry.get("name"), "").lower()
                asset_path = _as_str(entry.get("path"), "")
                if not name or not asset_path:
                    continue
                for prefix in prefixes:
                    lookup[f"{prefix}:{name}"] = asset_path
    _STAGED_ASSET_LOOKUP[cache_key] = lookup
    return lookup


def _parse_compliance(data: dict[str, Any], outline_dir: Path) -> ComplianceConfig:
    raw = data.get("compliance")
    if not isinstance(raw, dict):
        return ComplianceConfig(
            attribution_file=_default_attribution_file(outline_dir),
            require_attribution=False,
        )
    attribution_raw = _as_str(raw.get("attribution_file"), "assets/attribution.csv")
    attribution_file = _resolve_path(attribution_raw, outline_dir)
    require_attribution = bool(raw.get("require_attribution", False))
    return ComplianceConfig(
        attribution_file=attribution_file,
        require_attribution=require_attribution,
    )


def _resolve_palette(base_palette: dict[str, str], palette_key: str) -> dict[str, str]:
    if not palette_key:
        return dict(base_palette)
    key = palette_key.strip().lower()
    override = PALETTE_LIBRARY.get(key)
    if not override:
        return dict(base_palette)
    merged = dict(base_palette)
    merged.update(override)
    return merged


def _resolve_asset_path(value: Any, outline_dir: Path) -> Path | None:
    raw = str(value or "").strip()
    if not raw:
        return None
    normalized = raw.lower()
    if normalized.startswith(("asset:", "image:", "background:", "chart:", "generated:")):
        staged_path = _staged_asset_lookup(outline_dir).get(normalized)
        if not staged_path:
            raise FileNotFoundError(
                f"Staged asset alias not found: {raw}. "
                f"Expected {_default_staged_manifest(outline_dir)}"
            )
        return _resolve_path(staged_path, outline_dir)
    if raw.lower().startswith(("http://", "https://")):
        return None
    return _resolve_path(raw, outline_dir)


# Cache: original_svg_path -> converted_png_path (or None if conversion failed).
_SVG_CONVERSION_CACHE: dict[str, Path | None] = {}
_SVG_WARNING_EMITTED = False


def _convert_svg_to_png(svg_path: Path) -> Path | None:
    """Convert an SVG icon to PNG via cairosvg if available.

    Returns the PNG path on success, or None if cairosvg is unavailable or
    conversion fails. Caches results per-session so we only warn once.
    """
    global _SVG_WARNING_EMITTED
    key = str(svg_path.resolve())
    if key in _SVG_CONVERSION_CACHE:
        return _SVG_CONVERSION_CACHE[key]
    try:
        import cairosvg  # type: ignore
    except ImportError:
        if not _SVG_WARNING_EMITTED:
            print(
                "[build_deck] cairosvg not installed; SVG icons will be skipped. "
                "Install cairosvg or provide PNG icons instead.",
                file=sys.stderr,
            )
            _SVG_WARNING_EMITTED = True
        _SVG_CONVERSION_CACHE[key] = None
        return None
    try:
        png_path = svg_path.with_suffix(".converted.png")
        cairosvg.svg2png(
            url=str(svg_path),
            write_to=str(png_path),
            output_width=256,
            output_height=256,
        )
        _SVG_CONVERSION_CACHE[key] = png_path
        return png_path
    except Exception as exc:  # pragma: no cover - defensive
        print(f"[build_deck] SVG conversion failed for {svg_path}: {exc}", file=sys.stderr)
        _SVG_CONVERSION_CACHE[key] = None
        return None


def _resolve_icon_path(value: Any, outline_dir: Path) -> Path | None:
    """Resolve an `assets.icons[]` entry to a concrete file path.

    Resolution order:
      1. Absolute path: use as-is if file exists.
      2. Relative path with extension (.png/.svg/.jpg): resolve against
         `<outline_dir>/assets/icons/` first, then against outline_dir.
      3. Bare name (no extension, no slash): resolve against
         `<outline_dir>/assets/icons/<name>.png`, then `.svg`, then `.jpg`.

    Returns None (with a stderr warning) if no file exists. Returns the
    resolved PNG path on success; SVGs are converted via cairosvg if
    available, otherwise skipped with a warning.
    """
    raw = str(value or "").strip()
    if not raw:
        return None

    icons_dir = (outline_dir / "assets" / "icons").resolve()
    candidate: Path | None = None
    expected_msg = ""

    if raw.startswith("/"):
        # Case 1: absolute path.
        candidate = Path(raw).expanduser()
        expected_msg = str(candidate)
        if not candidate.exists():
            print(
                f"[build_deck] icon not found: {raw}, expected at {expected_msg}",
                file=sys.stderr,
            )
            return None
    else:
        has_ext = Path(raw).suffix.lower() in {".png", ".svg", ".jpg", ".jpeg"}
        if has_ext or "/" in raw or "\\" in raw:
            # Case 2: relative path — try icons_dir first, then outline_dir.
            primary = (icons_dir / raw).resolve()
            fallback = (outline_dir / raw).resolve()
            if primary.exists():
                candidate = primary
            elif fallback.exists():
                candidate = fallback
            else:
                expected_msg = str(primary)
        else:
            # Case 3: bare name — try .png, .svg, .jpg in icons_dir.
            for ext in (".png", ".svg", ".jpg", ".jpeg"):
                probe = (icons_dir / f"{raw}{ext}").resolve()
                if probe.exists():
                    candidate = probe
                    break
            if candidate is None:
                expected_msg = str((icons_dir / f"{raw}.png").resolve())

    if candidate is None:
        print(
            f"[build_deck] icon not found: {raw}, expected at {expected_msg}",
            file=sys.stderr,
        )
        return None

    if candidate.suffix.lower() == ".svg":
        converted = _convert_svg_to_png(candidate)
        if converted is None:
            print(
                f"[build_deck] icon skipped (SVG, no cairosvg): {candidate}",
                file=sys.stderr,
            )
            return None
        return converted

    return candidate


def _add_icon(
    slide: Any,
    *,
    path: Path,
    x: float,
    y: float,
    size: float,
) -> None:
    """Add a square icon image to the slide at (x, y) with given size (inches)."""
    try:
        slide.shapes.add_picture(
            str(path),
            Inches(x),
            Inches(y),
            width=Inches(size),
            height=Inches(size),
        )
    except Exception as exc:  # pragma: no cover - defensive
        print(f"[build_deck] failed to render icon {path}: {exc}", file=sys.stderr)


def _sources_footer_text(spec: dict[str, Any]) -> str:
    raw_sources = spec.get("sources")
    if not isinstance(raw_sources, list):
        return ""
    parts: list[str] = []
    for item in raw_sources:
        if isinstance(item, dict):
            text = _as_str(
                item.get("label")
                or item.get("title")
                or item.get("source")
                or item.get("source_page")
                or item.get("url"),
                "",
            )
        else:
            text = _as_str(item, "")
        if text:
            parts.append(text)
    if not parts:
        return ""
    if len(parts) > 3:
        shown = parts[:3]
        return f"Sources: {'; '.join(shown)}; +{len(parts) - 3} more"
    return f"Sources: {'; '.join(parts)}"


def _slide_footer_text(spec: dict[str, Any]) -> str:
    footer = _as_str(spec.get("footer"), "")
    sources = _sources_footer_text(spec)
    if footer and sources:
        text = f"{footer} | {sources}"
    else:
        text = footer or sources
    text = text.strip()
    if len(text) <= 180:
        return text
    return text[:177].rstrip() + "..."


def _metadata_candidates(asset_path: Path) -> list[Path]:
    return [
        asset_path.with_suffix(".metadata.json"),
        Path(f"{asset_path}.metadata.json"),
    ]


def _has_asset_metadata(asset_path: Path) -> bool:
    return any(candidate.exists() for candidate in _metadata_candidates(asset_path))


def _is_external_heuristic(asset_path: Path) -> bool:
    lowered = str(asset_path).lower()
    return (
        "tmp_assets" in lowered
        or "wikimedia" in lowered
        or "commons.wikimedia.org" in lowered
    )


def _looks_like_emoji_policy_disabled(style_preset_name: str) -> bool:
    lowered = style_preset_name.lower()
    return "boardroom" in lowered or "data-heavy" in lowered


def _strip_emojis(text: str) -> str:
    return EMOJI_PATTERN.sub("", text or "").strip()


def _limit_emojis(text: str, max_count: int) -> str:
    if max_count <= 0:
        return _strip_emojis(text)
    remaining = max_count
    output: list[str] = []
    for char in text:
        if EMOJI_PATTERN.match(char):
            if remaining <= 0:
                continue
            remaining -= 1
        output.append(char)
    return "".join(output).strip()


def _apply_emoji_policy(
    spec: dict[str, Any],
    *,
    slide_type: str,
    preset_name: str,
    emoji_mode: str,
) -> dict[str, Any]:
    if emoji_mode == "none" or _looks_like_emoji_policy_disabled(preset_name):
        sanitized = deepcopy(spec)
        for key in ("title", "subtitle", "body", "footer", "caption", "message"):
            if key in sanitized:
                sanitized[key] = _strip_emojis(str(sanitized.get(key, "")))
        if isinstance(sanitized.get("bullets"), list):
            sanitized["bullets"] = [
                _strip_emojis(str(item if not isinstance(item, dict) else item.get("text", "")))
                if not isinstance(item, dict)
                else {**item, "text": _strip_emojis(str(item.get("text", "")))}
                for item in sanitized["bullets"]
            ]
        return sanitized

    informal = _as_str(spec.get("tone"), "").lower() == "informal"
    allow = slide_type in {"title", "section"} or informal
    sanitized = deepcopy(spec)
    if not allow:
        return _apply_emoji_policy(
            sanitized,
            slide_type=slide_type,
            preset_name=preset_name,
            emoji_mode="none",
        )

    title = _as_str(sanitized.get("title"), "")
    sanitized["title"] = _limit_emojis(title, 1)
    remaining = 2 - len(EMOJI_PATTERN.findall(sanitized["title"]))
    remaining = max(0, remaining)
    for key in ("subtitle", "body", "footer", "caption", "message"):
        if key in sanitized:
            sanitized[key] = _limit_emojis(_as_str(sanitized.get(key), ""), remaining)
            remaining = max(0, remaining - len(EMOJI_PATTERN.findall(sanitized[key])))
    if isinstance(sanitized.get("bullets"), list):
        normalized_bullets: list[Any] = []
        for bullet in sanitized["bullets"]:
            if isinstance(bullet, dict):
                text = _limit_emojis(_as_str(bullet.get("text"), ""), remaining)
                remaining = max(0, remaining - len(EMOJI_PATTERN.findall(text)))
                normalized_bullets.append({**bullet, "text": text})
            else:
                text = _limit_emojis(_as_str(bullet, ""), remaining)
                remaining = max(0, remaining - len(EMOJI_PATTERN.findall(text)))
                normalized_bullets.append(text)
        sanitized["bullets"] = normalized_bullets
    return sanitized


def _coerce_render_mode(value: Any) -> str:
    # Backward compatibility and reliability-first default:
    # if render_mode is omitted, stay on deterministic reliable path.
    mode = _as_str(value, "reliable").lower()
    if mode not in RENDER_MODES:
        return "reliable"
    return mode


def _slide_density_level(spec: dict[str, Any]) -> str:
    bullets = spec.get("bullets")
    bullet_count = len(bullets) if isinstance(bullets, list) else 0
    paragraphs = spec.get("paragraphs")
    para_count = len(paragraphs) if isinstance(paragraphs, list) else 0
    cards = spec.get("cards")
    card_count = len(cards) if isinstance(cards, list) else 0
    total = bullet_count + para_count + card_count
    if total >= 9:
        return "high"
    if total <= 3:
        return "low"
    return "medium"


def _express_supported(spec: dict[str, Any], slide_type: str) -> bool:
    visual_intent = _as_str(spec.get("visual_intent"), "").lower()
    slide_intent = _as_str(spec.get("slide_intent"), "").lower()
    variant = _as_str(spec.get("variant"), "").lower()
    assets = _slide_assets(spec)
    if visual_intent == "flow" and not _as_str(assets.get("diagram"), ""):
        return False
    if slide_type in EXPRESS_ALLOWED_SLIDE_TYPES:
        return True
    if visual_intent in EXPRESS_ALLOWED_VISUALS:
        return True
    if slide_intent in EXPRESS_ALLOWED_INTENTS:
        return True
    if variant == "timeline":
        return True
    return variant == "split" and visual_intent == "comparison"


def _resolve_render_mode(
    spec: dict[str, Any],
    *,
    slide_type: str,
    deck_density: str,
) -> str:
    requested = _coerce_render_mode(spec.get("render_mode"))
    if requested == "reliable":
        return "reliable"
    supported = _express_supported(spec, slide_type)
    if requested == "express":
        return "express" if supported else "reliable"
    # auto mode
    intent = _as_str(spec.get("slide_intent"), "").lower()
    visual_intent = _as_str(spec.get("visual_intent"), "").lower()
    density = _slide_density_level(spec)
    if intent in {"section", "process"} or visual_intent in {"hero", "timeline", "comparison"}:
        if density != "high" and deck_density != "high" and supported:
            return "express"
    return "reliable"


def _collect_asset_paths(spec: dict[str, Any], outline_dir: Path) -> list[Path]:
    candidates: list[Path] = []
    bg = _resolve_asset_path(spec.get("background_image"), outline_dir)
    if bg:
        candidates.append(bg)
    thumbs = spec.get("thumbnails")
    if isinstance(thumbs, list):
        for item in thumbs:
            thumb = _resolve_asset_path(item, outline_dir)
            if thumb:
                candidates.append(thumb)
    assets = spec.get("assets")
    if isinstance(assets, dict):
        for key in ("hero_image", "diagram", "logo", "chart_data", "chart"):
            path = _resolve_asset_path(assets.get(key), outline_dir)
            if path:
                candidates.append(path)
        logos = assets.get("logos")
        if isinstance(logos, list):
            for item in logos:
                path = _resolve_asset_path(item, outline_dir)
                if path:
                    candidates.append(path)
    return candidates


def _read_attribution_rows(path: Path) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    if not path.exists():
        return rows
    with path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            rows.append({k: str(v or "").strip() for k, v in row.items()})
    return rows


def _attribution_contains(asset_path: Path, rows: list[dict[str, str]]) -> bool:
    basename = asset_path.name
    for row in rows:
        if basename in row.values() or str(asset_path) in row.values():
            return True
    return False


def _enforce_compliance(data: dict[str, Any], config: BuildConfig) -> None:
    slides = data.get("slides")
    if not isinstance(slides, list):
        return

    all_assets: list[Path] = []
    remote_refs: list[str] = []
    external_assets: list[Path] = []
    flagged_without_metadata: list[Path] = []

    for item in slides:
        if not isinstance(item, dict):
            continue
        assets = _collect_asset_paths(item, config.outline_dir)
        slide_assets = item.get("assets") if isinstance(item.get("assets"), dict) else {}
        required_assets: set[Path] = set()
        for key in ("hero_image", "diagram", "logo", "chart_data", "chart"):
            candidate = _resolve_asset_path(slide_assets.get(key), config.outline_dir)
            if candidate:
                required_assets.add(candidate)
        logos_block = slide_assets.get("logos")
        if isinstance(logos_block, list):
            for entry in logos_block:
                candidate = _resolve_asset_path(entry, config.outline_dir)
                if candidate:
                    required_assets.add(candidate)
        diagram_target = _resolve_asset_path(slide_assets.get("diagram"), config.outline_dir)
        mermaid_source = _resolve_asset_path(
            slide_assets.get("mermaid") or slide_assets.get("mermaid_source"),
            config.outline_dir,
        )
        all_assets.extend(assets)
        for asset in assets:
            if not asset.exists():
                if asset not in required_assets:
                    continue
                if (
                    diagram_target
                    and asset == diagram_target
                    and mermaid_source
                    and mermaid_source.exists()
                ):
                    # Flow diagrams can be generated from Mermaid at render time.
                    continue
                raise FileNotFoundError(f"Asset not found: {asset}")
            if _has_asset_metadata(asset):
                external_assets.append(asset)
            elif _is_external_heuristic(asset):
                flagged_without_metadata.append(asset)

        for raw in (item.get("background_image"),):
            raw_text = str(raw or "").strip().lower()
            if raw_text.startswith(("http://", "https://")):
                remote_refs.append(str(raw))
        if isinstance(slide_assets, dict):
            for key in ("hero_image", "diagram", "logo", "chart_data", "chart"):
                raw_text = str(slide_assets.get(key, "")).strip().lower()
                if raw_text.startswith(("http://", "https://")):
                    remote_refs.append(str(slide_assets.get(key)))
            logos = slide_assets.get("logos")
            if isinstance(logos, list):
                for entry in logos:
                    raw_text = str(entry or "").strip().lower()
                    if raw_text.startswith(("http://", "https://")):
                        remote_refs.append(str(entry))

    if remote_refs:
        raise ValueError(
            "Remote media URLs are not allowed. Use local files for logos and licensed assets."
        )

    if flagged_without_metadata:
        names = ", ".join(path.name for path in flagged_without_metadata[:5])
        raise ValueError(
            "External-looking assets are missing metadata sidecars. "
            f"Add <asset>.metadata.json or replace them. Examples: {names}"
        )

    needs_attribution = config.compliance.require_attribution or bool(external_assets)
    if not needs_attribution:
        return

    attribution_file = config.compliance.attribution_file
    rows = _read_attribution_rows(attribution_file)
    if not rows:
        raise ValueError(
            f"Attribution is required but file is missing/empty: {attribution_file}"
        )
    missing = [asset for asset in external_assets if not _attribution_contains(asset, rows)]
    if missing:
        names = ", ".join(path.name for path in missing[:8])
        raise ValueError(
            "External CC assets missing attribution entries: "
            f"{names}. Update {attribution_file}."
        )


def _load_build_config(data: dict[str, Any], outline_dir: Path) -> BuildConfig:
    deck_style = _parse_deck_style(data)
    compliance = _parse_compliance(data, outline_dir)
    return BuildConfig(
        outline_dir=outline_dir,
        deck_style=deck_style,
        compliance=compliance,
    )


def _layout_by_hint(prs: Presentation, hints: Iterable[str], fallback: int) -> Any:
    hint_list = [item.lower() for item in hints]
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if any(hint in name for hint in hint_list):
            return layout
    if len(prs.slide_layouts) == 0:
        raise ValueError("Presentation has no slide layouts.")
    index = max(0, min(fallback, len(prs.slide_layouts) - 1))
    return prs.slide_layouts[index]


def _blank_layout(prs: Presentation) -> Any:
    return _layout_by_hint(prs, ["blank"], fallback=max(0, len(prs.slide_layouts) - 1))


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _safe_float(value: Any, default: float | None = None) -> float | None:
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value or "").strip().replace(",", "")
    if not text:
        return default
    if text.endswith("%"):
        text = text[:-1]
    try:
        return float(text)
    except (TypeError, ValueError):
        return default


def _rgb(color_hex: str) -> RGBColor:
    value = (color_hex or "000000").strip().replace("#", "")
    if len(value) != 6:
        value = "000000"
    return RGBColor.from_string(value.upper())


def _slide_w(prs: Presentation) -> float:
    return prs.slide_width / EMU_PER_INCH


def _slide_h(prs: Presentation) -> float:
    return prs.slide_height / EMU_PER_INCH


def _font_title(preset: Any) -> str:
    return str(getattr(preset, "font_pair", {}).get("title", "Trebuchet MS"))


def _font_body(preset: Any) -> str:
    return str(getattr(preset, "font_pair", {}).get("body", "Calibri"))


def _font_caption(preset: Any) -> str:
    return str(getattr(preset, "font_pair", {}).get("caption", "Calibri"))


def _clamp_title_font(typography: Any, text: str, *, density: str) -> int:
    base = typography.title_max
    if density == "high":
        base = min(base, typography.title_min + 4)
    elif density == "medium":
        base = min(base, typography.title_min + 8)
    length = len(text.strip())
    if length > 90:
        base = max(typography.title_min, base - 6)
    elif length > 60:
        base = max(typography.title_min, base - 4)
    return max(typography.title_min, min(typography.title_max, base))


def _clamp_body_font(typography: Any, requested: int) -> int:
    min_readable = max(typography.body_min, 13)
    return max(min_readable, min(typography.body_max, requested))


def _set_background(slide: Any, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def _set_text_frame_defaults(frame: Any, margin: float = 0.05) -> None:
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(max(0.01, margin * 0.7))
    frame.margin_bottom = Inches(max(0.01, margin * 0.7))


def _set_paragraph_style(
    paragraph: Any,
    *,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    align: PP_ALIGN | None = None,
) -> None:
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = _rgb(color_hex)
    paragraph.font.bold = bold
    # Keep run-level styling in sync with paragraph-level styling so text
    # measurement and downstream QA detect the real rendered font settings.
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = _rgb(color_hex)
        run.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def _set_text_box(
    slide: Any,
    *,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    margin: float = 0.05,
    wrap: bool = True,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    _set_text_frame_defaults(frame, margin=margin)
    frame.word_wrap = wrap
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    _set_paragraph_style(
        paragraph,
        font_name=font_name,
        font_size=font_size,
        color_hex=color_hex,
        bold=bold,
        align=align,
    )
    return box


def _set_lines_box(
    slide: Any,
    *,
    lines: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    font_name: str,
    font_size: int,
    color_hex: str,
    bullet: bool = False,
    margin: float = 0.05,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    _set_text_frame_defaults(frame, margin=margin)
    cleaned_lines = [line.strip() for line in lines if line and line.strip()]
    if not cleaned_lines:
        return box

    def _apply_bullet(paragraph: Any, level: int = 0) -> None:
        paragraph.level = max(0, level)
        ppr = paragraph._element.get_or_add_pPr()
        for child in list(ppr):
            if child.tag.endswith("buChar") or child.tag.endswith("buNone") or child.tag.endswith("buAutoNum"):
                ppr.remove(child)
        # Use bullet formatting in paragraph properties (not unicode text symbols).
        bu_char = OxmlElement("a:buChar")
        bu_char.set("char", "•")
        ppr.append(bu_char)

    for index, line in enumerate(cleaned_lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        if bullet:
            _apply_bullet(paragraph, level=0)
        else:
            paragraph.level = 0
        paragraph.space_after = Pt(8 if bullet else 6)
        _set_paragraph_style(
            paragraph,
            font_name=font_name,
            font_size=font_size,
            color_hex=color_hex,
        )
    return box


def _add_line(
    slide: Any,
    *,
    x: float,
    y: float,
    w: float,
    color_hex: str,
    width_pt: float = 2.0,
) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x),
        Inches(y),
        Inches(x + w),
        Inches(y),
    )
    line.line.color.rgb = _rgb(color_hex)
    line.line.width = Pt(width_pt)


def _estimate_text_lines(text: str, *, width_in: float, font_size_pt: float) -> int:
    if not text.strip():
        return 0
    usable_width = max(0.4, width_in - 0.16)
    char_width_in = max(0.055, (font_size_pt / 72.0) * 0.55)
    chars_per_line = max(6, int(usable_width / char_width_in))
    total = 0
    for paragraph in text.splitlines() or [text]:
        stripped = paragraph.strip()
        total += max(1, int(math.ceil(len(stripped) / float(chars_per_line)))) if stripped else 1
    return total


def _estimate_text_height(lines: list[str], *, width_in: float, font_size_pt: float) -> float:
    text = "\n".join(line.strip() for line in lines if line and line.strip())
    if not text:
        return 0.0
    line_count = _estimate_text_lines(text, width_in=width_in, font_size_pt=font_size_pt)
    return line_count * (font_size_pt / 72.0) * 1.24


def _card_body_font(
    body_lines: list[str],
    *,
    width_in: float,
    typography: Any,
    available_height_in: float | None = None,
) -> int:
    text = " ".join(line.strip() for line in body_lines if line and line.strip())
    longest_line = max((len(line.strip()) for line in body_lines if line and line.strip()), default=0)
    requested = typography.body_max - max(0, (len(text) - 120) // 40)
    if width_in < 5.2:
        requested -= 1
    if width_in < 4.6:
        requested -= 1
    if width_in < 4.0:
        requested -= 1
    if longest_line > 34:
        requested -= 1
    if longest_line > 46:
        requested -= 1
    font = _clamp_body_font(typography, requested)
    if available_height_in is None:
        return font

    usable_width = max(0.6, width_in - 0.28)
    while font > typography.body_min:
        body_h = _estimate_text_height(body_lines[:6], width_in=usable_width, font_size_pt=font)
        if body_h <= available_height_in:
            break
        font -= 1
    # Fix 2a safety: when text still overflows at the preset body_min (e.g.
    # tight multi-card slides after the body_min bump from 13 -> 15), allow a
    # small emergency dip below body_min to prevent card overflow. Floor at
    # 12pt so nothing becomes illegible.
    emergency_floor = max(12, typography.body_min - 3)
    while font > emergency_floor:
        body_h = _estimate_text_height(body_lines[:6], width_in=usable_width, font_size_pt=font)
        if body_h <= available_height_in:
            break
        font -= 1
    return max(emergency_floor, font)


def _card_title_layout(title_text: str, width_in: float, typography: Any, preset: Any) -> tuple[int, float]:
    title_font_name = _font_title(preset)
    is_serif_title = any(
        token in title_font_name.lower()
        for token in ("georgia", "times", "palatino", "garamond", "cambria")
    )
    base_cap = 21 if is_serif_title else 24
    heading_font = max(typography.section_min - 2, min(typography.section_max - 2, base_cap))
    if width_in < 2.6:
        heading_font = min(heading_font, 17 if is_serif_title else 18)
    if width_in < 2.3:
        heading_font = min(heading_font, 16 if is_serif_title else 17)
    chars_per_line = max(8, int((width_in - 0.32) * 7.0))
    estimated_lines = max(1, min(2, int(math.ceil(len(title_text) / float(chars_per_line)))))
    if estimated_lines > 1:
        heading_font = max(
            typography.section_min - 5,
            heading_font - (3 if is_serif_title else 2),
        )
    title_h = min(1.02, max(0.60, 0.38 + estimated_lines * 0.24))
    return heading_font, title_h


def _preferred_card_height(
    *,
    title_text: str,
    body_lines: list[str],
    width_in: float,
    typography: Any,
    preset: Any,
    rail_h: float,
    min_h: float,
    max_h: float,
    icon_height: float = 0.0,
) -> float:
    heading_font, title_h = _card_title_layout(title_text, width_in, typography, preset)
    body_font = _card_body_font(
        body_lines,
        width_in=width_in,
        typography=typography,
    )
    body_h = _estimate_text_height(body_lines[:6], width_in=width_in - 0.28, font_size_pt=body_font)
    preferred = rail_h + icon_height + title_h + 0.16 + body_h + 0.20
    return min(max(max(min_h, preferred), 0.9), max_h)


def _add_card(
    slide: Any,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body_lines: list[str],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    rail_h: float,
    accent_key: str = "accent_primary",
    dark: bool = False,
    body_vertical_anchor: Any = None,
    icon_path: Path | None = None,
    icon_size: float = 0.0,
) -> None:
    fill_color = palette["bg_dark"] if dark else palette["surface"]
    line_color = palette["line"] if not dark else palette["accent_primary"]
    text_color = "FFFFFF" if dark else palette["text_primary"]
    muted_color = "CBD5E1" if dark else palette["text_muted"]

    # Keep rail + card geometry coherent: accent rails are rectangular, so
    # the card body must be rectangular too (no rounded-corner mismatch).
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(fill_color)
    card.line.color.rgb = _rgb(line_color)
    card.line.width = Pt(1.25)

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(rail_h),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = _rgb(palette.get(accent_key, palette["accent_primary"]))
    rail.line.fill.background()

    # When the caller passes an empty title explicitly, skip the title
    # row entirely — the slide's own header already frames the card and
    # a stale default like "Key Points" reads as auto-generated filler.
    title_text = title.strip()
    suppress_title = title_text == ""
    title_font_name = _font_title(preset)
    if suppress_title:
        heading_font = typography.section_min
        title_h = 0.0
    else:
        heading_font, title_h = _card_title_layout(title_text, w, typography, preset)
    body_gap = 0.16 if not suppress_title else 0.10

    # Icon block: positioned inside the card, below the rail, above the title.
    # Adds icon_size + icon_pad to the title-top so title/body shift down.
    icon_pad = 0.15 if (icon_path is not None and icon_size > 0) else 0.0
    icon_block_h = (icon_size + icon_pad) if (icon_path is not None and icon_size > 0) else 0.0
    if icon_path is not None and icon_size > 0:
        icon_top = y + rail_h + 0.08
        icon_left = x + (w - icon_size) / 2.0
        _add_icon(slide, path=icon_path, x=icon_left, y=icon_top, size=icon_size)

    title_top = y + rail_h + icon_block_h + 0.06
    body_y = y + rail_h + icon_block_h + title_h + body_gap
    body_h = max(0.5, h - rail_h - icon_block_h - title_h - body_gap - 0.08)
    body_font = _card_body_font(
        body_lines,
        width_in=w,
        typography=typography,
        available_height_in=body_h,
    )

    if not suppress_title:
        _set_text_box(
            slide,
            text=title_text,
            x=x + 0.14,
            y=title_top,
            w=w - 0.28,
            h=max(0.36, title_h - 0.02),
            font_name=title_font_name,
            font_size=heading_font,
            color_hex=text_color,
            bold=True,
            margin=0.01,
        )
    body_box = _set_lines_box(
        slide,
        lines=body_lines[:6],
        x=x + 0.14,
        y=body_y,
        w=w - 0.28,
        h=body_h,
        font_name=_font_body(preset),
        font_size=body_font,
        color_hex=muted_color if len(body_lines) > 1 else text_color,
        bullet=False,
        margin=0.01,
    )
    # Fix 1/2b: when caller requests a vertical anchor (e.g. MIDDLE), apply it
    # to the body text frame so expanded cards don't leave dead space below
    # short content.
    if body_vertical_anchor is not None and body_box is not None:
        try:
            body_box.text_frame.vertical_anchor = body_vertical_anchor
        except Exception:
            pass


def _normalize_bullets(items: Any) -> list[tuple[str, int]]:
    normalized: list[tuple[str, int]] = []
    if not isinstance(items, list):
        return normalized
    for item in items:
        if isinstance(item, str):
            text = item.strip()
            if text:
                normalized.append((text, 0))
            continue
        if isinstance(item, dict):
            text = str(item.get("text", "")).strip()
            if not text:
                continue
            level = max(0, min(_safe_int(item.get("level", 0), 0), 8))
            normalized.append((text, level))
    return normalized


def _extract_lines(spec: dict[str, Any]) -> list[str]:
    bullets = _normalize_bullets(spec.get("bullets"))
    if bullets:
        return [text for text, _ in bullets]
    paragraphs = spec.get("paragraphs")
    if isinstance(paragraphs, list):
        return [str(item).strip() for item in paragraphs if str(item).strip()]
    body = str(spec.get("body", "")).strip()
    if not body:
        return []
    parts = [part.strip() for part in body.replace("\n", ". ").split(".")]
    return [part for part in parts if part]


def _content_header(
    slide: Any,
    *,
    title: str,
    subtitle: str,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    margin_x: float,
    right_reserve: float = 0.0,
) -> float:
    header_w = max(6.0, 8.6 - right_reserve)
    title_text = title.strip() or "Content"
    title_cap = 28 if right_reserve >= 3.0 else 30
    title_font = max(typography.section_min, min(typography.section_max, title_cap))
    title_lines = max(1, min(3, _estimate_text_lines(title_text, width_in=header_w, font_size_pt=title_font)))
    title_h = min(1.18, max(0.62, 0.30 + title_lines * 0.32))
    if right_reserve >= 3.0:
        title_h = min(1.28, title_h + 0.10)
    title_y = 0.34

    # Editorial-minimal treatment: pure white slides carry one strong
    # color, used as a left accent bar flanking the title. The title
    # text itself stays in near-black — the colored bar is the branding
    # element, not the typography. This is the "heading title box with
    # color" look for editorial/research decks.
    preset_name = getattr(preset, "name", "")
    editorial = preset_name == "editorial-minimal"
    lab_report = preset_name == "lab-report"
    title_x = margin_x
    title_color_override: str | None = None
    subtitle_color_override: str | None = None

    if editorial:
        bar_w = 0.10
        bar_gap = 0.20
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(margin_x),
            Inches(title_y + 0.06),
            Inches(bar_w),
            Inches(max(0.50, title_h - 0.12)),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(palette.get("accent_primary", "D4461E"))
        bar.line.fill.background()
        title_x = margin_x + bar_w + bar_gap
        header_w = max(4.0, header_w - (bar_w + bar_gap))

    elif lab_report:
        # Clinical-data aesthetic inspired by (not cloned from) lab
        # comparison decks: dark header block + thin accent rule
        # beneath it signals technical precision. Pulls the actual
        # slide width from the presentation rather than assuming 13.33"
        # so this scales to any canvas.
        try:
            slide_w_actual = slide.part.package.presentation_part.presentation.slide_width / 914400.0
        except Exception:
            slide_w_actual = 13.33
        header_h = 0.94
        accent_h = 0.04
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(slide_w_actual),
            Inches(header_h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(palette.get("bg_dark", "0B2545"))
        bar.line.fill.background()
        # Thin accent rule beneath the header — skill-signature touch.
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(header_h),
            Inches(slide_w_actual),
            Inches(accent_h),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(palette.get("accent_secondary", "C9302C"))
        accent.line.fill.background()
        title_y = 0.14
        title_h = min(0.54, title_h)
        title_color_override = "FFFFFF"
        subtitle_color_override = "D1D5DB"

    _set_text_box(
        slide,
        text=title_text,
        x=title_x,
        y=title_y,
        w=header_w,
        h=title_h,
        font_name=_font_title(preset),
        font_size=title_font,
        color_hex=title_color_override or palette["text_primary"],
        bold=True,
        margin=0.01,
    )
    header_bottom = title_y + title_h
    if subtitle.strip():
        subtitle_w = max(6.2, 9.0 - right_reserve)
        subtitle_cap = 15 if right_reserve >= 3.0 else 16
        subtitle_font = max(typography.caption_min + 1, min(typography.caption_max + 2, subtitle_cap))
        subtitle_lines = max(
            1,
            min(2, _estimate_text_lines(subtitle.strip(), width_in=subtitle_w, font_size_pt=subtitle_font)),
        )
        subtitle_h = min(0.62, max(0.26, 0.12 + subtitle_lines * 0.22))
        if right_reserve >= 3.0:
            subtitle_h = min(0.70, subtitle_h + 0.08)
        subtitle_y = header_bottom + 0.04
        _set_text_box(
            slide,
            text=subtitle.strip(),
            x=title_x,
            y=subtitle_y,
            w=subtitle_w,
            h=subtitle_h,
            font_name=_font_caption(preset),
            font_size=subtitle_font,
            color_hex=subtitle_color_override or palette["text_muted"],
            margin=0.01,
        )
        header_bottom = subtitle_y + subtitle_h
    return header_bottom


def _footer_geometry(slide_h: float) -> tuple[float, float]:
    footer_h = 0.28
    footer_bottom_pad = 0.08
    footer_y = slide_h - footer_h - footer_bottom_pad
    return footer_y, footer_h


def _content_bottom(slide_h: float, *, has_footer: bool, layout: Any) -> float:
    if has_footer:
        footer_y, _ = _footer_geometry(slide_h)
        return footer_y - 0.10
    return slide_h - max(0.22, float(getattr(layout, "bottom_safe", 0.30)))


def _add_background_image(
    slide: Any,
    spec: dict[str, Any],
    *,
    slide_w: float,
    slide_h: float,
    palette: dict[str, str],
    outline_dir: Path,
) -> bool:
    candidate = _resolve_asset_path(spec.get("background_image"), outline_dir)
    if not candidate:
        return False
    if not candidate.exists() or not candidate.is_file():
        return False
    slide.shapes.add_picture(
        str(candidate),
        Inches(0),
        Inches(0),
        width=Inches(slide_w),
        height=Inches(slide_h),
    )
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(slide_w),
        Inches(slide_h),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = _rgb(palette["bg_primary"])
    try:
        overlay.fill.transparency = 20
    except Exception:
        pass
    overlay.line.fill.background()
    return True


def _add_content_motif(
    slide: Any,
    *,
    palette: dict[str, str],
    slide_w: float,
    margin_x: float,
) -> None:
    anchors = [
        (0.42, 0.30, "accent_primary", 85),
        (0.28, 0.78, "accent_secondary", 88),
        (0.20, 1.12, "line", 90),
    ]
    for size, y, color_key, transparency in anchors:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(slide_w - margin_x - size),
            Inches(y),
            Inches(size),
            Inches(size),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(palette.get(color_key, palette["accent_primary"]))
        try:
            shape.fill.transparency = transparency
        except Exception:
            pass
        shape.line.fill.background()


def _add_section_motif_enrichment(
    slide: Any,
    *,
    palette: dict[str, str],
    slide_w: float,
    slide_h: float,
    margin_x: float,
) -> None:
    """Quiet decorative footprint for content-free section dividers.

    Draws a horizontal accent rail ~60% down the canvas plus three
    circular motif dots sitting along that rail. Low-contrast on
    purpose -- this is enrichment, not content.
    """
    rail_y = slide_h * 0.60
    rail_height = 0.10
    rail_w = slide_w - margin_x * 2
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(margin_x),
        Inches(rail_y),
        Inches(rail_w),
        Inches(rail_height),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = _rgb(palette.get("accent_primary", palette.get("line", "BAE6FD")))
    try:
        rail.fill.transparency = 55
    except Exception:
        pass
    rail.line.fill.background()

    dot_diameters = [0.48, 0.36, 0.30]
    # Offsets from left margin where each dot centers.
    centers = [
        margin_x + rail_w * 0.18,
        margin_x + rail_w * 0.48,
        margin_x + rail_w * 0.82,
    ]
    color_keys = ["accent_primary", "accent_secondary", "accent_primary"]
    transparencies = [70, 65, 78]
    rail_center_y = rail_y + rail_height / 2.0
    for center_x, diameter, color_key, transparency in zip(
        centers, dot_diameters, color_keys, transparencies
    ):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - diameter / 2.0),
            Inches(rail_center_y - diameter / 2.0),
            Inches(diameter),
            Inches(diameter),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(
            palette.get(color_key, palette.get("accent_primary", "0EA5E9"))
        )
        try:
            dot.fill.transparency = transparency
        except Exception:
            pass
        dot.line.fill.background()


def _add_thumbnail_strip(
    slide: Any,
    *,
    thumbnails: list[str],
    slide_w: float,
    margin_x: float,
    palette: dict[str, str],
    outline_dir: Path,
) -> None:
    valid: list[Path] = []
    for item in thumbnails:
        candidate = _resolve_asset_path(item, outline_dir)
        if not candidate:
            continue
        if candidate.exists() and candidate.is_file():
            valid.append(candidate)
        if len(valid) >= 3:
            break
    if not valid:
        return
    thumb_w = 0.90
    thumb_h = 0.50
    gap = 0.10
    total_w = len(valid) * thumb_w + (len(valid) - 1) * gap
    start_x = max(margin_x + 6.5, slide_w - margin_x - total_w)
    y = 0.30
    for index, path in enumerate(valid):
        x = start_x + index * (thumb_w + gap)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(thumb_w),
            Inches(thumb_h),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = _rgb(palette["surface"])
        frame.line.color.rgb = _rgb(palette["line"])
        frame.line.width = Pt(0.8)
        try:
            frame.adjustments[0] = 0.08
        except Exception:
            pass
        slide.shapes.add_picture(
            str(path),
            Inches(x + 0.02),
            Inches(y + 0.02),
            width=Inches(thumb_w - 0.04),
            height=Inches(thumb_h - 0.04),
        )


def _slide_assets(spec: dict[str, Any]) -> dict[str, Any]:
    assets = spec.get("assets")
    if isinstance(assets, dict):
        return assets
    return {}


def _slide_icon_paths(
    spec: dict[str, Any], outline_dir: Path, count: int
) -> list[Path | None]:
    """Return a list of resolved icon paths sized to `count` (None for missing).

    Never raises. Icons are v1 supported on cards-2/cards-3, timeline, stats.
    Missing icons become None so callers can render the element without the
    icon. A stderr warning is emitted per miss by `_resolve_icon_path`.
    """
    assets = _slide_assets(spec)
    raw_icons = assets.get("icons")
    if not isinstance(raw_icons, list) or count <= 0:
        return [None] * max(0, count)
    resolved: list[Path | None] = []
    for idx in range(count):
        if idx >= len(raw_icons):
            resolved.append(None)
            continue
        resolved.append(_resolve_icon_path(raw_icons[idx], outline_dir))
    return resolved


def _asset_from_spec(spec: dict[str, Any], key: str, outline_dir: Path) -> Path | None:
    assets = _slide_assets(spec)
    return _resolve_asset_path(assets.get(key), outline_dir)


def _render_mermaid_asset_if_needed(spec: dict[str, Any], outline_dir: Path) -> Path | None:
    assets = _slide_assets(spec)
    diagram = _resolve_asset_path(assets.get("diagram"), outline_dir)
    if diagram and diagram.exists():
        return diagram
    mermaid_source = _resolve_asset_path(
        assets.get("mermaid") or assets.get("mermaid_source"),
        outline_dir,
    )
    if not mermaid_source or not mermaid_source.exists():
        return diagram

    target = diagram or mermaid_source.with_suffix(".png")
    render_script = Path(__file__).resolve().parent / "render_mermaid.py"
    if not render_script.exists():
        return diagram
    try:
        result = subprocess.run(
            [
                sys.executable,
                str(render_script),
                "--input",
                str(mermaid_source),
                "--output",
                str(target),
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            check=False,
        )
        if result.returncode == 0 and target.exists():
            return target
    except Exception:
        return diagram
    return diagram


def _load_chart_payload(spec: dict[str, Any], outline_dir: Path) -> dict[str, Any]:
    candidates = [
        spec.get("chart"),
        _slide_assets(spec).get("chart_data"),
        _slide_assets(spec).get("chart"),
    ]
    for candidate in candidates:
        if isinstance(candidate, dict):
            return deepcopy(candidate)
        if not candidate:
            continue
        path = _resolve_asset_path(candidate, outline_dir)
        if not path:
            continue
        if not path.exists():
            raise FileNotFoundError(f"Chart data file not found: {path}")
        payload = json.loads(path.read_text(encoding="utf-8"))
        if isinstance(payload, dict):
            return payload
        raise ValueError(f"Chart data must decode to a JSON object: {path}")
    return {}


def _normalize_chart_payload(spec: dict[str, Any], outline_dir: Path) -> dict[str, Any]:
    raw = _load_chart_payload(spec, outline_dir)
    if not raw:
        return {}

    chart_type = _as_str(raw.get("type"), "bar").lower()
    options = raw.get("options") if isinstance(raw.get("options"), dict) else {}
    # Fix 1: accept both series-level `labels` (legacy schema) and chart-level
    # `categories` (common shorthand). If series omit `labels`, copy
    # `categories` in as a fallback before validation.
    chart_level_categories = None
    raw_categories = raw.get("categories")
    if isinstance(raw_categories, list) and raw_categories:
        chart_level_categories = raw_categories
    elif isinstance(raw.get("labels"), list) and raw.get("labels"):
        chart_level_categories = raw.get("labels")

    def _error(reason: str) -> dict[str, Any]:
        """Return an error marker so the renderer can draw a red banner and QA can pick it up."""
        return {
            "__error__": reason,
            "type": chart_type,
            "title": _as_str(raw.get("title"), ""),
            "subtitle": _as_str(raw.get("subtitle"), ""),
            "notes": _as_str(raw.get("notes") or raw.get("message") or raw.get("caption"), ""),
            "sources": raw.get("sources") if isinstance(raw.get("sources"), list) else [],
            "facts": raw.get("facts") if isinstance(raw.get("facts"), list) else raw.get("stats"),
            "series": [],
            "options": options,
            "color1": _as_str(raw.get("color1"), ""),
            "color2": _as_str(raw.get("color2"), ""),
            "color3": _as_str(raw.get("color3"), ""),
            "color4": _as_str(raw.get("color4"), ""),
        }

    series_items = raw.get("series")
    normalized_series: list[dict[str, Any]] = []
    invalid_reasons: list[str] = []
    if isinstance(series_items, list) and series_items:
        for index, item in enumerate(series_items):
            if not isinstance(item, dict):
                invalid_reasons.append(
                    f"series[{index}] is not an object"
                )
                continue
            series_labels = item.get("labels") if isinstance(item.get("labels"), list) else None
            # Fallback priority: series.labels -> chart.categories -> chart.labels
            labels = series_labels if series_labels else chart_level_categories
            values = item.get("values")
            if not isinstance(labels, list) or not labels:
                invalid_reasons.append(
                    f"series[{index}] missing labels/categories"
                )
                continue
            if not isinstance(values, list) or not values:
                invalid_reasons.append(
                    f"series[{index}] missing values"
                )
                continue
            if len(labels) != len(values):
                invalid_reasons.append(
                    f"series[{index}] length mismatch: {len(labels)} labels vs {len(values)} values"
                )
                continue
            parsed_values = [_safe_float(value) for value in values]
            if any(value is None for value in parsed_values):
                invalid_reasons.append(
                    f"series[{index}] contains non-numeric values"
                )
                continue
            paired = [
                (str(label).strip(), value)
                for label, value in zip(labels, parsed_values)
                if str(label).strip() and value is not None
            ]
            if not paired:
                invalid_reasons.append(
                    f"series[{index}] has no usable (label, value) pairs"
                )
                continue
            normalized_series.append(
                {
                    "name": _as_str(item.get("name"), f"Series {index + 1}"),
                    "labels": [label for label, _ in paired],
                    "values": [float(value) for _, value in paired],
                }
            )
    else:
        # Flat schema: {labels: [...], values: [...]} or {categories: [...], values: [...]}
        labels = chart_level_categories
        values = raw.get("values")
        if isinstance(labels, list) and isinstance(values, list):
            if len(labels) != len(values):
                invalid_reasons.append(
                    f"flat chart length mismatch: {len(labels)} labels vs {len(values)} values"
                )
            else:
                parsed_values = [_safe_float(value) for value in values]
                if any(value is None for value in parsed_values):
                    invalid_reasons.append("flat chart contains non-numeric values")
                else:
                    paired = [
                        (str(label).strip(), value)
                        for label, value in zip(labels, parsed_values)
                        if str(label).strip() and value is not None
                    ]
                    if paired:
                        normalized_series.append(
                            {
                                "name": _as_str(raw.get("series_name"), "Series A"),
                                "labels": [label for label, _ in paired],
                                "values": [float(value) for _, value in paired],
                            }
                        )
        elif series_items is None:
            invalid_reasons.append("chart payload has no series and no flat labels/values")

    if not normalized_series:
        if invalid_reasons:
            return _error("; ".join(invalid_reasons))
        # Raw payload exists but contained nothing actionable.
        return _error("chart payload present but produced no series")

    return {
        "type": chart_type,
        "title": _as_str(raw.get("title"), ""),
        "subtitle": _as_str(raw.get("subtitle"), ""),
        "notes": _as_str(raw.get("notes") or raw.get("message") or raw.get("caption"), ""),
        "sources": raw.get("sources") if isinstance(raw.get("sources"), list) else [],
        "facts": raw.get("facts") if isinstance(raw.get("facts"), list) else raw.get("stats"),
        "series": normalized_series,
        "options": options,
        "color1": _as_str(raw.get("color1"), ""),
        "color2": _as_str(raw.get("color2"), ""),
        "color3": _as_str(raw.get("color3"), ""),
        "color4": _as_str(raw.get("color4"), ""),
    }


def _normalize_fact_items(spec: dict[str, Any], chart_payload: dict[str, Any] | None = None) -> list[dict[str, str]]:
    collections: list[Any] = []
    for owner in (spec, chart_payload or {}):
        for key in ("facts", "stats", "evidence"):
            value = owner.get(key)
            if isinstance(value, list) and value:
                collections = value
                break
        if collections:
            break

    normalized: list[dict[str, str]] = []
    for item in collections:
        if isinstance(item, str):
            text = item.strip()
            if not text:
                continue
            normalized.append(
                {
                    "value": "",
                    "label": text,
                    "detail": "",
                    "source": "",
                    "accent": "",
                }
            )
            continue
        if not isinstance(item, dict):
            continue
        value = _as_str(item.get("value") or item.get("metric"), "")
        label = _as_str(item.get("label") or item.get("title"), "")
        detail = _as_str(item.get("detail") or item.get("body") or item.get("text"), "")
        source = _as_str(item.get("source") or item.get("citation"), "")
        accent = _as_str(item.get("accent"), "")
        if not any((value, label, detail, source)):
            continue
        if not label:
            label = detail or value or "Key Fact"
            if detail and value:
                detail = detail
            elif detail == label:
                detail = ""
        normalized.append(
            {
                "value": value,
                "label": label,
                "detail": detail if detail != label else "",
                "source": source,
                "accent": accent,
            }
        )
    return normalized[:4]


def _add_visual_assets(
    slide: Any,
    *,
    spec: dict[str, Any],
    slide_w: float,
    margin_x: float,
    outline_dir: Path,
) -> None:
    hero = _asset_from_spec(spec, "hero_image", outline_dir)
    if hero and hero.exists():
        box_w = 3.25
        box_h = 2.05
        x = slide_w - margin_x - box_w
        y = 1.42
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(box_w),
            Inches(box_h),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = _rgb("FFFFFF")
        frame.line.fill.background()
        try:
            frame.adjustments[0] = 0.08
        except Exception:
            pass
        slide.shapes.add_picture(
            str(hero),
            Inches(x + 0.02),
            Inches(y + 0.02),
            width=Inches(box_w - 0.04),
            height=Inches(box_h - 0.04),
        )

    logo = _asset_from_spec(spec, "logo", outline_dir)
    if logo and logo.exists():
        slide.shapes.add_picture(
            str(logo),
            Inches(slide_w - margin_x - 1.10),
            Inches(0.20),
            width=Inches(0.95),
            height=Inches(0.34),
        )


def _add_flow_diagram_layout(
    prs: Presentation,
    spec: dict[str, Any],
    *,
    preset: Any,
    outline_dir: Path,
) -> Any:
    typography = preset.typography
    palette = preset.palette
    layout = preset.layout
    slide = prs.slides.add_slide(_blank_layout(prs))
    sw = _slide_w(prs)
    sh = _slide_h(prs)
    mx = layout.margin_x
    _set_background(slide, palette["bg_primary"])
    _content_header(
        slide,
        title=_as_str(spec.get("title"), "Flow"),
        subtitle=_as_str(spec.get("subtitle"), ""),
        palette=palette,
        typography=typography,
        preset=preset,
        margin_x=mx,
    )

    diagram = _render_mermaid_asset_if_needed(spec, outline_dir)
    diagram_y = 1.45
    diagram_h = 4.95
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(mx),
        Inches(diagram_y),
        Inches(sw - mx * 2),
        Inches(diagram_h),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(palette["surface"])
    frame.line.color.rgb = _rgb(palette["line"])
    frame.line.width = Pt(1.2)
    try:
        frame.adjustments[0] = 0.03
    except Exception:
        pass

    if diagram and diagram.exists():
        slide.shapes.add_picture(
            str(diagram),
            Inches(mx + 0.08),
            Inches(diagram_y + 0.08),
            width=Inches(sw - mx * 2 - 0.16),
            height=Inches(diagram_h - 0.16),
        )
    else:
        _set_text_box(
            slide,
            text="Diagram asset missing. Provide assets.diagram for flow slides.",
            x=mx + 0.2,
            y=diagram_y + 2.1,
            w=sw - mx * 2 - 0.4,
            h=0.8,
            font_name=_font_body(preset),
            font_size=_clamp_body_font(typography, typography.body_max),
            color_hex=palette["text_muted"],
            align=PP_ALIGN.CENTER,
        )

    message = _as_str(spec.get("message"), "")
    caption = _as_str(spec.get("caption"), "")
    if not message:
        message = "Decision: Confirm the process owner and transition criteria."
    if not caption:
        caption = "Diagram source should map directly to the above decision."
    _set_text_box(
        slide,
        text=message,
        x=mx,
        y=6.48,
        w=sw - mx * 2,
        h=0.35,
        font_name=_font_title(preset),
        font_size=max(typography.caption_min + 1, 12),
        color_hex=palette["text_primary"],
        bold=True,
        margin=0.01,
    )
    _set_text_box(
        slide,
        text=caption,
        x=mx,
        y=6.83,
        w=sw - mx * 2,
        h=0.35,
        font_name=_font_caption(preset),
        font_size=max(typography.caption_min, 11),
        color_hex=palette["text_muted"],
        margin=0.01,
    )
    return slide


def _add_summary_callout(
    slide: Any,
    *,
    text: str,
    palette: dict[str, str],
    preset: Any,
    slide_w: float,
    slide_h: float,
    margin_x: float,
    has_footer: bool,
) -> float:
    """Render a rounded-rectangle summary callout anchored at the bottom.

    Inspired by the TB-LAMP-Seq / lab-data deck pattern: a full-width
    rounded box just above the footer carrying a single bold prose line
    that answers "so what?" for the slide. Universal — sits below any
    variant. Returns the Y-coordinate the callout occupies so callers
    can account for the space (most callers don't need to, since content
    zone is sized independently and the callout overlays the slide's
    bottom padding area).
    """
    body = (text or "").strip()
    if not body:
        return 0.0
    # Bottom anchoring: callout sits just above the footer if present,
    # otherwise just above the slide bottom-safe margin.
    footer_reserve = 0.50 if has_footer else 0.25
    callout_h = 0.78
    callout_y = slide_h - footer_reserve - callout_h
    callout_w = slide_w - margin_x * 2.2
    callout_x = margin_x * 1.1
    # Use ROUNDED_RECTANGLE so the shape reads as a pill/callout, not a card.
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(callout_x),
        Inches(callout_y),
        Inches(callout_w),
        Inches(callout_h),
    )
    # Deep accent fill + light text — the callout should visually punch
    # above the body text, like a pull-quote.
    accent_color = palette.get("accent_primary", "14B8A6")
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(accent_color)
    box.line.fill.background()
    _set_text_box(
        slide,
        text=body,
        x=callout_x + 0.25,
        y=callout_y + 0.08,
        w=callout_w - 0.50,
        h=callout_h - 0.16,
        font_name=_font_body(preset),
        font_size=max(13, min(16, preset.typography.body_max)),
        color_hex="FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0.04,
    )
    return callout_h


def _add_footer(
    slide: Any,
    *,
    text: str,
    palette: dict[str, str],
    preset: Any,
    slide_h: float,
    margin_x: float,
    dark: bool = False,
    page_index: int | None = None,
    page_total: int | None = None,
    slide_w: float = 13.33,
) -> None:
    footer_y, footer_h = _footer_geometry(slide_h)
    color = "94A3B8" if dark else palette["text_muted"]
    show_page = (
        isinstance(page_index, int)
        and isinstance(page_total, int)
        and page_total > 0
        and 1 <= page_index <= page_total
    )
    if not text.strip() and not show_page:
        return

    # Reserve a right-aligned slot for the page number when present so the
    # source line doesn't overlap it. Two text boxes: source on left,
    # page indicator on right. When only one is present, it spans the
    # full footer width.
    if show_page:
        page_text = f"{page_index} / {page_total}"
        page_w = 1.2
        page_x = slide_w - margin_x - page_w
        _set_text_box(
            slide,
            text=page_text,
            x=page_x,
            y=footer_y,
            w=page_w,
            h=footer_h,
            font_name=_font_caption(preset),
            font_size=10,
            color_hex=color,
            margin=0.0,
            align=PP_ALIGN.RIGHT,
            wrap=False,
        )
        left_w = max(1.0, page_x - margin_x - 0.20)
    else:
        left_w = slide_w - margin_x * 2

    if text.strip():
        font_size = 10 if len(text.strip()) <= 84 else 9
        _set_text_box(
            slide,
            text=text.strip(),
            x=margin_x,
            y=footer_y,
            w=min(10.5, left_w),
            h=footer_h,
            font_name=_font_caption(preset),
            font_size=font_size,
            color_hex=color,
            margin=0.0,
            wrap=False,
        )


# ---------------------------------------------------------------------------
# Fix 3: title-slide motif variety. All motifs share the same top-right
# footprint (~1.2" wide x ~1.0" tall, anchored around x ~= sw - 2.6, y ~= 0.5)
# and use only palette accent_primary/accent_secondary so they stay quiet.
# ---------------------------------------------------------------------------


def _draw_title_motif_orbit(slide: Any, *, palette: dict[str, str], slide_w: float) -> None:
    """Original three-orbit-dots motif (kept for backward compatibility)."""
    for index, scale in enumerate([0.9, 0.72, 0.55]):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(slide_w - (2.6 - index * 0.65)),
            Inches(0.7 + index * 0.42),
            Inches(scale),
            Inches(scale),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(
            palette["accent_primary"] if index != 1 else palette["accent_secondary"]
        )
        circle.line.fill.background()


def _draw_title_motif_stripe(slide: Any, *, palette: dict[str, str], slide_w: float) -> None:
    """Three horizontal accent stripes, descending widths, top-right corner."""
    stripes = [
        (1.20, 0.55, "accent_primary"),
        (0.90, 0.78, "accent_secondary"),
        (0.60, 1.01, "accent_primary"),
    ]
    right_edge = slide_w - 0.55
    for width, y_off, color_key in stripes:
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(right_edge - width),
            Inches(y_off),
            Inches(width),
            Inches(0.12),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _rgb(palette[color_key])
        stripe.line.fill.background()


def _draw_title_motif_wave(slide: Any, *, palette: dict[str, str], slide_w: float) -> None:
    """Quiet wave made from three small rounded-rectangle segments."""
    base_x = slide_w - 2.40
    # Three short stubby rounded rects at staggered heights to evoke a wave.
    segments = [
        (0.00, 0.80, 0.55, "accent_primary"),
        (0.55, 0.55, 0.55, "accent_secondary"),
        (1.10, 0.80, 0.55, "accent_primary"),
    ]
    for dx, y_off, width, color_key in segments:
        seg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(base_x + dx),
            Inches(y_off),
            Inches(width),
            Inches(0.22),
        )
        try:
            seg.adjustments[0] = 0.5
        except Exception:
            pass
        seg.fill.solid()
        seg.fill.fore_color.rgb = _rgb(palette[color_key])
        seg.line.fill.background()
    # Thin connector tint to tie them together
    tint = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(base_x),
        Inches(1.12),
        Inches(1.65),
        Inches(0.06),
    )
    tint.fill.solid()
    tint.fill.fore_color.rgb = _rgb(palette["accent_secondary"])
    tint.line.fill.background()


def _draw_title_motif_arc(slide: Any, *, palette: dict[str, str], slide_w: float) -> None:
    """Three concentric quarter-arcs (stacked ovals) in the top-right corner."""
    # Large to small ovals give a concentric-ring feel with no fills colliding.
    rings = [
        (1.30, "accent_primary"),
        (0.95, "accent_secondary"),
        (0.60, "accent_primary"),
    ]
    anchor_x = slide_w - 2.45
    anchor_y = 0.55
    for size, color_key in rings:
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(anchor_x + (1.30 - size) / 2),
            Inches(anchor_y + (1.30 - size) / 2),
            Inches(size),
            Inches(size),
        )
        ring.fill.background()
        ring.line.color.rgb = _rgb(palette[color_key])
        ring.line.width = Pt(2.25)


def _draw_title_motif_geometric(
    slide: Any, *, palette: dict[str, str], slide_w: float
) -> None:
    """Three overlapping small triangle/square shapes, top-right corner."""
    base_x = slide_w - 2.35
    # Triangle + square + triangle, overlapping slightly.
    tri1 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(base_x),
        Inches(0.55),
        Inches(0.65),
        Inches(0.72),
    )
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = _rgb(palette["accent_primary"])
    tri1.line.fill.background()

    sq = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(base_x + 0.48),
        Inches(0.72),
        Inches(0.55),
        Inches(0.55),
    )
    sq.fill.solid()
    sq.fill.fore_color.rgb = _rgb(palette["accent_secondary"])
    sq.line.fill.background()

    tri2 = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(base_x + 0.98),
        Inches(0.55),
        Inches(0.60),
        Inches(0.72),
    )
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = _rgb(palette["accent_primary"])
    tri2.line.fill.background()


_TITLE_MOTIFS: dict[str, Any] = {
    "orbit": _draw_title_motif_orbit,
    "stripe": _draw_title_motif_stripe,
    "wave": _draw_title_motif_wave,
    "arc": _draw_title_motif_arc,
    "geometric": _draw_title_motif_geometric,
}


# Keyword → motif mapping. First match wins, stable iteration order.
# Word-boundary matched against lowercased title_text + subtitle_text.
_TITLE_MOTIF_KEYWORDS: list[tuple[str, tuple[str, ...]]] = [
    ("wave", (
        "water", "flow", "ocean", "river", "stream", "tide", "current",
        "energy", "wave", "transition", "change", "shift",
    )),
    ("arc", (
        "journey", "arc", "bridge", "transform", "evolution", "trajectory",
        "growth", "path", "pivot",
    )),
    ("orbit", (
        "system", "network", "ecosystem", "platform", "cycle", "loop",
        "orbit", "galaxy", "cluster",
    )),
    ("stripe", (
        "process", "pipeline", "stages", "phases", "roadmap", "workflow",
        "sequence", "timeline", "stripe", "strategy",
    )),
    ("geometric", (
        "framework", "structure", "matrix", "grid", "architecture",
        "model", "blueprint", "schema",
    )),
]


def _pick_title_motif(
    *,
    title_text: str,
    deck_style: DeckStyleConfig,
    subtitle_text: str = "",
) -> str:
    """Pick motif name using tiered resolution:

    1. Explicit override (`deck_style.title_motif`) — "none" skips motif
       entirely, any valid choice wins.
    2. Keyword match against title+subtitle (word-boundary regex).
    3. Fallback: crc32(title_text) % 5 — keeps decks with generic titles
       from rendering a barren title slide.

    Returns "" when the caller should skip drawing a motif.
    """
    override = (deck_style.title_motif or "").strip().lower()
    if override == "none":
        return ""
    if override in TITLE_MOTIF_CHOICES:
        return override

    combined = f"{title_text or ''} {subtitle_text or ''}".lower()
    if combined.strip():
        # Word-boundary match avoids false hits like "processor" → stripe.
        for motif_name, words in _TITLE_MOTIF_KEYWORDS:
            pattern = r"\b(?:" + "|".join(re.escape(w) for w in words) + r")\b"
            if re.search(pattern, combined):
                return motif_name

    key = (title_text or "").encode("utf-8")
    idx = zlib.crc32(key) % len(TITLE_MOTIF_CHOICES)
    return TITLE_MOTIF_CHOICES[idx]


def _add_title_slide(
    prs: Presentation,
    spec: dict[str, Any],
    preset: Any,
    *,
    config: BuildConfig,
) -> Any:
    typography = preset.typography
    palette = preset.palette
    slide = prs.slides.add_slide(_blank_layout(prs))
    sw = _slide_w(prs)
    sh = _slide_h(prs)
    margin_x = preset.layout.margin_x

    _set_background(slide, palette["bg_dark"])

    left_rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.16),
        Inches(sh),
    )
    left_rail.fill.solid()
    left_rail.fill.fore_color.rgb = _rgb(palette["accent_primary"])
    left_rail.line.fill.background()

    title_text = str(spec.get("title", "Untitled Presentation")).strip()
    _set_text_box(
        slide,
        text=title_text,
        x=margin_x + 0.1,
        y=2.0,
        w=sw - margin_x * 2 - 2.7,
        h=2.0,
        font_name=_font_title(preset),
        font_size=_clamp_title_font(
            typography,
            title_text,
            density=config.deck_style.visual_density,
        ),
        color_hex="FFFFFF",
        bold=True,
        margin=0.01,
    )

    subtitle = str(spec.get("subtitle", "")).strip()
    if subtitle:
        subtitle_w = sw - margin_x * 2 - 2.8
        subtitle_font = max(typography.section_min, 22)
        subtitle_h = min(
            1.15,
            max(
                0.6,
                _estimate_text_height(
                    [subtitle],
                    width_in=subtitle_w,
                    font_size_pt=subtitle_font,
                )
                + 0.10,
            ),
        )
        _set_text_box(
            slide,
            text=subtitle,
            x=margin_x + 0.1,
            y=4.55,
            w=subtitle_w,
            h=subtitle_h,
            font_name=_font_body(preset),
            font_size=subtitle_font,
            color_hex=palette["accent_primary"],
            bold=True,
            margin=0.01,
        )

    # Fix 3: hash-rotated title motif variety. If a hero image is present,
    # the hero code path has already drawn it (or will draw it) and we skip
    # the decorative motif entirely so it doesn't compete. Otherwise the
    # motif is picked deterministically from the deck title, with optional
    # explicit override from `deck_style.title_motif`.
    hero_asset = _asset_from_spec(spec, "hero_image", config.outline_dir)
    has_hero = bool(hero_asset and hero_asset.exists())
    if not has_hero:
        motif_choice = _pick_title_motif(
            title_text=title_text,
            deck_style=config.deck_style,
            subtitle_text=str(spec.get("subtitle", "")),
        )
        # Empty string = explicit opt-out via deck_style.title_motif: "none".
        if motif_choice:
            motif_fn = _TITLE_MOTIFS.get(motif_choice, _draw_title_motif_orbit)
            motif_fn(slide, palette=palette, slide_w=sw)

    # Bug 4 fix: Only draw footer when author actually provided one.
    # Do not fall back to "Prepared deck" or similar tells.
    title_footer_text = _slide_footer_text(spec)
    if not title_footer_text:
        date_value = _as_str(spec.get("date"), "")
        if date_value:
            title_footer_text = date_value
    if title_footer_text:
        _add_footer(
            slide,
            text=title_footer_text,
            palette=palette,
            preset=preset,
            slide_h=sh,
            margin_x=margin_x + 0.1,
            dark=True,
        )
    return slide


def _add_section_slide(
    prs: Presentation,
    spec: dict[str, Any],
    preset: Any,
    *,
    config: BuildConfig,
) -> Any:
    typography = preset.typography
    palette = preset.palette
    slide = prs.slides.add_slide(_blank_layout(prs))
    sw = _slide_w(prs)
    sh = _slide_h(prs)
    mx = preset.layout.margin_x

    _set_background(slide, palette["bg_primary"])
    _add_content_motif(slide, palette=palette, slide_w=sw, margin_x=mx)

    title_text = str(spec.get("title", "Section")).strip()
    _set_text_box(
        slide,
        text=title_text,
        x=mx,
        y=1.45,
        w=sw - mx * 2,
        h=1.0,
        font_name=_font_title(preset),
        font_size=max(
            typography.title_min,
            _clamp_title_font(typography, title_text, density=config.deck_style.visual_density) - 2,
        ),
        color_hex=palette["text_primary"],
        bold=True,
        margin=0.01,
    )
    subtitle = str(spec.get("subtitle", "")).strip()
    if subtitle:
        _set_text_box(
            slide,
            text=subtitle,
            x=mx,
            y=2.55,
            w=sw - mx * 2,
            h=0.7,
            font_name=_font_body(preset),
            font_size=_clamp_body_font(typography, typography.body_max),
            color_hex=palette["text_muted"],
            margin=0.01,
        )

    # Bug 6 enrichment: section dividers tend to leave the bottom ~75% of
    # the canvas empty. When the author has not provided bullets, body,
    # or caption, drop in a quiet decorative footprint (accent rail +
    # orbit dots) below the subtitle stack. This is enrichment, not
    # content, so the styling is deliberately low-contrast.
    has_section_content = (
        bool(str(spec.get("body", "")).strip())
        or bool(str(spec.get("caption", "")).strip())
        or (isinstance(spec.get("bullets"), list) and any(
            str(item if not isinstance(item, dict) else item.get("text", "")).strip()
            for item in spec.get("bullets", [])
        ))
    )
    if not has_section_content:
        _add_section_motif_enrichment(
            slide,
            palette=palette,
            slide_w=sw,
            slide_h=sh,
            margin_x=mx,
        )

    _add_footer(
        slide,
        text=_slide_footer_text(spec),
        palette=palette,
        preset=preset,
        slide_h=sh,
        margin_x=mx,
    )
    return slide


def _add_standard_content(
    slide: Any,
    *,
    lines: list[str],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    intro_paragraph: str | None = None,
    card_heading: str = "",
) -> None:
    x = layout.margin_x
    y = content_top
    w = slide_w - x * 2
    intro_text = (intro_paragraph or "").strip()

    # Fix 2b: single-card standard variant should expand the card to fill
    # the available content zone (content_top -> content_bottom). This gives
    # the slide a confident, deliberate block rather than a short card
    # floating above dead space. Fix 1: pair with MIDDLE vertical anchor so
    # bullets don't orphan at the top when they don't fill the card.
    # We cap at 85% of the available height to leave breathing room below
    # the card and stay within the preset's max_density budget (the density
    # lint counts card area against whole-slide area; a 100% fill plus header
    # and footer pushes density past most presets' 0.80 cap).
    available_h = max(1.2, content_bottom - y)
    expand_cap = max(2.5, min(available_h * 0.85, 4.6))

    # Card-heading resolution: outline can set `card_heading` explicitly.
    # Default: "" (no card heading) — the slide's own header already
    # frames the card, and a stale default like "Key Points" reads as
    # auto-generated filler. The caller in _add_content_slide passes
    # `card_heading` through from the outline spec.
    card_heading = (card_heading or "").strip()
    if intro_text:
        preferred = _preferred_card_height_with_intro(
            title_text=card_heading or "Key Points",  # measure uses a non-empty string
            intro_text=intro_text,
            body_lines=lines[:6],
            width_in=w,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            min_h=2.25,
            max_h=expand_cap,
        )
        # Fix 2b: grow to fill the zone, capping at the zone height minus a
        # small breathing margin so the card doesn't crowd the bottom edge.
        h = max(preferred, expand_cap)
        _add_card_with_intro(
            slide,
            x=x,
            y=y,
            w=w,
            h=h,
            title=card_heading,
            intro_text=intro_text,
            body_lines=lines[:6],
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key="accent_primary",
            body_vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    preferred = _preferred_card_height(
        title_text=card_heading or "Key Points",
        body_lines=lines[:6],
        width_in=w,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        min_h=2.25,
        max_h=expand_cap,
    )
    h = max(preferred, expand_cap)
    _add_card(
        slide,
        x=x,
        y=y,
        w=w,
        h=h,
        title=card_heading,
        body_lines=lines[:6],
        palette=palette,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        accent_key="accent_primary",
        body_vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def _preferred_card_height_with_intro(
    *,
    title_text: str,
    intro_text: str,
    body_lines: list[str],
    width_in: float,
    typography: Any,
    preset: Any,
    rail_h: float,
    min_h: float,
    max_h: float,
) -> float:
    heading_font, title_h = _card_title_layout(title_text, width_in, typography, preset)
    usable_width = max(0.6, width_in - 0.28)
    intro_font = _clamp_body_font(typography, typography.body_max - 1)
    intro_h = _estimate_text_height([intro_text], width_in=usable_width, font_size_pt=intro_font)
    body_font = _card_body_font(body_lines, width_in=width_in, typography=typography)
    body_h = _estimate_text_height(body_lines[:6], width_in=usable_width, font_size_pt=body_font)
    preferred = rail_h + title_h + 0.16 + intro_h + 0.14 + body_h + 0.24
    return min(max(max(min_h, preferred), 0.9), max_h)


def _add_card_with_intro(
    slide: Any,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    intro_text: str,
    body_lines: list[str],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    rail_h: float,
    accent_key: str = "accent_primary",
    dark: bool = False,
    body_vertical_anchor: Any = None,
) -> None:
    fill_color = palette["bg_dark"] if dark else palette["surface"]
    line_color = palette["line"] if not dark else palette["accent_primary"]
    text_color = "FFFFFF" if dark else palette["text_primary"]
    muted_color = "CBD5E1" if dark else palette["text_muted"]

    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(fill_color)
    card.line.color.rgb = _rgb(line_color)
    card.line.width = Pt(1.25)

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(rail_h),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = _rgb(palette.get(accent_key, palette["accent_primary"]))
    rail.line.fill.background()

    # Same empty-title-suppression behavior as _add_card: an empty title
    # skips the title row so the card's top decoration is just the rail.
    # Callers use this when the slide's own header already frames the
    # card (e.g., standard variant with a subtitle present).
    title_text = title.strip()
    suppress_title = title_text == ""
    title_font_name = _font_title(preset)
    if suppress_title:
        title_h = 0.0
    else:
        heading_font, title_h = _card_title_layout(title_text, w, typography, preset)
        _set_text_box(
            slide,
            text=title_text,
            x=x + 0.14,
            y=y + rail_h + 0.06,
            w=w - 0.28,
            h=max(0.36, title_h - 0.02),
            font_name=title_font_name,
            font_size=heading_font,
            color_hex=text_color,
            bold=True,
            margin=0.01,
        )

    usable_width = max(0.6, w - 0.28)
    intro_font = _clamp_body_font(typography, typography.body_max - 1)
    intro_h_est = _estimate_text_height([intro_text], width_in=usable_width, font_size_pt=intro_font)
    available_after_title = max(0.5, h - rail_h - title_h - 0.16 - 0.24)
    intro_h = min(max(0.32, intro_h_est + 0.04), available_after_title * 0.55)
    intro_y = y + rail_h + (title_h + 0.16 if not suppress_title else 0.12)
    _set_text_box(
        slide,
        text=intro_text,
        x=x + 0.14,
        y=intro_y,
        w=usable_width,
        h=intro_h,
        font_name=_font_body(preset),
        font_size=intro_font,
        color_hex=text_color,
        margin=0.01,
    )

    body_y = intro_y + intro_h + 0.10
    body_h = max(0.5, (y + h) - body_y - 0.12)
    body_font = _card_body_font(
        body_lines,
        width_in=w,
        typography=typography,
        available_height_in=body_h,
    )
    body_box = _set_lines_box(
        slide,
        lines=body_lines[:6],
        x=x + 0.14,
        y=body_y,
        w=usable_width,
        h=body_h,
        font_name=_font_body(preset),
        font_size=body_font,
        color_hex=muted_color,
        bullet=True,
        margin=0.01,
    )
    # Fix 1/2b: optionally center bullets vertically when the card has been
    # expanded to fill the content zone and the bullets don't fill the gap.
    if body_vertical_anchor is not None and body_box is not None:
        try:
            body_box.text_frame.vertical_anchor = body_vertical_anchor
        except Exception:
            pass


def _add_hero_content(
    slide: Any,
    *,
    lines: list[str],
    hero_image: Path | None,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    image_w = 3.25
    gutter = layout.gutter
    text_w = max(4.8, total_w - image_w - gutter)
    text_lines = (lines or ["Add core mission update."])[:6]
    card_h = _preferred_card_height(
        title_text="Key Message",
        body_lines=text_lines,
        width_in=text_w,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        min_h=2.35,
        max_h=max(2.8, min(content_bottom - y, 3.55)),
    )
    _add_card(
        slide,
        x=x,
        y=y,
        w=text_w,
        h=card_h,
        title="Key Message",
        body_lines=text_lines,
        palette=palette,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        accent_key="accent_primary",
    )

    image_h = min(max(2.2, card_h), max(2.2, content_bottom - y))
    frame_x = x + text_w + gutter
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(frame_x),
        Inches(y),
        Inches(image_w),
        Inches(image_h),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(palette["surface"])
    frame.line.color.rgb = _rgb(palette["line"])
    frame.line.width = Pt(1.0)
    try:
        frame.adjustments[0] = 0.05
    except Exception:
        pass
    if hero_image and hero_image.exists():
        slide.shapes.add_picture(
            str(hero_image),
            Inches(frame_x + 0.03),
            Inches(y + 0.03),
            width=Inches(image_w - 0.06),
            height=Inches(image_h - 0.06),
        )
    else:
        _set_text_box(
            slide,
            text="Provide assets.hero_image or asset:<alias> for the hero visual.",
            x=frame_x + 0.14,
            y=y + max(0.6, image_h / 2.0 - 0.25),
            w=image_w - 0.28,
            h=0.5,
            font_name=_font_body(preset),
            font_size=max(typography.caption_min, 10),
            color_hex=palette["text_muted"],
            align=PP_ALIGN.CENTER,
            margin=0.01,
        )


def _image_contain_geometry(image_path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    try:
        from PIL import Image as PILImage

        with PILImage.open(image_path) as im:
            img_w, img_h = im.size
        if img_w <= 0 or img_h <= 0:
            return x, y, w, h
        box_ratio = w / max(h, 0.01)
        img_ratio = img_w / img_h
        if img_ratio >= box_ratio:
            fit_w = w
            fit_h = w / img_ratio
        else:
            fit_h = h
            fit_w = h * img_ratio
        return x + (w - fit_w) / 2.0, y + (h - fit_h) / 2.0, fit_w, fit_h
    except Exception:
        return x, y, w, h


def _generated_image_metadata(image_path: Path | None, spec: dict[str, Any]) -> dict[str, str]:
    metadata: dict[str, str] = {}
    if image_path is not None:
        metadata_path = Path(f"{image_path}.metadata.json")
        if metadata_path.exists():
            try:
                raw = json.loads(metadata_path.read_text(encoding="utf-8"))
                if isinstance(raw, dict):
                    for key in ("model", "prompt", "revised_prompt", "purpose", "edit_note", "created_at"):
                        value = raw.get(key)
                        if value is not None:
                            metadata[key] = str(value)
            except Exception:
                pass
    slide_meta = spec.get("image_generation")
    if isinstance(slide_meta, dict):
        for key in ("model", "prompt", "revised_prompt", "purpose", "edit_note", "created_at"):
            value = slide_meta.get(key)
            if value is not None and str(value).strip():
                metadata[key] = str(value)
    return metadata


def _add_generated_image_content(
    slide: Any,
    *,
    image_path: Path | None,
    spec: dict[str, Any],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    panel_w = 3.0
    gutter = layout.gutter
    image_w = total_w - panel_w - gutter
    content_h = max(2.7, content_bottom - y)

    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(image_w),
        Inches(content_h),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(palette["surface"])
    frame.line.color.rgb = _rgb(palette["line"])
    frame.line.width = Pt(1.0)

    if image_path and image_path.exists():
        ix, iy, iw, ih = _image_contain_geometry(image_path, x + 0.08, y + 0.08, image_w - 0.16, content_h - 0.16)
        slide.shapes.add_picture(str(image_path), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))
    else:
        _set_text_box(
            slide,
            text="Generated image asset missing. Rebuild with --allow-generated-images or replace this slide.",
            x=x + 0.35,
            y=y + content_h / 2.0 - 0.35,
            w=image_w - 0.70,
            h=0.70,
            font_name=_font_body(preset),
            font_size=max(typography.body_min, 14),
            color_hex=palette["text_muted"],
            align=PP_ALIGN.CENTER,
            margin=0.01,
        )

    meta = _generated_image_metadata(image_path, spec)
    panel_x = x + image_w + gutter
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(panel_x),
        Inches(y),
        Inches(panel_w),
        Inches(content_h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(palette["bg_dark"])
    panel.line.fill.background()

    label_color = palette["accent_primary"]
    _set_text_box(
        slide,
        text="GENERATED VISUAL",
        x=panel_x + 0.20,
        y=y + 0.22,
        w=panel_w - 0.40,
        h=0.30,
        font_name=_font_title(preset),
        font_size=12,
        color_hex=label_color,
        bold=True,
        margin=0.01,
    )
    details = [
        f"Model: {meta.get('model', 'OpenAI image model')}",
        f"Purpose: {meta.get('purpose', 'Concept visual')}",
        "Delete this slide if source-backed imagery is preferred.",
    ]
    edit_note = meta.get("edit_note", "").strip()
    if edit_note:
        details.append(f"Edit note: {edit_note}")
    prompt = meta.get("prompt") or meta.get("revised_prompt") or ""
    if prompt:
        details.append(f"Prompt: {prompt[:280]}")
    _set_lines_box(
        slide,
        lines=details,
        x=panel_x + 0.20,
        y=y + 0.66,
        w=panel_w - 0.40,
        h=content_h - 0.88,
        font_name=_font_body(preset),
        font_size=12,
        color_hex="FFFFFF",
        bullet=False,
        margin=0.01,
    )


def _add_split_content(
    slide: Any,
    *,
    lines: list[str],
    highlights: list[str],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    left_w = total_w * 0.58
    right_w = total_w - left_w - layout.gutter
    left_lines = (lines or ["Add main narrative point."])[:6]
    right_lines = (highlights or lines[:4] or ["Add checklist item"])[:6]
    # Fix 2b: the left column of split is a single "primary" card and should
    # fill the available content zone. The right column (checklist) keeps its
    # compact, measured-height behavior. Like the standard variant we leave
    # breathing room so the preset's density budget still passes.
    available_h = max(1.6, content_bottom - y)
    expand_cap = max(2.8, min(available_h * 0.85, 4.6))
    left_preferred = _preferred_card_height(
        title_text="Main Narrative",
        body_lines=left_lines,
        width_in=left_w,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        min_h=2.1,
        max_h=expand_cap,
    )
    right_h = _preferred_card_height(
        title_text="Checklist",
        body_lines=right_lines,
        width_in=right_w,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        min_h=1.5,
        max_h=max(2.25, min(content_bottom - y, 3.1)),
    )
    left_box_h = max(left_preferred, expand_cap)
    compact_right = len(right_lines) <= 2 and max((len(line) for line in right_lines), default=0) <= 34
    if compact_right:
        right_box_h = min(max(1.55, right_h), max(1.8, left_box_h * 0.62))
    else:
        # Let the right card grow to match the left when it has substantive
        # content; cap it at the available zone so it still reads as a card.
        right_box_h = min(left_box_h, max(right_h, min(left_box_h, available_h)))
    # Fix 2: anchor body text to the TOP of split cards so the gap between the
    # card title and the body matches the normal padding. Previously the left
    # card used MSO_ANCHOR.MIDDLE, which leaves a visible half-card gap between
    # title and body when the card was sized for a density budget taller than
    # the body text. Both cards now anchor-top for a consistent rhythm.
    _add_card(
        slide,
        x=x,
        y=y,
        w=left_w,
        h=left_box_h,
        title="Main Narrative",
        body_lines=left_lines,
        palette=palette,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        accent_key="accent_primary",
        body_vertical_anchor=MSO_ANCHOR.TOP,
    )
    _add_card(
        slide,
        x=x + left_w + layout.gutter,
        y=y,
        w=right_w,
        h=right_box_h,
        title="Checklist",
        body_lines=right_lines,
        palette=palette,
        typography=typography,
        preset=preset,
        rail_h=layout.rail_height,
        accent_key="accent_secondary",
        body_vertical_anchor=MSO_ANCHOR.TOP,
    )


def _add_cards_grid_promoted(
    slide: Any,
    *,
    cards: list[dict[str, Any]],
    promote_index: int,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    icon_paths: list[Path | None] | None = None,
) -> None:
    """Asymmetric cards-3 layout: one card doubles as a left-column anchor,
    the other two stack on the right. Breaks the symmetric 3-up grid.
    """
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    gutter = layout.gutter
    # Left column (promoted): 2/3 width - gutter/2; right column: 1/3 width.
    left_w = total_w * 0.60 - gutter / 2.0
    right_w = total_w * 0.40 - gutter / 2.0
    total_h = max(2.6, content_bottom - y)
    small_h = (total_h - gutter) / 2.0
    icons = icon_paths or []
    icon_size = 0.5
    icon_pad = 0.15

    promote_index = max(0, min(2, promote_index))
    order = [promote_index] + [i for i in range(3) if i != promote_index]
    positions = [
        (x, y, left_w, total_h),                                   # big left card
        (x + left_w + gutter, y, right_w, small_h),                # top right
        (x + left_w + gutter, y + small_h + gutter, right_w, small_h),  # bottom right
    ]
    dark_bg = palette["bg_primary"].upper() == palette["bg_dark"].upper()

    for slot, idx in enumerate(order):
        card = cards[idx] if idx < len(cards) else {}
        title = str(card.get("title", f"Pillar {idx + 1}")).strip()
        body = card.get("body_lines") or card.get("body") or card.get("text") or ""
        if isinstance(body, list):
            body_lines = [str(item).strip() for item in body if str(item).strip()]
        else:
            body_lines = [
                part.strip() for part in str(body).replace("\n", ". ").split(".")
                if part.strip()
            ]
        if not body_lines:
            body_lines = [f"Define {title.lower()} execution detail."]
        accent_raw = str(card.get("accent", "")).strip()
        accent_key = accent_raw if accent_raw in palette else (
            "accent_primary" if slot == 0 else "accent_secondary"
        )
        icon_path = icons[idx] if idx < len(icons) else None
        cx, cy, cw, ch = positions[slot]
        # Promoted card gets more body lines; the small cards get tighter copy.
        max_lines = 7 if slot == 0 else 4
        _add_card(
            slide,
            x=cx,
            y=cy,
            w=cw,
            h=ch,
            title=title,
            body_lines=body_lines[:max_lines],
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key=accent_key,
            dark=dark_bg,
            icon_path=icon_path,
            icon_size=icon_size if icon_path else 0.0,
        )


def _add_cards_grid(
    slide: Any,
    *,
    cards: list[dict[str, Any]],
    columns: int,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    icon_paths: list[Path | None] | None = None,
    promote_card: int | None = None,
) -> None:
    columns = max(2, min(3, columns))
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    gutter = layout.gutter
    card_w = (total_w - gutter * (columns - 1)) / columns

    # Promoted-card asymmetric layout (cards-3 only). The promoted card
    # occupies the full left column (2×card_w + gutter); the other two
    # stack in the right column. Breaks the 3-up grid rigidity without
    # adding a new variant.
    if (
        columns == 3
        and promote_card is not None
        and 0 <= promote_card < len(cards)
        and len(cards) >= 3
    ):
        _add_cards_grid_promoted(
            slide,
            cards=cards,
            promote_index=int(promote_card),
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=slide_w,
            content_bottom=content_bottom,
            content_top=content_top,
            icon_paths=icon_paths,
        )
        return
    card_entries: list[tuple[str, list[str], str]] = []
    preferred_heights: list[float] = []
    min_card_h = 1.2 if columns == 2 else 1.15
    max_card_h = max(1.5, min(content_bottom - y, 2.6 if columns == 2 else 2.3))
    icon_size = 0.5
    icons = icon_paths or []
    icon_pad = 0.15
    for idx in range(columns):
        card = cards[idx] if idx < len(cards) else {}
        title = str(card.get("title", f"Pillar {idx + 1}")).strip()
        body = card.get("body_lines") or card.get("body") or card.get("text") or ""
        if isinstance(body, list):
            body_lines = [str(item).strip() for item in body if str(item).strip()]
        else:
            body_lines = [part.strip() for part in str(body).replace("\n", ". ").split(".") if part.strip()]
        if not body_lines:
            body_lines = [f"Define {title.lower()} execution detail."]
        accent = str(card.get("accent", "accent_primary")).strip() or "accent_primary"
        card_entries.append((title, body_lines[:5], accent if accent in palette else "accent_primary"))
        icon_here = idx < len(icons) and icons[idx] is not None
        icon_block = (icon_size + icon_pad) if icon_here else 0.0
        preferred_heights.append(
            _preferred_card_height(
                title_text=title,
                body_lines=body_lines[:5],
                width_in=card_w,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                min_h=min_card_h,
                max_h=max_card_h,
                icon_height=icon_block,
            )
        )
    card_h = max(preferred_heights) if preferred_heights else max(min_card_h, max_card_h)
    for idx, (title, body_lines, accent_key) in enumerate(card_entries):
        icon_path = icons[idx] if idx < len(icons) else None
        _add_card(
            slide,
            x=x + idx * (card_w + gutter),
            y=y,
            w=card_w,
            h=card_h,
            title=title,
            body_lines=body_lines,
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key=accent_key,
            dark=palette["bg_primary"].upper() == palette["bg_dark"].upper(),
            icon_path=icon_path,
            icon_size=icon_size if icon_path else 0.0,
        )


def _fact_body_lines(item: dict[str, str], *, include_source: bool = True) -> list[str]:
    lines: list[str] = []
    detail = _as_str(item.get("detail"), "")
    source = _as_str(item.get("source"), "")
    if detail:
        lines.extend([part.strip() for part in detail.replace("\n", ". ").split(".") if part.strip()])
    if include_source and source:
        lines.append(f"Source: {source}")
    return lines[:3]


def _preferred_fact_card_height(
    item: dict[str, str],
    *,
    width_in: float,
    typography: Any,
    preset: Any,
    rail_h: float,
    min_h: float,
    max_h: float,
    include_source: bool = True,
    icon_height: float = 0.0,
) -> float:
    label = _as_str(item.get("label"), "Key Fact")
    detail_lines = _fact_body_lines(item, include_source=include_source)
    label_font = max(typography.body_min, typography.caption_min + 2, min(typography.section_min, 15))
    detail_font = _card_body_font(
        detail_lines or [label],
        width_in=width_in,
        typography=typography,
    )
    value = _as_str(item.get("value"), "")
    value_h = 0.74 if value else 0.0
    label_h = _estimate_text_height([label], width_in=max(0.8, width_in - 0.28), font_size_pt=label_font)
    detail_h = _estimate_text_height(detail_lines, width_in=max(0.8, width_in - 0.28), font_size_pt=detail_font)
    preferred = rail_h + icon_height + 0.14 + value_h + label_h + detail_h + 0.24
    return min(max(max(min_h, preferred), 1.0), max_h)


def _add_fact_card(
    slide: Any,
    *,
    item: dict[str, str],
    x: float,
    y: float,
    w: float,
    h: float,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    rail_h: float,
    accent_key: str,
    include_source: bool = True,
    icon_path: Path | None = None,
    icon_size: float = 0.0,
) -> None:
    accent = accent_key if accent_key in palette else "accent_primary"
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(palette["surface"])
    card.line.color.rgb = _rgb(palette["line"])
    card.line.width = Pt(1.2)

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(rail_h),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = _rgb(palette[accent])
    rail.line.fill.background()

    cursor_y = y + rail_h + 0.10

    # Icon: rendered centered above the value, inside the tile (below rail).
    icon_pad = 0.15 if (icon_path is not None and icon_size > 0) else 0.0
    if icon_path is not None and icon_size > 0:
        icon_left = x + (w - icon_size) / 2.0
        _add_icon(slide, path=icon_path, x=icon_left, y=cursor_y, size=icon_size)
        cursor_y += icon_size + icon_pad

    value = _as_str(item.get("value"), "")
    if value:
        value_font = max(typography.section_max - 1, typography.title_min)
        if w < 4.0:
            value_font = min(value_font, 20)
        _set_text_box(
            slide,
            text=value,
            x=x + 0.14,
            y=cursor_y,
            w=w - 0.28,
            h=0.62,
            font_name=_font_title(preset),
            font_size=value_font,
            color_hex=palette[accent],
            bold=True,
            margin=0.01,
        )
        cursor_y += 0.72

    label = _as_str(item.get("label"), "Key Fact")
    label_font = max(typography.body_min, typography.caption_min + 2, min(typography.section_min, 15))
    label_lines = max(1, min(2, _estimate_text_lines(label, width_in=max(0.8, w - 0.28), font_size_pt=label_font)))
    label_h = min(0.50, max(0.22, 0.12 + label_lines * 0.16))
    _set_text_box(
        slide,
        text=label,
        x=x + 0.14,
        y=cursor_y,
        w=w - 0.28,
        h=label_h,
        font_name=_font_body(preset),
        font_size=label_font,
        color_hex=palette["text_primary"],
        bold=True,
        margin=0.01,
    )
    cursor_y += label_h + 0.08

    detail_lines = _fact_body_lines(item, include_source=include_source)
    if detail_lines:
        body_h = max(0.30, h - (cursor_y - y) - 0.12)
        body_font = _card_body_font(
            detail_lines,
            width_in=max(0.8, w - 0.28),
            typography=typography,
            available_height_in=max(0.24, body_h - 0.02),
        )
        _set_lines_box(
            slide,
            lines=detail_lines,
            x=x + 0.14,
            y=cursor_y,
            w=w - 0.28,
            h=body_h,
            font_name=_font_body(preset),
            font_size=body_font,
            color_hex=palette["text_muted"],
            bullet=False,
            margin=0.01,
        )


def _add_stats_content(
    slide: Any,
    *,
    facts: list[dict[str, str]],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    icon_paths: list[Path | None] | None = None,
) -> None:
    items = facts[:4] if facts else [
        {"value": "42%", "label": "Primary outcome", "detail": "Add sourced evidence detail.", "source": ""},
        {"value": "6 mo", "label": "Time window", "detail": "Add timing or throughput context.", "source": ""},
        {"value": "$12M", "label": "Envelope", "detail": "Add bound or planning detail.", "source": ""},
    ]
    count = len(items)
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    gutter = layout.gutter
    columns = 2 if count >= 4 else max(1, count)
    rows = max(1, int(math.ceil(count / float(columns))))
    card_w = (total_w - gutter * (columns - 1)) / columns
    total_h = max(1.6, content_bottom - y)
    max_card_h = (total_h - gutter * (rows - 1)) / rows

    icons = icon_paths or []
    icon_size = 0.4
    icon_pad = 0.15

    row_heights = [1.25 for _ in range(rows)]
    prepared: list[tuple[int, int, dict[str, str]]] = []
    for index, item in enumerate(items):
        row = index // columns
        col = index % columns
        prepared.append((row, col, item))
        icon_here = index < len(icons) and icons[index] is not None
        icon_block = (icon_size + icon_pad) if icon_here else 0.0
        row_heights[row] = max(
            row_heights[row],
            _preferred_fact_card_height(
                item,
                width_in=card_w,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                min_h=1.22,
                max_h=max(1.45, min(2.2, max_card_h)),
                icon_height=icon_block,
            ),
        )

    row_h = min(max_card_h, max(row_heights))
    for idx_item, (row, col, item) in enumerate(prepared):
        accent = _as_str(item.get("accent"), "")
        if accent not in palette:
            accent = "accent_primary" if (row * columns + col) % 2 == 0 else "accent_secondary"
        icon_path = icons[idx_item] if idx_item < len(icons) else None
        _add_fact_card(
            slide,
            item=item,
            x=x + col * (card_w + gutter),
            y=y + row * (row_h + gutter),
            w=card_w,
            h=row_h,
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key=accent,
            icon_path=icon_path,
            icon_size=icon_size if icon_path else 0.0,
        )


def _chart_colors(chart_payload: dict[str, Any], palette: dict[str, str]) -> list[str]:
    colors: list[str] = []
    options = chart_payload.get("options") if isinstance(chart_payload.get("options"), dict) else {}
    raw_colors = options.get("chartColors")
    if isinstance(raw_colors, list):
        for value in raw_colors:
            color = _as_str(value, "").replace("#", "")
            if len(color) == 6:
                colors.append(color.upper())
    for key in ("color1", "color2", "color3", "color4"):
        color = _as_str(chart_payload.get(key), "").replace("#", "")
        if len(color) == 6:
            colors.append(color.upper())
    fallback = [
        palette["accent_primary"],
        palette["accent_secondary"],
        palette["text_muted"],
        palette["line"],
    ]
    colors.extend([item for item in fallback if item not in colors])
    return colors


def _chart_type_from_payload(chart_payload: dict[str, Any]) -> XL_CHART_TYPE:
    chart_type = _as_str(chart_payload.get("type"), "bar").lower()
    if chart_type == "line":
        return XL_CHART_TYPE.LINE_MARKERS
    if chart_type == "pie":
        return XL_CHART_TYPE.PIE
    return XL_CHART_TYPE.COLUMN_CLUSTERED


def _legend_position(value: str) -> XL_LEGEND_POSITION:
    token = value.strip().lower()
    if token == "l":
        return XL_LEGEND_POSITION.LEFT
    if token == "t":
        return XL_LEGEND_POSITION.TOP
    if token == "b":
        return XL_LEGEND_POSITION.BOTTOM
    return XL_LEGEND_POSITION.RIGHT


def _add_chart_frame(
    slide: Any,
    *,
    chart_payload: dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
) -> None:
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _rgb(palette["surface"])
    frame.line.color.rgb = _rgb(palette["line"])
    frame.line.width = Pt(1.2)

    # Fix 1: when chart payload arrived but failed schema validation, render
    # a bold red banner so the slide visibly flags the problem instead of
    # silently showing a generic "provide chart data" placeholder.
    error_reason = chart_payload.get("__error__") if isinstance(chart_payload, dict) else None
    if error_reason:
        banner_h = min(0.72, max(0.48, h * 0.22))
        banner_y = y + max(0.10, (h - banner_h) / 2.0)
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.20),
            Inches(banner_y),
            Inches(max(1.0, w - 0.40)),
            Inches(banner_h),
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = _rgb("B91C1C")
        banner.line.color.rgb = _rgb("7F1D1D")
        banner.line.width = Pt(1.0)
        _set_text_box(
            slide,
            text="Chart data malformed \u2014 see QA report",
            x=x + 0.28,
            y=banner_y + 0.02,
            w=max(0.8, w - 0.56),
            h=banner_h - 0.04,
            font_name=_font_body(preset),
            font_size=max(typography.caption_min + 2, 13),
            color_hex="FFFFFF",
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0.02,
        )
        return

    if not chart_payload.get("series"):
        _set_text_box(
            slide,
            text="Provide chart data as an inline chart object or a staged chart:<alias> JSON file.",
            x=x + 0.20,
            y=y + max(0.9, h / 2.0 - 0.22),
            w=w - 0.40,
            h=0.44,
            font_name=_font_body(preset),
            font_size=max(typography.caption_min + 1, 11),
            color_hex=palette["text_muted"],
            align=PP_ALIGN.CENTER,
            margin=0.01,
        )
        return

    data = ChartData()
    primary_series = chart_payload["series"][0]
    data.categories = primary_series["labels"]
    if _chart_type_from_payload(chart_payload) == XL_CHART_TYPE.PIE:
        data.add_series(primary_series["name"], primary_series["values"])
    else:
        for item in chart_payload["series"]:
            data.add_series(item["name"], item["values"])

    chart_box_y = y + 0.14
    chart_box_h = max(1.6, h - 0.24)
    graphic_frame = slide.shapes.add_chart(
        _chart_type_from_payload(chart_payload),
        Inches(x + 0.16),
        Inches(chart_box_y),
        Inches(max(2.2, w - 0.30)),
        Inches(chart_box_h),
        data,
    )
    chart = graphic_frame.chart
    chart.has_title = False
    options = chart_payload.get("options") if isinstance(chart_payload.get("options"), dict) else {}
    colors = _chart_colors(chart_payload, palette)
    show_legend = bool(options.get("showLegend", len(chart.series) > 1 or _chart_type_from_payload(chart_payload) == XL_CHART_TYPE.PIE))
    chart.has_legend = show_legend
    if chart.has_legend:
        chart.legend.position = _legend_position(_as_str(options.get("legendPos"), "r"))
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.name = _font_body(preset)
            chart.legend.font.size = Pt(max(typography.caption_min, 10))
        except Exception:
            pass

    if _chart_type_from_payload(chart_payload) == XL_CHART_TYPE.PIE:
        series = chart.series[0]
        series.has_data_labels = True
        try:
            chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
        for index, point in enumerate(series.points):
            color = colors[index % len(colors)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color)
    else:
        for index, series in enumerate(chart.series):
            color = colors[index % len(colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(color)
            series.format.line.color.rgb = _rgb(color)
            try:
                series.format.line.width = Pt(2.2)
            except Exception:
                pass
        try:
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = _rgb(palette["line"])
            chart.value_axis.tick_labels.font.name = _font_caption(preset)
            chart.value_axis.tick_labels.font.size = Pt(max(typography.caption_min, 10))
            chart.category_axis.tick_labels.font.name = _font_caption(preset)
            chart.category_axis.tick_labels.font.size = Pt(max(typography.caption_min, 10))
            cat_axis_title = _as_str(options.get("catAxisTitle"), "")
            if cat_axis_title:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = cat_axis_title
            val_axis_title = _as_str(options.get("valAxisTitle"), "")
            if val_axis_title:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = val_axis_title
        except Exception:
            pass


def _add_chart_content(
    slide: Any,
    *,
    chart_payload: dict[str, Any],
    facts: list[dict[str, str]],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    note_text: str,
) -> None:
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    gutter = layout.gutter
    fact_items = facts[:3]
    note_h = 0.0
    facts_h = 0.0
    if fact_items:
        columns = len(fact_items)
        card_w = (total_w - gutter * (columns - 1)) / columns
        pref_h = [
            _preferred_fact_card_height(
                item,
                width_in=card_w,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                min_h=1.05,
                max_h=1.95,
                include_source=False,
            )
            for item in fact_items
        ]
        facts_h = min(1.95, max(pref_h)) if pref_h else 1.35
    if note_text and not fact_items:
        note_h = 0.50
    chart_h = max(2.2, content_bottom - y - facts_h - note_h - (gutter if fact_items else 0.0) - (0.10 if note_h else 0.0))
    _add_chart_frame(
        slide,
        chart_payload=chart_payload,
        x=x,
        y=y,
        w=total_w,
        h=chart_h,
        palette=palette,
        typography=typography,
        preset=preset,
    )
    if fact_items:
        columns = len(fact_items)
        card_w = (total_w - gutter * (columns - 1)) / columns
        fact_y = y + chart_h + gutter
        for index, item in enumerate(fact_items):
            accent = _as_str(item.get("accent"), "")
            if accent not in palette:
                accent = "accent_primary" if index % 2 == 0 else "accent_secondary"
            _add_fact_card(
                slide,
                item=item,
                x=x + index * (card_w + gutter),
                y=fact_y,
                w=card_w,
                h=facts_h,
                palette=palette,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                accent_key=accent,
                include_source=False,
            )
        return

    if note_text and note_h > 0.0:
        _set_text_box(
            slide,
            text=note_text,
            x=x,
            y=y + chart_h + 0.08,
            w=total_w,
            h=note_h,
            font_name=_font_caption(preset),
            font_size=max(typography.caption_min, 11),
            color_hex=palette["text_muted"],
            margin=0.01,
        )


def _add_timeline_content(
    slide: Any,
    *,
    milestones: list[dict[str, Any]],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    icon_paths: list[Path | None] | None = None,
) -> None:
    x = layout.margin_x
    y = max(1.5, content_top + 0.04)
    total_w = slide_w - x * 2
    items = milestones[:4] if milestones else [
        {"label": "Q1", "title": "Discover", "body": "Define baseline"},
        {"label": "Q2", "title": "Build", "body": "Pilot delivery"},
        {"label": "Q3", "title": "Scale", "body": "Expand coverage"},
        {"label": "Q4", "title": "Optimize", "body": "Harden operations"},
    ]
    count = len(items)
    gutter = layout.gutter
    card_w = (total_w - gutter * (count - 1)) / max(1, count)
    card_w = max(2.15, min(2.65, card_w))
    content_w = card_w * count + gutter * (count - 1)
    start_x = x + max(0.0, (total_w - content_w) / 2.0)
    icons = icon_paths or []
    has_any_icon = any(p is not None for p in icons)
    icon_size = 0.5
    icon_pad = 0.15
    # When icons are present, shift the timeline line down so icons fit above
    # the dot label without overlapping the slide header.
    icon_offset = (icon_size + icon_pad) if has_any_icon else 0.0
    line_y = y + 1.1 + icon_offset
    _add_line(
        slide,
        x=start_x,
        y=line_y,
        w=content_w,
        color_hex=palette["line"],
        width_pt=5.0,
    )

    # Render icons centered above each dot (above the line_y marker).
    if has_any_icon:
        for idx in range(count):
            icon_path = icons[idx] if idx < len(icons) else None
            if icon_path is None:
                continue
            card_x = start_x + idx * (card_w + gutter)
            cx = card_x + card_w / 2.0
            icon_top = line_y - icon_offset
            icon_left = cx - icon_size / 2.0
            _add_icon(slide, path=icon_path, x=icon_left, y=icon_top, size=icon_size)

    # Bug 3 fix: the timeline used to render the milestone label twice —
    # once via the dot marker (drawn below) and again as the first body
    # line inside the card. Keep only the dot label, write the label
    # text inside the dot, and drop it from the card body.
    card_h = 1.8
    for idx, item in enumerate(items):
        title = str(item.get("title", item.get("label", f"Step {idx + 1}"))).strip()
        body = str(item.get("body", item.get("text", ""))).strip()
        body_lines = [part.strip() for part in body.split(".") if part.strip()]
        card_h = max(
            card_h,
            _preferred_card_height(
                title_text=title,
                body_lines=body_lines[:4],
                width_in=card_w,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                min_h=1.55,
                max_h=max(1.7, min(2.2, content_bottom - (line_y + 0.34))),
            ),
        )
    card_h = min(card_h, max(1.7, min(2.2, content_bottom - (line_y + 0.34))))
    marker_height = 0.44
    label_font_size = max(10, min(14, int(typography.caption_max)))
    for idx, item in enumerate(items):
        card_x = start_x + idx * (card_w + gutter)
        cx = card_x + card_w / 2.0
        accent_hex = palette["accent_primary"] if idx % 2 == 0 else palette["accent_secondary"]

        label = str(item.get("label", f"{idx + 1}")).strip()
        # Size the marker to fit the label. Single-digit "1"/"Q1" →
        # circle. Multi-char "3 PM"/"Step 1" → pill. Width scales with
        # label length; height stays constant so the baseline aligns.
        label_chars = len(label)
        if label_chars <= 2:
            marker_w = marker_height  # circle
            shape = MSO_SHAPE.OVAL
        else:
            # ~0.10" per char + 0.24" padding, clamped to card width.
            marker_w = min(card_w - 0.20, max(marker_height, 0.24 + label_chars * 0.10))
            shape = MSO_SHAPE.ROUNDED_RECTANGLE
        marker = slide.shapes.add_shape(
            shape,
            Inches(cx - marker_w / 2.0),
            Inches(line_y - marker_height / 2.0),
            Inches(marker_w),
            Inches(marker_height),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = _rgb(accent_hex)
        marker.line.color.rgb = _rgb(palette["bg_primary"])
        marker.line.width = Pt(1)

        if label:
            # Render the label as readable white text centered in the marker.
            label_box = _set_text_box(
                slide,
                text=label,
                x=cx - marker_w / 2.0,
                y=line_y - marker_height / 2.0,
                w=marker_w,
                h=marker_height,
                font_name=_font_title(preset),
                font_size=label_font_size,
                color_hex="FFFFFF",
                bold=True,
                align=PP_ALIGN.CENTER,
                margin=0.02,
            )
            try:
                label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

        body_lines = [
            part.strip()
            for part in str(item.get("body", item.get("text", ""))).strip().split(".")
            if part.strip()
        ][:4]
        _add_card(
            slide,
            x=card_x,
            y=line_y + 0.34,
            w=card_w,
            h=card_h,
            title=str(item.get("title", item.get("label", f"Step {idx + 1}"))).strip(),
            body_lines=body_lines,
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key="accent_primary" if idx % 2 == 0 else "accent_secondary",
        )


def _add_matrix_content(
    slide: Any,
    *,
    quadrants: list[dict[str, Any]],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    icon_paths: list[Path | None] | None = None,
) -> None:
    x = layout.margin_x
    y = content_top
    total_w = slide_w - x * 2
    total_h = max(2.8, content_bottom - y)
    gutter = layout.gutter
    card_w = (total_w - gutter) / 2
    max_card_h = (total_h - gutter) / 2

    items = quadrants[:4] if quadrants else [
        {"title": "Strengths", "body": "Current advantages and proof points"},
        {"title": "Weaknesses", "body": "Known gaps and constraints"},
        {"title": "Opportunities", "body": "High-probability expansion paths"},
        {"title": "Threats", "body": "External risks to mitigate"},
    ]
    icons = list(icon_paths) if icon_paths else []
    row_pref = [1.35, 1.35]
    prepared: list[tuple[int, int, str, list[str]]] = []
    for idx in range(4):
        item = items[idx] if idx < len(items) else {}
        row = idx // 2
        col = idx % 2
        title = str(item.get("title", f"Quadrant {idx + 1}")).strip()
        body = str(item.get("body", item.get("text", ""))).strip()
        body_lines = [part.strip() for part in body.replace("\n", ". ").split(".") if part.strip()]
        if not body_lines:
            body_lines = [f"Add {title.lower()} content."]
        prepared.append((row, col, title, body_lines[:4]))
        row_pref[row] = max(
            row_pref[row],
            _preferred_card_height(
                title_text=title,
                body_lines=body_lines[:4],
                width_in=card_w,
                typography=typography,
                preset=preset,
                rail_h=layout.rail_height,
                min_h=1.35,
                max_h=max(1.5, min(2.0, max_card_h)),
            ),
        )
    card_h = min(max_card_h, max(row_pref))
    # Matrix cards are ~2x2 and shorter than cards-3 columns, so icons use a
    # smaller footprint (0.4") than the 0.5" used in cards-3.
    icon_size = 0.4
    for idx, (row, col, title, body_lines) in enumerate(prepared):
        card_x = x + col * (card_w + gutter)
        card_y = y + row * (card_h + gutter)
        icon_path = icons[idx] if idx < len(icons) else None
        _add_card(
            slide,
            x=card_x,
            y=card_y,
            w=card_w,
            h=card_h,
            title=title,
            body_lines=body_lines,
            palette=palette,
            typography=typography,
            preset=preset,
            rail_h=layout.rail_height,
            accent_key="accent_primary" if ((row * 2) + col) % 2 == 0 else "accent_secondary",
            icon_path=icon_path,
            icon_size=icon_size if icon_path else 0.0,
        )


def _luminance(hex_color: str) -> float:
    """WCAG-style relative luminance [0..1] for a 6-char hex color."""
    try:
        h = hex_color.lstrip("#")
        if len(h) != 6:
            return 0.5
        r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    except (ValueError, TypeError):
        return 0.5

    def _chan(v: float) -> float:
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

    return 0.2126 * _chan(r) + 0.7152 * _chan(g) + 0.0722 * _chan(b)


def _contrast_ok(fg: str, bg: str, threshold: float = 0.35) -> bool:
    """True when fg and bg luminances differ by more than `threshold`.
    Not full WCAG — just a "will you see it at all" check that rejects
    near-black-on-black or near-white-on-white.
    """
    if not fg or not bg:
        return False
    return abs(_luminance(fg) - _luminance(bg)) >= threshold


def _kpi_value_font_size(value_text: str) -> int:
    """Length-aware autosize for kpi-hero value. Prevents overflow on long values.

    ≤4 chars (e.g., '42%'): 120pt
    5-6 chars (e.g., '98.5%', '$1.2M'): 96pt
    7-8 chars (e.g., '2,456ms'): 72pt
    ≥9 chars: 60pt
    """
    n = len(value_text.strip())
    if n <= 4:
        return 120
    if n <= 6:
        return 96
    if n <= 8:
        return 72
    return 60


def _add_kpi_hero_content(
    slide: Any,
    *,
    value: str,
    label: str,
    context: str,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
    dark: bool = True,
) -> None:
    """Single-KPI poster layout.

    Renders with a dark background by default — kpi-hero's job is to break
    the light-content rhythm, not just the layout rhythm. Set `dark=False`
    in the outline (`"theme": "light"`) to keep it on the deck's default bg.
    """
    value_text = str(value).strip() or "?"
    label_text = str(label).strip()
    context_text = str(context).strip()

    # Resolve colors against the (potentially inverted) palette locally.
    # Only the 4 keys kpi-hero touches — keeps the palette-swap scope
    # minimal per earlier review. Background is set by the caller in
    # _add_content_slide based on the same `dark` decision.
    if dark:
        # Pick the accent that actually reads on a dark bg. Most presets
        # have a bright accent_secondary (amber, orange) — that's the
        # default. But presets like editorial-minimal use accent_secondary
        # as a dark neutral, which would render invisible on pure black.
        # Fall back to accent_primary when the secondary is too dark.
        bg_color = palette.get("bg_dark", "0F172A")
        sec = palette.get("accent_secondary", "")
        prim = palette.get("accent_primary", "F59E0B")
        value_color = sec if _contrast_ok(sec, bg_color) else prim
        if not _contrast_ok(value_color, bg_color):
            value_color = "F59E0B"  # last-resort amber always reads on dark
        label_color = "FFFFFF"
        context_color = "CBD5E1"
    else:
        value_color = palette["accent_primary"]
        label_color = palette["text_primary"]
        context_color = palette["text_muted"]

    content_h = max(2.5, content_bottom - content_top)
    value_font = _kpi_value_font_size(value_text)
    value_h = min(2.6, max(1.6, value_font / 72.0 * 1.35))
    label_font = min(typography.section_max + 6, 28)
    label_h = 0.55
    context_font = min(typography.caption_max + 1, 16)
    context_h = 0.42

    total = value_h + 0.18 + (label_h if label_text else 0.0) + \
            (0.10 + context_h if context_text else 0.0)
    value_y = content_top + max(0.10, (content_h - total) / 2.0)

    _set_text_box(
        slide,
        text=value_text,
        x=layout.margin_x,
        y=value_y,
        w=slide_w - layout.margin_x * 2,
        h=value_h,
        font_name=_font_title(preset),
        font_size=value_font,
        color_hex=value_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0.02,
    )

    next_y = value_y + value_h + 0.18
    if label_text:
        _set_text_box(
            slide,
            text=label_text,
            x=layout.margin_x,
            y=next_y,
            w=slide_w - layout.margin_x * 2,
            h=label_h,
            font_name=_font_title(preset),
            font_size=label_font,
            color_hex=label_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0.02,
        )
        next_y += label_h + 0.10

    if context_text:
        _set_text_box(
            slide,
            text=context_text,
            x=layout.margin_x + 1.2,
            y=next_y,
            w=slide_w - (layout.margin_x + 1.2) * 2,
            h=context_h,
            font_name=_font_caption(preset),
            font_size=context_font,
            color_hex=context_color,
            align=PP_ALIGN.CENTER,
            margin=0.02,
        )


def _add_comparison_content(
    slide: Any,
    *,
    left: dict[str, Any],
    right: dict[str, Any],
    verdict: str,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    """Two-column comparison with dividing rule. Optional verdict strip at bottom."""
    margin_x = layout.margin_x
    gutter = 0.45
    usable_w = slide_w - margin_x * 2
    col_w = (usable_w - gutter) / 2.0

    has_verdict = bool(verdict and verdict.strip())
    verdict_h = 0.6 if has_verdict else 0.0
    verdict_gap = 0.20 if has_verdict else 0.0

    col_top = content_top + 0.05
    col_bottom = content_bottom - verdict_h - verdict_gap
    col_h = max(2.4, col_bottom - col_top)

    def _render_column(side_spec: dict[str, Any], x: float, accent_key: str) -> None:
        side_title = str(side_spec.get("title", "")).strip()
        body = side_spec.get("body", "")
        if isinstance(body, list):
            body_lines = [str(item).strip() for item in body if str(item).strip()]
        else:
            body_lines = [
                part.strip() for part in str(body).replace("\n", ". ").split(".")
                if part.strip()
            ]

        # No accent rule above the title — that pattern ("thin colored line
        # under/over titles") is a documented AI-slide tell. Instead the
        # column gets its identity from (1) oversized, colored title, and
        # (2) the vertical divider between columns.
        title_y = col_top
        title_font = min(typography.section_max + 4, 28)
        title_h = 0.72
        _set_text_box(
            slide,
            text=side_title or "—",
            x=x,
            y=title_y,
            w=col_w,
            h=title_h,
            font_name=_font_title(preset),
            font_size=title_font,
            color_hex=palette[accent_key],
            bold=True,
            margin=0.01,
        )

        body_y = title_y + title_h + 0.14
        body_h = max(0.8, col_h - (body_y - col_top) - 0.08)
        body_font = _card_body_font(
            body_lines if body_lines else [""],
            width_in=col_w,
            typography=typography,
            available_height_in=body_h,
        )
        _set_lines_box(
            slide,
            lines=body_lines or ["Add comparison content."],
            x=x,
            y=body_y,
            w=col_w,
            h=body_h,
            font_name=_font_body(preset),
            font_size=body_font,
            color_hex=palette["text_primary"],
            bullet=True,
            margin=0.04,
        )

    _render_column(left, margin_x, "accent_primary")
    _render_column(right, margin_x + col_w + gutter, "accent_secondary")

    # Vertical divider rule between columns — rendered as a thin filled
    # rectangle (not a connector) so it stays visible at 1.5pt and reads
    # clearly as "A vs B" rather than two parallel lists.
    divider_x = margin_x + col_w + gutter / 2.0 - 0.02
    divider_y_top = col_top + 0.06
    divider_height = col_h - 0.12
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(divider_x),
        Inches(divider_y_top),
        Inches(0.04),
        Inches(divider_height),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _rgb(palette.get("line", palette["text_muted"]))
    divider.line.fill.background()

    if has_verdict:
        verdict_y = col_bottom + verdict_gap
        verdict_x = margin_x + 0.6
        verdict_w = usable_w - 1.2
        # Dark strip gives the synthesis visual weight that reads
        # immediately — a subtle surface panel gets lost. Left accent
        # bar reinforces this is the "answer" of the comparison.
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(verdict_x),
            Inches(verdict_y),
            Inches(verdict_w),
            Inches(verdict_h),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = _rgb(palette.get("bg_dark", "0F172A"))
        strip.line.fill.background()

        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(verdict_x),
            Inches(verdict_y),
            Inches(0.08),
            Inches(verdict_h),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = _rgb(palette.get("accent_primary", "14B8A6"))
        accent_bar.line.fill.background()

        _set_text_box(
            slide,
            text=str(verdict).strip(),
            x=verdict_x + 0.24,
            y=verdict_y + 0.04,
            w=verdict_w - 0.40,
            h=verdict_h - 0.08,
            font_name=_font_body(preset),
            font_size=min(typography.body_max + 1, 16),
            color_hex="FFFFFF",
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0.04,
        )


def _add_table_content(
    slide: Any,
    *,
    headers: list[str],
    rows: list[list[Any]],
    caption: str,
    column_weights: list[float] | None,
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    """Native OOXML table renderer. Used when parallel-field data reads
    better as a grid than as cards. Header row is accent-colored; body
    rows alternate surface for readability.
    """
    margin_x = layout.margin_x
    usable_w = slide_w - margin_x * 2
    has_caption = bool(caption and caption.strip())
    caption_gap = 0.14
    caption_h = 0.42 if has_caption else 0.0

    available_h = max(1.4, content_bottom - content_top - caption_h - caption_gap)

    col_count = max(1, len(headers))
    row_count = len(rows) + 1  # + header row

    # Row height: header slightly taller than body rows.
    # Target total table height: row_count * ~0.5" but clamp to available.
    ideal_row_h = 0.55
    ideal_header_h = 0.62
    total_ideal = ideal_header_h + (row_count - 1) * ideal_row_h
    scale = min(1.0, available_h / total_ideal) if total_ideal > 0 else 1.0
    header_h = max(0.42, ideal_header_h * scale)
    body_row_h = max(0.36, ideal_row_h * scale)
    table_h = header_h + (row_count - 1) * body_row_h

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(margin_x),
        Inches(content_top),
        Inches(usable_w),
        Inches(table_h),
    )
    table = table_shape.table

    # Column widths — proportional to weights, or equal if not provided.
    weights = column_weights if column_weights and len(column_weights) == col_count else None
    if weights:
        total = sum(weights)
        if total <= 0:
            weights = None
    if weights:
        running = 0.0
        for i, w in enumerate(weights):
            col_w = usable_w * (w / sum(weights))
            table.columns[i].width = Inches(col_w)
            running += col_w
    else:
        col_w = usable_w / col_count
        for i in range(col_count):
            table.columns[i].width = Inches(col_w)

    # Row heights
    table.rows[0].height = Inches(header_h)
    for i in range(1, row_count):
        table.rows[i].height = Inches(body_row_h)

    header_fill = palette.get("accent_primary", "14B8A6")
    header_text = "FFFFFF"
    body_text = palette.get("text_primary", "0F172A")
    body_fill_a = palette.get("surface", "FFFFFF")
    body_fill_b = palette.get("bg_primary", "F4F8FB")
    header_font_size = min(typography.body_max + 1, 15)
    body_font_size = min(typography.body_max, 13)

    font_title = _font_title(preset)
    font_body = _font_body(preset)

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_fill)
        cell.text = ""  # clear default
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(header).strip()
        _set_paragraph_style(
            p,
            font_name=font_title,
            font_size=header_font_size,
            color_hex=header_text,
            bold=True,
        )

    for row_idx, row in enumerate(rows, start=1):
        fill_color = body_fill_a if (row_idx % 2) == 1 else body_fill_b
        for col_idx in range(col_count):
            value = row[col_idx] if col_idx < len(row) else ""
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill_color)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(value).strip()
            _set_paragraph_style(
                p,
                font_name=font_body,
                font_size=body_font_size,
                color_hex=body_text,
            )

    if has_caption:
        caption_y = content_top + table_h + caption_gap
        _set_text_box(
            slide,
            text=caption.strip(),
            x=margin_x,
            y=caption_y,
            w=usable_w,
            h=caption_h,
            font_name=_font_caption(preset),
            font_size=min(typography.caption_max + 1, 12),
            color_hex=palette.get("text_muted", "475569"),
            margin=0.01,
        )


def _add_image_sidebar_content(
    slide: Any,
    *,
    image_path: Path | None,
    image_side: str,
    sidebar_sections: list[dict[str, Any]],
    palette: dict[str, str],
    typography: Any,
    preset: Any,
    layout: Any,
    slide_w: float,
    content_bottom: float,
    content_top: float,
) -> None:
    """Image on one side, sidebar with 2-4 labeled sections on the other.

    The dominant lab-data / research-deck composition: a chart, figure,
    or photo takes the visual weight; the sidebar gives the reader
    labeled sections of bullets explaining it. Inspired by the TB
    LAMP-Seq results slides (Run 1/2/3) and the CDC workflow slide.
    """
    margin_x = layout.margin_x
    usable_w = slide_w - margin_x * 2
    content_h = max(2.8, content_bottom - content_top)
    gap = 0.32

    side = (image_side or "left").strip().lower()
    if side not in ("left", "right"):
        side = "left"

    # Image takes ~60% of usable width when present, sidebar takes rest.
    if image_path is not None and image_path.exists():
        image_w = usable_w * 0.58
        sidebar_w = usable_w - image_w - gap
    else:
        image_w = 0.0
        sidebar_w = usable_w

    if side == "left":
        image_x = margin_x
        sidebar_x = margin_x + image_w + (gap if image_w else 0.0)
    else:
        sidebar_x = margin_x
        image_x = margin_x + sidebar_w + gap

    # Place the image — aspect-ratio preserving within the allotted box.
    if image_path is not None and image_path.exists() and image_w > 0:
        try:
            from PIL import Image as PILImage
            with PILImage.open(image_path) as im:
                iw, ih = im.size
            ratio = ih / iw if iw > 0 else 0.75
        except Exception:
            ratio = 0.6
        img_h_fit = min(content_h, image_w * ratio)
        img_w_fit = img_h_fit / ratio if ratio > 0 else image_w
        img_w_fit = min(image_w, img_w_fit)
        img_h_fit = img_w_fit * ratio
        # Center vertically in content zone
        img_y = content_top + (content_h - img_h_fit) / 2.0
        slide.shapes.add_picture(
            str(image_path),
            Inches(image_x),
            Inches(img_y),
            width=Inches(img_w_fit),
            height=Inches(img_h_fit),
        )

    # Sidebar sections — bold label + bullet list per section.
    sections = [s for s in (sidebar_sections or []) if isinstance(s, dict)][:4]
    if not sections:
        return
    n = len(sections)
    # Allocate vertical space per section, account for inter-section gap.
    section_gap = 0.18
    total_gaps = section_gap * max(0, n - 1)
    section_h = (content_h - total_gaps) / n
    label_font = min(typography.section_max, 18)
    body_font = max(typography.body_min, min(typography.body_max, 13))

    for idx, section in enumerate(sections):
        sy = content_top + idx * (section_h + section_gap)
        label = str(section.get("title", "")).strip()
        body_raw = section.get("body", "")
        if isinstance(body_raw, list):
            lines = [str(item).strip() for item in body_raw if str(item).strip()]
        else:
            lines = [
                part.strip()
                for part in str(body_raw).replace("\n", "\n").split("\n")
                if part.strip()
            ]

        if label:
            _set_text_box(
                slide,
                text=label,
                x=sidebar_x,
                y=sy,
                w=sidebar_w,
                h=0.40,
                font_name=_font_title(preset),
                font_size=label_font,
                color_hex=palette.get("accent_primary", "14B8A6"),
                bold=True,
                margin=0.01,
            )
            body_y = sy + 0.48
        else:
            body_y = sy

        body_h = max(0.4, section_h - (body_y - sy) - 0.04)
        if lines:
            _set_lines_box(
                slide,
                lines=lines,
                x=sidebar_x,
                y=body_y,
                w=sidebar_w,
                h=body_h,
                font_name=_font_body(preset),
                font_size=body_font,
                color_hex=palette.get("text_primary", "0F172A"),
                bullet=True,
                margin=0.01,
            )


def _add_content_slide(
    prs: Presentation,
    spec: dict[str, Any],
    preset: Any,
    *,
    config: BuildConfig,
    render_mode: str,
    page_index: int | None = None,
    page_total: int | None = None,
) -> Any:
    typography = preset.typography
    palette = preset.palette
    layout = preset.layout
    visual_intent = _as_str(spec.get("visual_intent"), "").lower()
    if visual_intent == "flow":
        return _add_flow_diagram_layout(
            prs,
            spec,
            preset=preset,
            outline_dir=config.outline_dir,
        )

    slide = prs.slides.add_slide(_blank_layout(prs))
    sw = _slide_w(prs)
    sh = _slide_h(prs)
    chart_payload = _normalize_chart_payload(spec, config.outline_dir)
    facts = _normalize_fact_items(spec, chart_payload)
    footer_spec = spec
    if not spec.get("sources") and chart_payload.get("sources"):
        footer_spec = deepcopy(spec)
        footer_spec["sources"] = chart_payload["sources"]
    footer = _slide_footer_text(footer_spec)
    has_footer = bool(footer)
    thumbnails = spec.get("thumbnails")
    thumb_items = [str(item).strip() for item in thumbnails] if isinstance(thumbnails, list) else []
    lines = _extract_lines(spec)
    hero_asset = _asset_from_spec(spec, "hero_image", config.outline_dir)
    variant = str(spec.get("variant", "")).strip().lower()
    if visual_intent == "hero" and hero_asset and hero_asset.exists():
        variant = "hero"

    if not variant:
        if isinstance(spec.get("cards"), list):
            variant = "cards-3" if len(spec["cards"]) >= 3 else "cards-2"
        elif isinstance(spec.get("milestones"), list):
            variant = "timeline"
        elif isinstance(spec.get("quadrants"), list):
            variant = "matrix"
        elif chart_payload:
            variant = "chart"
        elif facts:
            variant = "stats"
        elif isinstance(spec.get("highlights"), list):
            variant = "split"
        elif visual_intent == "comparison":
            variant = "split"
        else:
            variant = "standard"

    right_reserve = 0.0
    if thumb_items:
        right_reserve = max(right_reserve, 3.05)
    if variant == "hero":
        right_reserve = max(right_reserve, 3.55)
    content_bottom = _content_bottom(sh, has_footer=has_footer, layout=layout)

    # Rhythm-breaking dark treatment for kpi-hero. A single dark content
    # slide in a light-bg deck snaps the eye's rhythm far more than any
    # new layout on the same background. Only 4 palette keys actually get
    # swapped (bg, title text, subtitle text, footer text) — the variant
    # body renders its own colors internally.
    kpi_dark = variant == "kpi-hero" and _as_str(spec.get("theme"), "dark").lower() != "light"
    if kpi_dark:
        _set_background(slide, palette["bg_dark"])
    else:
        _set_background(slide, palette["bg_primary"])
    _add_background_image(
        slide,
        spec,
        slide_w=sw,
        slide_h=sh,
        palette=palette,
        outline_dir=config.outline_dir,
    )
    if render_mode == "express" and not kpi_dark:
        _add_content_motif(slide, palette=palette, slide_w=sw, margin_x=layout.margin_x)
    header_palette = palette
    if kpi_dark:
        header_palette = dict(palette)
        header_palette["text_primary"] = "FFFFFF"
        header_palette["text_muted"] = "CBD5E1"
    header_bottom = _content_header(
        slide,
        title=str(spec.get("title", "Content")),
        subtitle=str(spec.get("subtitle", "")),
        palette=header_palette,
        typography=typography,
        preset=preset,
        margin_x=layout.margin_x,
        right_reserve=right_reserve,
    )
    content_top = max(1.35, header_bottom + 0.12)
    if thumb_items:
        _add_thumbnail_strip(
            slide,
            thumbnails=thumb_items,
            slide_w=sw,
            margin_x=layout.margin_x,
            palette=palette,
            outline_dir=config.outline_dir,
        )

    # Icon support (v1): cards-2/cards-3, timeline, stats.
    # Other variants log a debug warning and ignore icons rather than crash.
    has_icons_requested = isinstance(_slide_assets(spec).get("icons"), list) and bool(
        _slide_assets(spec).get("icons")
    )
    ICON_SUPPORTED_VARIANTS = {"cards-2", "cards-3", "timeline", "stats", "matrix"}
    if has_icons_requested and variant not in ICON_SUPPORTED_VARIANTS:
        print(
            f"[build_deck] debug: assets.icons ignored on variant '{variant}' "
            f"(supported: cards-2, cards-3, timeline, stats, matrix)",
            file=sys.stderr,
        )

    if variant == "hero":
        _add_hero_content(
            slide,
            lines=lines,
            hero_image=hero_asset,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    elif variant == "generated-image":
        generated_asset = _asset_from_spec(spec, "generated_image", config.outline_dir)
        if generated_asset is None:
            generated_asset = _asset_from_spec(spec, "hero_image", config.outline_dir)
        if generated_asset is None:
            generated_asset = _asset_from_spec(spec, "image", config.outline_dir)
        _add_generated_image_content(
            slide,
            image_path=generated_asset,
            spec=spec,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    elif variant in {"cards-2", "cards-3"}:
        cards = spec.get("cards")
        if not isinstance(cards, list):
            cards = []
        if not cards and lines:
            for idx, line in enumerate(lines[:3]):
                cards.append(
                    {
                        "title": f"Pillar {idx + 1}",
                        "body": line,
                        "accent": "accent_primary" if idx % 2 == 0 else "accent_secondary",
                    }
                )
        columns = 2 if variant == "cards-2" else 3
        icon_paths = _slide_icon_paths(spec, config.outline_dir, columns)
        promote_raw = spec.get("promote_card")
        promote_index: int | None = None
        if columns == 3 and isinstance(promote_raw, int) and 0 <= promote_raw < len(cards):
            promote_index = promote_raw
        _add_cards_grid(
            slide,
            cards=cards,
            columns=columns,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            icon_paths=icon_paths,
            promote_card=promote_index,
        )
    elif variant == "split":
        highlights = spec.get("highlights")
        if not isinstance(highlights, list):
            highlights = []
        _add_split_content(
            slide,
            lines=lines,
            highlights=[str(item) for item in highlights],
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    elif variant == "timeline":
        milestones = spec.get("milestones")
        if not isinstance(milestones, list):
            milestones = []
        milestone_count = min(4, len(milestones)) if milestones else 4
        icon_paths = _slide_icon_paths(spec, config.outline_dir, milestone_count)
        _add_timeline_content(
            slide,
            milestones=milestones,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            icon_paths=icon_paths,
        )
    elif variant == "matrix":
        quadrants = spec.get("quadrants")
        if not isinstance(quadrants, list):
            quadrants = []
        icon_paths = _slide_icon_paths(spec, config.outline_dir, 4)
        _add_matrix_content(
            slide,
            quadrants=quadrants,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            icon_paths=icon_paths,
        )
    elif variant == "stats":
        fact_count = min(4, len(facts)) if facts else 3
        icon_paths = _slide_icon_paths(spec, config.outline_dir, fact_count)
        _add_stats_content(
            slide,
            facts=facts,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            icon_paths=icon_paths,
        )
    elif variant == "chart":
        _add_chart_content(
            slide,
            chart_payload=chart_payload,
            facts=facts,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            note_text=_as_str(spec.get("message") or spec.get("caption"), "") or _as_str(chart_payload.get("notes"), ""),
        )
    elif variant == "kpi-hero":
        _add_kpi_hero_content(
            slide,
            value=_as_str(spec.get("value"), ""),
            label=_as_str(spec.get("label"), ""),
            context=_as_str(spec.get("context"), ""),
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            dark=kpi_dark,
        )
    elif variant == "comparison-2col":
        left_spec = spec.get("left") if isinstance(spec.get("left"), dict) else {}
        right_spec = spec.get("right") if isinstance(spec.get("right"), dict) else {}
        _add_comparison_content(
            slide,
            left=left_spec,
            right=right_spec,
            verdict=_as_str(spec.get("verdict"), ""),
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    elif variant == "image-sidebar":
        image_asset = _asset_from_spec(spec, "hero_image", config.outline_dir)
        if image_asset is None:
            image_asset = _asset_from_spec(spec, "image", config.outline_dir)
        sections = spec.get("sidebar_sections") if isinstance(spec.get("sidebar_sections"), list) else []
        _add_image_sidebar_content(
            slide,
            image_path=image_asset,
            image_side=_as_str(spec.get("image_side"), "left"),
            sidebar_sections=sections,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    elif variant == "table":
        headers = spec.get("headers") if isinstance(spec.get("headers"), list) else []
        rows = spec.get("rows") if isinstance(spec.get("rows"), list) else []
        cw_raw = spec.get("column_weights")
        column_weights = (
            [float(w) for w in cw_raw]
            if isinstance(cw_raw, list) and all(isinstance(w, (int, float)) for w in cw_raw)
            else None
        )
        _add_table_content(
            slide,
            headers=[str(h) for h in headers],
            rows=[list(r) if isinstance(r, list) else [r] for r in rows],
            caption=_as_str(spec.get("caption"), ""),
            column_weights=column_weights,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
        )
    else:
        # Bug 1 fix: on standard variant, when the outline supplies both
        # `body` and `bullets`, pass `body` through as an intro paragraph
        # rendered above the bullets inside the card. `_extract_lines`
        # already prefers bullets when both are present, so `lines` holds
        # the bullets and the body text stays separate.
        raw_body = _as_str(spec.get("body"), "").strip()
        raw_bullets = spec.get("bullets")
        has_bullets = isinstance(raw_bullets, list) and any(
            str(item if not isinstance(item, dict) else item.get("text", "")).strip()
            for item in raw_bullets
        )
        intro_paragraph = raw_body if raw_body and has_bullets else None
        card_heading_override = _as_str(spec.get("card_heading"), "").strip()
        _add_standard_content(
            slide,
            lines=lines,
            palette=palette,
            typography=typography,
            preset=preset,
            layout=layout,
            slide_w=sw,
            content_bottom=content_bottom,
            content_top=content_top,
            intro_paragraph=intro_paragraph,
            card_heading=card_heading_override,
        )

    if render_mode == "express":
        _add_visual_assets(
            slide,
            spec=spec,
            slide_w=sw,
            margin_x=layout.margin_x,
            outline_dir=config.outline_dir,
        )

    # Universal summary callout (the "oval box for key points" pattern
    # from lab-data decks). Sits above the footer on any variant when
    # the outline sets `summary_callout`. Skip on variants that already
    # carry their own bottom emphasis (kpi-hero is the callout itself;
    # comparison-2col already has a verdict strip; pull-quote is a
    # single-element variant).
    callout_text = _as_str(spec.get("summary_callout"), "")
    callout_conflicts = {
        "kpi-hero",
        "generated-image",
        "comparison-2col" if _as_str(spec.get("verdict"), "") else "",
        "pull-quote",
    }
    if callout_text and variant not in callout_conflicts:
        _add_summary_callout(
            slide,
            text=callout_text,
            palette=palette,
            preset=preset,
            slide_w=sw,
            slide_h=sh,
            margin_x=layout.margin_x,
            has_footer=bool(footer),
        )

    show_page = bool(config.deck_style.show_page_numbers)
    if footer or show_page:
        _add_footer(
            slide,
            text=footer or "",
            palette=palette,
            preset=preset,
            slide_h=sh,
            margin_x=layout.margin_x,
            dark=kpi_dark,
            page_index=page_index if show_page else None,
            page_total=page_total if show_page else None,
            slide_w=sw,
        )
    return slide


def _build(prs: Presentation, data: dict[str, Any], preset: Any, *, config: BuildConfig) -> int:
    slides = data.get("slides")
    count = 0

    if not isinstance(slides, list) or len(slides) == 0:
        title = str(data.get("title", "")).strip()
        subtitle = str(data.get("subtitle", "")).strip()
        if title or subtitle:
            _add_title_slide(
                prs,
                {"title": title or "Presentation", "subtitle": subtitle},
                preset,
                config=config,
            )
            return 1
        raise ValueError("Outline JSON must include a non-empty slides array.")

    for index, item in enumerate(slides, start=1):
        if not isinstance(item, dict):
            continue
        slide_type = str(item.get("type", "content")).strip().lower()
        spec = _apply_emoji_policy(
            item,
            slide_type=slide_type,
            preset_name=preset.name,
            emoji_mode=config.deck_style.emoji_mode,
        )
        resolved_mode = _resolve_render_mode(
            spec,
            slide_type=slide_type,
            deck_density=config.deck_style.visual_density,
        )
        try:
            if slide_type == "title":
                slide = _add_title_slide(prs, spec, preset, config=config)
                if resolved_mode == "express":
                    _add_visual_assets(
                        slide,
                        spec=spec,
                        slide_w=_slide_w(prs),
                        margin_x=preset.layout.margin_x,
                        outline_dir=config.outline_dir,
                    )
            elif slide_type == "section":
                slide = _add_section_slide(prs, spec, preset, config=config)
                if resolved_mode == "express":
                    _add_visual_assets(
                        slide,
                        spec=spec,
                        slide_w=_slide_w(prs),
                        margin_x=preset.layout.margin_x,
                        outline_dir=config.outline_dir,
                    )
            else:
                _add_content_slide(
                    prs,
                    spec,
                    preset,
                    config=config,
                    render_mode=resolved_mode,
                    page_index=index,
                    page_total=len(slides),
                )
        except Exception as exc:
            if resolved_mode != "express":
                raise
            # Hard fallback: express path must never break full deck generation.
            print(
                f"[warn] slide {index}: express render failed ({exc}); "
                "falling back to reliable mode."
            )
            if slide_type == "title":
                _add_title_slide(prs, spec, preset, config=config)
            elif slide_type == "section":
                _add_section_slide(prs, spec, preset, config=config)
            else:
                _add_content_slide(
                    prs,
                    spec,
                    preset,
                    config=config,
                    render_mode="reliable",
                    page_index=index,
                    page_total=len(slides),
                )
        count += 1

    return count


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a .pptx deck from JSON outline.")
    parser.add_argument("--outline", required=True, help="Path to outline JSON file")
    parser.add_argument("--output", required=True, help="Output .pptx file")
    parser.add_argument("--template", help="Optional template .pptx")
    parser.add_argument(
        "--style-preset",
        default="executive-clinical",
        help="Style preset for typography defaults",
    )
    parser.add_argument(
        "--layout-16x9",
        action="store_true",
        help="Force 16:9 wide layout (13.333x7.5) when no template is used",
    )
    parser.add_argument(
        "--font-pair",
        help="Optional font pair key override (for example: editorial_serif_v1)",
    )
    parser.add_argument(
        "--palette-key",
        help="Optional palette key override (for example: climate_coastal_v1)",
    )
    parser.add_argument("--overwrite", action="store_true", help="Overwrite output if it exists")
    return parser.parse_args()


def main() -> int:
    args = _args()
    outline_path = Path(args.outline).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()
    template_path = Path(args.template).expanduser().resolve() if args.template else None
    base_preset = get_style_preset(args.style_preset)

    if not outline_path.exists():
        raise FileNotFoundError(f"Outline file not found: {outline_path}")
    if template_path and not template_path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")
    if output_path.exists() and not args.overwrite:
        raise FileExistsError(f"Output exists: {output_path}. Use --overwrite to replace it.")

    with outline_path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    if not isinstance(data, dict):
        raise ValueError("Outline JSON root must be an object.")
    build_config = _load_build_config(data, outline_path.parent)

    deck_style = build_config.deck_style
    if args.font_pair:
        deck_style = DeckStyleConfig(
            font_pair=str(args.font_pair).strip().lower(),
            palette_key=deck_style.palette_key,
            visual_density=deck_style.visual_density,
            emoji_mode=deck_style.emoji_mode,
            title_motif=deck_style.title_motif,
        )
    if args.palette_key:
        deck_style = DeckStyleConfig(
            font_pair=deck_style.font_pair,
            palette_key=str(args.palette_key).strip().lower(),
            visual_density=deck_style.visual_density,
            emoji_mode=deck_style.emoji_mode,
            title_motif=deck_style.title_motif,
        )
    build_config = BuildConfig(
        outline_dir=build_config.outline_dir,
        deck_style=deck_style,
        compliance=build_config.compliance,
    )

    try:
        font_tokens = get_font_pair(build_config.deck_style.font_pair)
    except ValueError as exc:
        print(f"[warn] {exc} Falling back to {DEFAULT_FONT_PAIR_KEY}.")
        font_tokens = get_font_pair(DEFAULT_FONT_PAIR_KEY)

    runtime_preset = RuntimePreset(
        name=base_preset.name,
        palette=_resolve_palette(base_preset.palette, build_config.deck_style.palette_key),
        typography=base_preset.typography,
        layout=base_preset.layout,
        font_pair={
            "title": font_tokens.title,
            "body": font_tokens.body,
            "caption": font_tokens.caption,
        },
    )
    _enforce_compliance(data, build_config)

    prs = Presentation(str(template_path)) if template_path else Presentation()
    if not template_path and args.layout_16x9:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif not template_path and not args.layout_16x9:
        # Default to 16:9 for reproducible deck QA in this workspace.
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slide_count = _build(prs, data, runtime_preset, config=build_config)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(
        f"Wrote {slide_count} slide(s) to {output_path} "
        f"with style preset '{runtime_preset.name}', "
        f"font pair '{build_config.deck_style.font_pair}', "
        f"palette '{build_config.deck_style.palette_key or 'preset-default'}'."
    )
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
