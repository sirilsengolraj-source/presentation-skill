"""Edit an existing .pptx file in place-style (input -> output).

Three subcommands cover the most common edits on decks you didn't generate:

- `replace-text`: find/replace across every shape text frame and run.
- `list-slides`: print a 1-based index + title summary of each slide.
- `delete-slide`: drop a slide (1-based) and its presentation relationship.

These operations use python-pptx only. For XML-level edits python-pptx can't
express, use `unpack_pptx.py` / `pack_pptx.py` instead.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation


def _resolve_input(path_str: str) -> Path:
    path = Path(path_str).expanduser().resolve()
    if not path.exists():
        print(f"Error: input file not found: {path}", file=sys.stderr)
        raise SystemExit(2)
    if not path.is_file():
        print(f"Error: input is not a file: {path}", file=sys.stderr)
        raise SystemExit(2)
    return path


def _resolve_output(path_str: str) -> Path:
    path = Path(path_str).expanduser().resolve()
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _guard_output(output_path: Path, input_path: Path, overwrite: bool) -> None:
    """Refuse clobbering and same-path writes. Raises SystemExit(3) on failure."""
    if output_path.resolve() == input_path.resolve():
        print(
            f"Error: --output resolves to the same path as --input: {output_path}. "
            f"Point --output to a different file.",
            file=sys.stderr,
        )
        raise SystemExit(3)
    if output_path.exists() and not overwrite:
        print(
            f"Error: output exists (use --overwrite to replace): {output_path}",
            file=sys.stderr,
        )
        raise SystemExit(3)


def _title_text(slide: Any) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            return (title_shape.text or "").strip()
    except Exception:
        pass
    return ""


def _iter_text_frames(shape: Any) -> Any:
    """Yield text frames from a shape.

    Descends into group shapes and table cells. Does NOT visit chart,
    SmartArt, or slide-master text (see references/editing.md).
    """
    # Group shape: recurse into children. Groups can themselves contain
    # tables, so we route through the same dispatcher.
    if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from _iter_text_frames(child)
        return
    # Table: visit every cell's text frame.
    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                cell_frame = getattr(cell, "text_frame", None)
                if cell_frame is not None:
                    yield cell_frame
        return
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame


def _replace_in_runs(text_frame: Any, find: str, replace: str) -> int:
    """Replace occurrences of `find` inside every run of `text_frame`.

    Run-level replacement preserves formatting. Matches that span run
    boundaries will not be caught; this mirrors how most tooling treats
    python-pptx replacements.
    """
    count = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text and find in run.text:
                count += run.text.count(find)
                run.text = run.text.replace(find, replace)
    return count


def _cmd_replace_text(args: argparse.Namespace) -> int:
    if args.find == "":
        print("Error: --find must be a non-empty string", file=sys.stderr)
        return 2

    input_path = _resolve_input(args.input)
    output_path = _resolve_output(args.output)
    _guard_output(output_path, input_path, args.overwrite)

    presentation = Presentation(str(input_path))
    total = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            for text_frame in _iter_text_frames(shape):
                total += _replace_in_runs(text_frame, args.find, args.replace)
        if getattr(slide, "has_notes_slide", False):
            notes_frame = slide.notes_slide.notes_text_frame
            total += _replace_in_runs(notes_frame, args.find, args.replace)

    if total == 0 and args.require_match:
        print(
            f"Error: --require-match set but 0 occurrence(s) of {args.find!r} found; "
            f"output not written.",
            file=sys.stderr,
        )
        return 4

    presentation.save(str(output_path))
    print(f"Replaced {total} occurrence(s) of {args.find!r} -> {args.replace!r}")
    print(f"Wrote {output_path}")
    return 0


def _cmd_list_slides(args: argparse.Namespace) -> int:
    input_path = _resolve_input(args.input)
    presentation = Presentation(str(input_path))
    print(f"File: {input_path}")
    print(f"Slides: {len(presentation.slides)}")
    for idx, slide in enumerate(presentation.slides, start=1):
        title = _title_text(slide) or "(No title)"
        print(f"  {idx:>3}. {title}")
    return 0


def _cmd_delete_slide(args: argparse.Namespace) -> int:
    input_path = _resolve_input(args.input)
    output_path = _resolve_output(args.output)
    _guard_output(output_path, input_path, args.overwrite)

    presentation = Presentation(str(input_path))
    slides = presentation.slides
    total = len(slides)
    index = int(args.index)
    if index < 1 or index > total:
        print(
            f"Error: --index {index} out of range (deck has {total} slide(s))",
            file=sys.stderr,
        )
        return 2

    # python-pptx has no public delete API. We remove the sldId entry from the
    # presentation's sldIdLst and drop the matching relationship. The slide
    # part itself is left in the package; PowerPoint ignores orphaned parts.
    zero_based = index - 1
    slide_id_list = slides._sldIdLst  # type: ignore[attr-defined]
    sld_id_elements = list(slide_id_list)
    target = sld_id_elements[zero_based]
    rId = target.rId
    slide_id_list.remove(target)
    presentation.part.drop_rel(rId)

    presentation.save(str(output_path))
    print(f"Deleted slide {index} of {total}; wrote {output_path}")
    return 0


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Edit an existing .pptx deck (text replace, slide list, slide delete).",
    )
    sub = parser.add_subparsers(dest="command", required=True)

    replace = sub.add_parser(
        "replace-text",
        help="Replace text across shapes, table cells, groups, and notes.",
    )
    replace.add_argument("--input", required=True, help="Input .pptx path")
    replace.add_argument("--output", required=True, help="Output .pptx path")
    replace.add_argument("--find", required=True, help="Exact substring to find")
    replace.add_argument("--replace", required=True, help="Replacement substring")
    replace.add_argument(
        "--overwrite",
        action="store_true",
        help="Allow overwriting an existing --output file",
    )
    replace.add_argument(
        "--require-match",
        action="store_true",
        help="Exit non-zero (4) if 0 occurrences are replaced",
    )
    replace.set_defaults(func=_cmd_replace_text)

    listing = sub.add_parser(
        "list-slides",
        help="Print a 1-based index + title for every slide.",
    )
    listing.add_argument("--input", required=True, help="Input .pptx path")
    listing.set_defaults(func=_cmd_list_slides)

    delete = sub.add_parser(
        "delete-slide",
        help="Delete a slide by 1-based index.",
    )
    delete.add_argument("--input", required=True, help="Input .pptx path")
    delete.add_argument("--output", required=True, help="Output .pptx path")
    delete.add_argument(
        "--index",
        required=True,
        type=int,
        help="1-based slide index to delete",
    )
    delete.add_argument(
        "--overwrite",
        action="store_true",
        help="Allow overwriting an existing --output file",
    )
    delete.set_defaults(func=_cmd_delete_slide)

    return parser


def main() -> int:
    parser = _build_parser()
    args = parser.parse_args()
    return int(args.func(args))


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except SystemExit:
        raise
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}", file=sys.stderr)
        raise SystemExit(1)
