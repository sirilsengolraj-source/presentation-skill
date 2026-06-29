#!/usr/bin/env python3
"""Extract a lightweight PPTX text inventory and flag likely text issues.

This implementation is intentionally limited to the open-source deck-builder
workflow in this repository. It focuses on:
- text overflow estimated from box size, font size, and copy length
- text box overlap between independent text shapes on the same slide

The JSON structure stays compatible with the repository QA tools:
{
  "slide-01": {
    "shape-003": {
      "name": "TextBox 3",
      "text": "Example",
      "bounds_in": {"x": 1.0, "y": 2.0, "w": 3.0, "h": 1.2},
      "overflow": {"overflow_inches": 0.18},
      "overlap": [{"with": "shape-004", "overlap_inches": 0.09}]
    }
  }
 }
"""

from __future__ import annotations

import argparse
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_INCH = 914400.0
MIN_OVERLAP_IN = 0.08
MIN_OVERFLOW_IN = 0.12


@dataclass(frozen=True)
class ParagraphData:
    text: str
    font_size: float | None


@dataclass
class ShapeData:
    shape: Any
    name: str
    text: str
    width: float
    height: float
    paragraphs: list[ParagraphData]
    frame_overflow_bottom: float | None
    slide_overflow_right: float | None
    slide_overflow_bottom: float | None
    slide_height_emu: int | None


InventoryData = dict[str, dict[str, ShapeData]]


def _inches(emu: int | float) -> float:
    return round(float(emu) / EMU_PER_INCH, 3)


def _shape_bounds(shape: Any) -> dict[str, float]:
    return {
        "x": _inches(shape.left),
        "y": _inches(shape.top),
        "w": _inches(shape.width),
        "h": _inches(shape.height),
    }


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").strip()


def _font_size_pt(shape: Any) -> float:
    if not getattr(shape, "has_text_frame", False):
        return 12.0
    for paragraph in shape.text_frame.paragraphs:
        paragraph_font = getattr(paragraph, "font", None)
        if paragraph_font is not None and paragraph_font.size is not None:
            return float(paragraph_font.size.pt)
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return 12.0


def _paragraphs(shape: Any) -> list[ParagraphData]:
    if not getattr(shape, "has_text_frame", False):
        return []
    rows: list[ParagraphData] = []
    for paragraph in shape.text_frame.paragraphs:
        font_size = None
        if getattr(paragraph, "font", None) is not None and paragraph.font.size is not None:
            font_size = float(paragraph.font.size.pt)
        elif paragraph.runs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    font_size = float(run.font.size.pt)
                    break
        rows.append(
            ParagraphData(
                text=(paragraph.text or "").strip(),
                font_size=font_size,
            )
        )
    return rows


def _estimate_lines(text: str, width_in: float, font_size_pt: float) -> int:
    if not text.strip():
        return 0
    usable_width = max(0.4, width_in - 0.16)
    char_width_in = max(0.05, (font_size_pt / 72.0) * 0.44)
    chars_per_line = max(6, int(usable_width / char_width_in))
    total = 0
    for paragraph in text.splitlines() or [text]:
        stripped = paragraph.strip()
        total += max(1, math.ceil(len(stripped) / chars_per_line)) if stripped else 1
    return total


def _estimate_text_height(text: str, width_in: float, font_size_pt: float) -> float:
    lines = _estimate_lines(text, width_in, font_size_pt)
    line_height_in = (font_size_pt / 72.0) * 1.22
    return round(lines * line_height_in + 0.08, 3)


def _overflow_amount(shape: Any, text: str) -> float:
    bounds = _shape_bounds(shape)
    if bounds["h"] <= 0.0:
        return 0.0
    font_size_pt = _font_size_pt(shape)
    estimated_h = _estimate_text_height(text, bounds["w"], font_size_pt)
    overflow = max(0.0, estimated_h - max(0.0, bounds["h"] - 0.06))
    return round(overflow, 3)


def _slide_overflow_amount(shape: Any, prs: Presentation) -> tuple[float, float]:
    right = max(0.0, (shape.left + shape.width - prs.slide_width) / EMU_PER_INCH)
    bottom = max(0.0, (shape.top + shape.height - prs.slide_height) / EMU_PER_INCH)
    return round(right, 3), round(bottom, 3)


def _overlap_amount(a: dict[str, float], b: dict[str, float]) -> float:
    overlap_x = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
    overlap_y = min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"])
    if overlap_x <= MIN_OVERLAP_IN or overlap_y <= MIN_OVERLAP_IN:
        return 0.0
    return round(min(overlap_x, overlap_y), 3)


def _inventory(prs: Presentation, *, issues_only: bool) -> dict[str, dict[str, Any]]:
    payload: dict[str, dict[str, Any]] = {}

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{slide_index:02d}"
        shapes: list[tuple[str, Any, str, dict[str, float]]] = []
        slide_payload: dict[str, Any] = {}

        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if not text:
                continue
            shape_key = f"shape-{shape_index:03d}"
            bounds = _shape_bounds(shape)
            entry = {
                "name": getattr(shape, "name", shape_key),
                "text": text,
                "bounds_in": bounds,
            }
            overflow = _overflow_amount(shape, text)
            if overflow > MIN_OVERFLOW_IN:
                entry["overflow"] = {"overflow_inches": overflow}
            slide_payload[shape_key] = entry
            shapes.append((shape_key, shape, text, bounds))

        for index, (shape_key, _, _, bounds) in enumerate(shapes):
            overlaps: list[dict[str, Any]] = []
            for other_key, _, _, other_bounds in shapes[index + 1 :]:
                overlap = _overlap_amount(bounds, other_bounds)
                if overlap <= 0:
                    continue
                overlaps.append(
                    {
                        "with": other_key,
                        "overlap_inches": overlap,
                    }
                )
                slide_payload[other_key].setdefault("overlap", []).append(
                    {
                        "with": shape_key,
                        "overlap_inches": overlap,
                    }
                )
            if overlaps:
                slide_payload[shape_key]["overlap"] = overlaps

        if issues_only:
            slide_payload = {
                key: value
                for key, value in slide_payload.items()
                if "overflow" in value or "overlap" in value
            }

        if slide_payload:
            payload[slide_key] = slide_payload

    return payload


def extract_text_inventory(input_pptx: Path | str, prs: Presentation | None = None) -> InventoryData:
    presentation = prs or Presentation(str(input_pptx))
    payload: InventoryData = {}

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_key = f"slide-{slide_index:02d}"
        slide_payload: dict[str, ShapeData] = {}
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if not text:
                continue
            shape_key = f"shape-{shape_index:03d}"
            bounds = _shape_bounds(shape)
            overflow = _overflow_amount(shape, text)
            overflow_right, overflow_bottom = _slide_overflow_amount(shape, presentation)
            slide_payload[shape_key] = ShapeData(
                shape=shape,
                name=getattr(shape, "name", shape_key),
                text=text,
                width=bounds["w"],
                height=bounds["h"],
                paragraphs=_paragraphs(shape),
                frame_overflow_bottom=overflow if overflow > 0 else 0.0,
                slide_overflow_right=overflow_right if overflow_right > 0 else 0.0,
                slide_overflow_bottom=overflow_bottom if overflow_bottom > 0 else 0.0,
                slide_height_emu=int(presentation.slide_height),
            )
        if slide_payload:
            payload[slide_key] = slide_payload

    return payload


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract PPTX text inventory")
    parser.add_argument("input_pptx", help="Input .pptx file")
    parser.add_argument("output_json", help="Output JSON path")
    parser.add_argument(
        "--issues-only",
        action="store_true",
        help="Only write text shapes with overflow or overlap issues",
    )
    return parser.parse_args()


def main() -> int:
    args = _args()
    input_path = Path(args.input_pptx).expanduser().resolve()
    output_path = Path(args.output_json).expanduser().resolve()

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    print(f"Extracting text inventory from: {input_path}")
    if args.issues_only:
        print("Filtering to include only text shapes with issues (overflow/overlap)")

    prs = Presentation(str(input_path))
    payload = _inventory(prs, issues_only=args.issues_only)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    print(f"Output saved to: {output_path}")

    if args.issues_only:
        issue_shapes = sum(
            len(slide_shapes)
            for slide_shapes in payload.values()
            if isinstance(slide_shapes, dict)
        )
        if issue_shapes == 0:
            print("No issues discovered")
        else:
            print(f"Found {issue_shapes} text elements with issues in {len(payload)} slides")
    else:
        shape_count = sum(
            len(slide_shapes)
            for slide_shapes in payload.values()
            if isinstance(slide_shapes, dict)
        )
        print(f"Indexed {shape_count} text elements across {len(payload)} slides")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
